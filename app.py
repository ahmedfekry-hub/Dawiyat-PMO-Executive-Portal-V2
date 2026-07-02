
import base64
import hashlib
import hmac
import time
import io
import json
import re
import shutil
import smtplib
from datetime import datetime, timezone, timedelta
from email.message import EmailMessage
from pathlib import Path
from typing import Any, Dict, List, Mapping, Tuple

import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from pptx.dml.color import RGBColor


BASE_DIR = Path(__file__).parent
DASHBOARD_PATH = BASE_DIR / "dashboard" / "dashboard.html"
DATA_DIR = BASE_DIR / "data"
BACKUP_DIR = BASE_DIR / "backups"
BACKUP_DIR.mkdir(exist_ok=True)

WO_PATH = DATA_DIR / "u_osp_work_order.csv"
PROJECT_UPDATES_PATH = DATA_DIR / "project_updates.csv"
CHANGE_LOG_PATH = DATA_DIR / "change_log.csv"
NOTIFICATIONS_PATH = DATA_DIR / "notifications.csv"
NOTIFICATION_ACCESS_PATH = DATA_DIR / "notification_access.csv"
DAILY_DIGEST_PATH = DATA_DIR / "daily_digest.csv"
WHATSAPP_OUTBOX_PATH = DATA_DIR / "whatsapp_outbox.csv"
MASTER_OPERATIONAL_PATH = DATA_DIR / "master_operational_data.csv"
CHANGE_IMPACT_PATH = DATA_DIR / "change_impact.csv"
KPI_IMPACT_PATH = DATA_DIR / "kpi_impact.csv"
SNAPSHOT_DIR = DATA_DIR / "snapshots"
PENALTIES_PATH = DATA_DIR / "Penalties.csv"
DOC_STATUS_CACHE_PATH = DATA_DIR / "Dawiyat_Document_Status.csv"
DISTRICT_PATH = DATA_DIR / "District.csv"
PERMISSIONS_XLSX_PATH = DATA_DIR / "permissions.xlsx"

DATA_FILES = {
    "u_osp_work_order.csv": WO_PATH,
    "project_updates.csv": PROJECT_UPDATES_PATH,
    "change_log.csv": CHANGE_LOG_PATH,
    "notifications.csv": NOTIFICATIONS_PATH,
    "notification_access.csv": NOTIFICATION_ACCESS_PATH,
    "daily_digest.csv": DAILY_DIGEST_PATH,
    "whatsapp_outbox.csv": WHATSAPP_OUTBOX_PATH,
    "master_operational_data.csv": MASTER_OPERATIONAL_PATH,
    "change_impact.csv": CHANGE_IMPACT_PATH,
    "kpi_impact.csv": KPI_IMPACT_PATH,
    "Penalties.csv": PENALTIES_PATH,
    "District.csv": DISTRICT_PATH,
}

ASSETS_DIR = BASE_DIR / "assets"
DAWIYAT_LOGO_PATH = ASSETS_DIR / "dawiyat_logo.jpg"
MET_LOGO_PATH = ASSETS_DIR / "met_logo.jpg"
PPT_COVER_PATH = ASSETS_DIR / "ppt_cover.png"

DOCUMENT_TYPES = ["Design", "Permit", "Photos", "PAT", "AsBuilt", "Handover", "Commercial"]
DOCUMENT_FOLDER_MAP = {
    "Design": "01 Design",
    "Permit": "02 Permit",
    "Photos": "03 Photos",
    "PAT": "04 PAT",
    "AsBuilt": "05 AsBuilt",
    "Handover": "06 Handover",
    "Commercial": "07 Commercial",
}
DOCUMENT_EXTENSIONS = {
    "Design": ["pdf", "zip", "dwg", "dxf", "xlsx", "xls", "docx"],
    "Permit": ["pdf", "jpg", "jpeg", "png", "xlsx", "xls", "docx"],
    "Photos": ["jpg", "jpeg", "png", "webp", "zip", "rar", "7z"],
    "PAT": ["pdf", "xlsx", "xls", "docx", "jpg", "jpeg", "png"],
    "AsBuilt": ["pdf", "zip", "dwg", "dxf", "xlsx", "xls", "docx"],
    "Handover": ["pdf", "xlsx", "xls", "docx", "jpg", "jpeg", "png", "zip"],
    "Commercial": ["pdf", "xlsx", "xls", "docx", "zip"],
}
GOOGLE_DRIVE_FOLDER_MIME = "application/vnd.google-apps.folder"

# Editable update center configuration
PROJECT_UPDATE_EDITABLE_COLUMNS = [
    "Fiber Status", "Civil Status", "FULL WO STATUS",
    "Invoice Status", "SOR Status", "1st 50 Invoice Status",
    "1st 50 Invoice Cost Amount", "PAT Status", "AsBuilt  Status",
    "DCR_Status", "RFS Certificate", "As-built BOQ", "Redline",
    "Handover O&M _Status", "Handover Consultant _Status", "SOR Reference Number",
    "Asbuilt Final Amount",
]

PROJECT_UPDATE_STATUS_OPTIONS = {
    "Fiber Status": ["Not Started", "Completed", "In Progress"],
    "Civil Status": ["Not Started", "Completed", "In Progress"],
    "FULL WO STATUS": ["Not Started", "Completed", "In Progress"],
    "Invoice Status": ["Civil Not Start", "As-Built Stage", "Fibre Materials", "Waiting SOR", "SOR 50% Stage", "PT Stage", "Civil On-Hold", "PT Stage - POP Issue", "Civil In Progress"],
    "SOR Status": ["Created", "Not Create", "Requested"],
    "1st 50 Invoice Status": ["SOR not Create", "Submitted", "Not Start", "In Progress"],
    "PAT Status": ["Not Started", "Scheduled", "Under Progress", "Re-Test Required", "Under Review", "Accepted", "Failed"],
    "AsBuilt  Status": ["Not Started", "Under Preparation", "Under Review", "Submitted", "Approved", "Rejected"],
    "DCR_Status": ["Not Started", "Under Preparation", "Submitted", "Under Review", "Approved", "Missing Documents", "Rejected"],
    "RFS Certificate": ["Not Started", "Under Preparation", "Submitted", "Under Review", "Approved", "Missing Documents", "Rejected"],
    "As-built BOQ": ["Not Started", "Under Preparation", "Submitted", "Under Review", "Approved", "Missing Documents", "Rejected"],
    "Redline": ["Not Started", "Under Preparation", "Submitted", "Under Review", "Approved", "Missing Documents", "Rejected"],
    "Handover O&M _Status": ["Not Started", "Scheduled", "Waiting OILS Sheet from MO", "Approved", "Rejected", "Under Clearing Remarks"],
    "Handover Consultant _Status": ["Not Started", "Scheduled", "Waiting OILS Sheet from MO", "Approved", "Rejected", "Under Clearing Remarks"],
}

PROJECT_UPDATE_FINANCE_COLUMNS = ["Invoice Status", "SOR Status", "1st 50 Invoice Status", "1st 50 Invoice Cost Amount", "SOR Reference Number"]
PROJECT_UPDATE_PM_COLUMNS = ["Fiber Status", "Civil Status", "FULL WO STATUS", "PAT Status", "AsBuilt  Status", "DCR_Status", "RFS Certificate", "As-built BOQ", "Redline", "Handover O&M _Status", "Handover Consultant _Status", "Asbuilt Final Amount"]
PROJECT_MASTER_ADMIN_FILES = {"District.csv": DISTRICT_PATH, "Penalties.csv": PENALTIES_PATH}

KSA_TZ = timezone(timedelta(hours=3))

def ksa_now() -> datetime:
    return datetime.now(KSA_TZ)


def image_to_base64(path: Path) -> str:
    if not path.exists():
        return ""
    return base64.b64encode(path.read_bytes()).decode("ascii")



ROLE_PERMISSIONS = {
    "admin": {
        "dashboard": True, "assistant": True, "alerts": True, "reports": True,
        "admin": True, "upload": True, "email": True, "documents": True, "ppt_builder": True,
        "export": True, "export_excel": True, "export_pdf": True, "export_ppt": True,
        "pages": ["Dashboard", "AI Executive Assistant", "Smart Alerts", "Executive Reports", "📊 Executive PPT Builder", "Upload CSV", "📤 Document Upload Center", "Admin Board"],
        "dashboard_tabs": ["overview", "tables", "pmo", "performance", "perf-explanation", "decision", "reports"],
        "hide_buttons": [], "hide_tables": [],
    },
    "pmo": {
        "dashboard": True, "assistant": True, "alerts": True, "reports": True,
        "admin": False, "upload": True, "email": False, "documents": True, "ppt_builder": True,
        "export": True, "export_excel": True, "export_pdf": True, "export_ppt": True,
        "pages": ["Dashboard", "AI Executive Assistant", "Smart Alerts", "Executive Reports", "📊 Executive PPT Builder", "Upload CSV", "📤 Document Upload Center"],
        "dashboard_tabs": ["overview", "tables", "pmo", "performance", "perf-explanation", "decision", "reports"],
        "hide_buttons": [], "hide_tables": [],
    },
    "board": {
        "dashboard": True, "assistant": False, "alerts": False, "reports": True,
        "admin": False, "upload": False, "email": False, "documents": False, "ppt_builder": True,
        "export": True, "export_excel": False, "export_pdf": True, "export_ppt": True,
        "pages": ["Dashboard", "Executive Reports", "📊 Executive PPT Builder"],
        "dashboard_tabs": ["overview", "performance", "decision", "reports"],
        "hide_buttons": ["Export Excel", "Upload", "Delete", "Import"],
        "hide_tables": ["PMO Audit", "Missing MET Actual", "Raw Data"],
    },
    "finance": {
        "dashboard": True, "assistant": False, "alerts": True, "reports": True,
        "admin": False, "upload": False, "email": False, "documents": False, "ppt_builder": False,
        "export": True, "export_excel": True, "export_pdf": True, "export_ppt": False,
        "pages": ["Dashboard", "Smart Alerts", "Executive Reports"],
        "dashboard_tabs": ["overview", "tables", "reports"],
        "hide_buttons": ["Upload", "Delete", "Import"],
        "hide_tables": ["PMO Audit", "Missing MET Actual", "PM Review"],
    },
    "audit": {
        "dashboard": True, "assistant": False, "alerts": True, "reports": True,
        "admin": False, "upload": False, "email": False, "documents": False, "ppt_builder": False,
        "export": True, "export_excel": True, "export_pdf": True, "export_ppt": False,
        "pages": ["Dashboard", "Smart Alerts", "Executive Reports"],
        "dashboard_tabs": ["pmo", "reports"],
        "hide_buttons": ["Upload", "Delete", "Import"],
        "hide_tables": ["Executive Financial Report"],
    },
    "operations": {
        "dashboard": True, "assistant": False, "alerts": True, "reports": True,
        "admin": False, "upload": False, "email": False, "documents": True, "ppt_builder": False,
        "export": True, "export_excel": True, "export_pdf": True, "export_ppt": False,
        "pages": ["Dashboard", "Smart Alerts", "Executive Reports", "📤 Document Upload Center"],
        "dashboard_tabs": ["overview", "tables", "performance", "reports"],
        "hide_buttons": ["Delete", "Import"],
        "hide_tables": ["PMO Audit", "Executive Financial Report"],
    },
    "viewer": {
        "dashboard": True, "assistant": False, "alerts": False, "reports": False,
        "admin": False, "upload": False, "email": False, "documents": False, "ppt_builder": False,
        "project_updates": False,
        "export": False, "export_excel": False, "export_pdf": False, "export_ppt": False,
        "pages": ["Dashboard"],
        "dashboard_tabs": ["overview"],
        "hide_buttons": ["Export", "Upload", "Delete", "Import"],
        "hide_tables": ["PMO Audit", "Tables & Exports", "Raw Data"],
    },
}

ROLE_DISPLAY_NAMES = {
    "admin": "Admin",
    "pmo": "PMO",
    "board": "Board",
    "finance": "Finance",
    "audit": "Audit",
    "operations": "Operations",
    "viewer": "Viewer",
}


st.set_page_config(
    page_title="Dawiyat PMO Executive Portal",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed",
)


# Streamlit pages that must stay hidden from sidebar navigation.
# They are opened only through controlled action buttons on the Dashboard.
HIDDEN_ACTION_PAGES = {
    "📤 Document Upload Center", "📊 Executive PPT Builder", "Admin Board",
    "Project Updates Center", "Data Update Agent", "Notification Center 🔔",
    "Executive Daily Digest", "WhatsApp Agent",
}


PORTAL_CSS = """
<style>
.block-container {
    padding-top: 0.6rem;
    padding-left: 0.7rem;
    padding-right: 0.7rem;
    max-width: 100%;
}
header, footer {visibility: hidden;}
/* Best-effort hide Streamlit Cloud developer toolbar/Manage App from the portal UI.
   Note: Streamlit may still show Manage App to app owners/collaborators outside app DOM. */
[data-testid="stToolbar"],
[data-testid="stDecoration"],
[data-testid="stStatusWidget"],
[data-testid="manage-app-button"],
.stDeployButton,
button[title="Manage app"],
a[title="Manage app"],
div[aria-label="Manage app"],
iframe[title="streamlit runtime"],
.viewerBadge_container__1QSob {
    display: none !important;
    visibility: hidden !important;
    pointer-events: none !important;
}

[data-testid="stSidebar"] {
    background: linear-gradient(180deg,#0f2746 0%,#10223a 100%);
}
[data-testid="stSidebar"] * { color: #f8fafc; }
.portal-login {
    max-width: 1360px;
    width: min(98vw, 1360px);
    margin: 1vh auto 0.8vh;
    padding: 0;
    border-radius: 38px;
    background:
        radial-gradient(circle at top left, rgba(20,184,166,.18), transparent 32%),
        radial-gradient(circle at top right, rgba(245,158,11,.18), transparent 34%),
        linear-gradient(135deg,#ffffff 0%,#f8fbff 100%);
    border: 1px solid #d9e3ef;
    box-shadow: 0 32px 85px rgba(15,23,42,.18);
    overflow: hidden;
    text-align: center;
}
.portal-login-top {
    height: 8px;
    background: linear-gradient(90deg,#f97316,#14b8a6,#2563eb);
}
.portal-login-inner {
    padding: 28px 68px 24px;
}
.logo-row {
    display:grid;
    grid-template-columns: 1.15fr .85fr;
    gap:34px;
    align-items:center;
    margin-bottom:18px;
}
.logo-card {
    height:138px;
    border:1px solid #e6edf5;
    border-radius:28px;
    display:flex;
    align-items:center;
    justify-content:center;
    background:#fff;
    box-shadow: 0 12px 28px rgba(15,23,42,.08);
    padding:18px 26px;
}
.logo-card img {
    max-width:100%;
    max-height:116px;
    object-fit:contain;
}
.logo-card.met-logo img {
    max-height:112px;
}
.portal-badge {
    display:inline-flex;
    padding:8px 16px;
    border-radius:999px;
    background:#10223a;
    color:#facc15;
    font-size:12px;
    font-weight:900;
    letter-spacing:.12em;
    text-transform:uppercase;
    margin-bottom:10px;
}
.portal-title {
    font-size: 48px;
    font-weight: 1000;
    color: #10223a;
    margin-bottom: 8px;
    letter-spacing:-.04em;
}
.portal-subtitle {
    color: #526985;
    font-size: 20px;
    line-height: 1.8;
    max-width:720px;
    margin:0 auto;
}
.portal-feature-row {
    display:grid;
    grid-template-columns: repeat(3,1fr);
    gap:12px;
    margin-top:18px;
}
.portal-feature {
    border:1px solid #e6edf5;
    border-radius:18px;
    background:rgba(255,255,255,.75);
    padding:15px 16px;
    color:#10223a;
    font-size:15px;
    font-weight:800;
}

.account-access-row {
    width: 100%;
    display: flex;
    justify-content: center;
    align-items: center;
    margin: 18px auto 18px;
    text-align: center;
}
.account-access-pill {
    display:flex;
    flex-direction:column;
    align-items:center;
    justify-content:center;
    width:min(1100px,95vw);
    margin: 0 auto;
    padding:18px 40px;
    border-radius:999px;
    background:linear-gradient(135deg,#111827,#1e3a8a);
    color:#facc15 !important;
    font-size:20px;
    line-height:1.45;
    font-weight:1000;
    text-align:center;
    box-shadow:0 14px 28px rgba(15,23,42,.24);
    border:1px solid rgba(255,255,255,.18);
}
.account-access-pill .admin-contact-name {
    display:block;
    color:#ffffff !important;
    font-size:26px;
    font-weight:1000;
    letter-spacing:.01em;
    line-height:1.35;
}
.login-form-wrap {
    max-width: 880px;
    margin: 0 auto 18px;
    padding: 18px 20px 20px;
    border: 6px solid #0b0f19;
    border-radius: 4px;
    background: rgba(255,255,255,.84);
    box-shadow: 0 16px 36px rgba(15,23,42,.14);
}
.login-form-wrap .login-title {
    color: #10223a;
    font-size: 18px;
    font-weight: 1000;
    text-align: center;
    margin-bottom: 10px;
}
div[data-testid="stForm"] {
    max-width: 880px;
    margin: 0 auto 8px;
    padding: 18px 20px 20px;
    border: 6px solid #0b0f19;
    border-radius: 4px;
    background: rgba(255,255,255,.86);
    box-shadow: 0 16px 36px rgba(15,23,42,.14);
}
/* Enlarged login form controls */
div[data-testid="stTextInput"] label {
    font-size: 18px !important;
    font-weight: 900 !important;
    color: #10223a !important;
}
div[data-testid="stTextInput"] input {
    min-height: 54px !important;
    font-size: 18px !important;
    border-radius: 14px !important;
}
div[data-testid="stButton"] button {
    min-height: 50px !important;
    font-size: 18px !important;
    font-weight: 900 !important;
    border-radius: 14px !important;
}

.upload-center-hero {
    border:1px solid #bfdbfe;
    border-radius:18px;
    background:linear-gradient(135deg,#eaf4ff 0%,#f8fbff 100%);
    padding:16px 18px;
    margin:8px 0 10px;
    box-shadow:0 10px 22px rgba(37,99,235,.08);
}
.upload-center-hero .uc-title {
    color:#0f3b73;
    font-size:24px;
    font-weight:1000;
    margin-bottom:5px;
    letter-spacing:-.02em;
}
.upload-center-hero .uc-subtitle {
    color:#315b89;
    font-size:14px;
    line-height:1.6;
    font-weight:700;
}

@media(max-width:760px){
    .portal-login { margin:2vh auto; }
    .portal-login-inner { padding:24px 18px; }
    .logo-row { grid-template-columns:1fr; }
    .portal-title { font-size:30px; }
    .portal-subtitle { font-size:14px; }
    .portal-feature-row { grid-template-columns:1fr; }
}
.exec-card {
    border:1px solid #dbe5f1;
    border-radius:20px;
    background:#fff;
    padding:16px;
    box-shadow:0 10px 24px rgba(15,23,42,.06);
}
.exec-card h3 {margin:0 0 8px;color:#10223a;font-weight:900;}
.exec-card .small {color:#64748b;font-size:13px;}
.metric-tile {
    border:1px solid #dbe5f1;
    border-radius:18px;
    background:#fff;
    padding:16px;
    min-height:120px;
    box-shadow:0 8px 22px rgba(15,23,42,.05);
}
.metric-tile .label {
    color:#64748b;
    font-size:11px;
    font-weight:900;
    text-transform:uppercase;
    letter-spacing:.08em;
}
.metric-tile .value {
    color:#10223a;
    font-size:26px;
    font-weight:900;
    margin-top:8px;
}
.priority-high {background:#fee2e2;color:#991b1b;border-radius:999px;padding:5px 10px;font-weight:900;}
.priority-med {background:#fef3c7;color:#92400e;border-radius:999px;padding:5px 10px;font-weight:900;}
.priority-low {background:#dcfce7;color:#166534;border-radius:999px;padding:5px 10px;font-weight:900;}
.arabic-box {
    direction: rtl;
    text-align: right;
    border:1px solid #dbe5f1;
    border-radius:18px;
    padding:18px;
    background:#fff;
    line-height:1.8;
}

.quick-actions-panel {
    border:1px solid #d9e3ef;
    border-radius:18px;
    padding:14px 16px;
    background:#ffffff;
    margin:10px 0 14px;
    box-shadow:0 8px 22px rgba(15,23,42,.05);
}
.quick-actions-title { font-weight:900; color:#10223a; font-size:18px; margin-bottom:4px; }
.quick-actions-subtitle { color:#64748b; font-size:12px; margin-bottom:10px; }
body.dark-ui .quick-actions-panel { background:#111f34 !important; border-color:#2b3d5a !important; }
body.dark-ui .quick-actions-title { color:#eaf2ff !important; }
body.dark-ui .quick-actions-subtitle { color:#9fb0c7 !important; }
</style>
"""
st.markdown(PORTAL_CSS, unsafe_allow_html=True)


def _default_users() -> Dict[str, Dict[str, str]]:
    """Safe fallback users used only when no valid users are configured in Streamlit Secrets."""
    return {
        "ahmedfekry": {"password": "20020099", "role": "admin"},
        "pmo_team": {"password": "PMO12345", "role": "pmo"},
        "board": {"password": "Met_12345", "role": "board"},
        "finance": {"password": "Finance12345", "role": "finance"},
        "viewer": {"password": "Viewer12345", "role": "viewer"},
    }


def _secret_get(data: Any, key: str, default: str = "") -> str:
    """Read from dict / Streamlit AttrDict / object safely."""
    try:
        if isinstance(data, Mapping):
            return str(data.get(key, default))
        if hasattr(data, "get"):
            return str(data.get(key, default))
        return str(getattr(data, key, default))
    except Exception:
        return str(default)



def _session_secret() -> str:
    """Return a stable secret used only for browser-refresh login persistence.

    Optional Streamlit Secrets:
    [session]
    secret = "PUT_RANDOM_LONG_TEXT_HERE"

    If not configured, the app uses a deterministic fallback from the project name.
    For production, adding [session].secret is recommended but not mandatory.
    """
    try:
        raw = st.secrets.get("session", {})
        secret = _secret_get(raw, "secret", "").strip()
        if secret:
            return secret
    except Exception:
        pass
    try:
        email = _secret_get(st.secrets.get("gcp_service_account", {}), "client_email", "")
        if email:
            return f"dawiyat-pmo-session::{email}"
    except Exception:
        pass
    return "dawiyat-pmo-portal-local-session-secret"


def _make_session_signature(username: str, role: str, password: str, issued_at: str) -> str:
    payload = f"{username}|{role}|{password}|{issued_at}".encode("utf-8")
    return hmac.new(_session_secret().encode("utf-8"), payload, hashlib.sha256).hexdigest()


def _format_login_time(timestamp_str: str | None = None) -> str:
    """Return a readable local login time for the user session."""
    try:
        ts = int(timestamp_str) if timestamp_str else int(time.time())
        return datetime.fromtimestamp(ts, KSA_TZ).strftime("%d-%b-%Y %I:%M %p")
    except Exception:
        return ksa_now().strftime("%d-%b-%Y %I:%M %p")



def _clear_login_query_params() -> None:
    """Remove persisted login parameters from the URL without touching other app state."""
    try:
        for k in ["auth_user", "auth_role", "auth_iat", "auth_sig"]:
            if k in st.query_params:
                del st.query_params[k]
    except Exception:
        pass


def _persist_login(username: str, role: str, password: str, issued_at: str | None = None) -> None:
    """Persist authenticated login across browser refreshes.

    This prevents pressing the browser Refresh button from sending the user back
    to the login screen. Logout explicitly removes these parameters.
    """
    issued_at = issued_at or str(int(time.time()))
    signature = _make_session_signature(username, role, password, issued_at)
    try:
        st.query_params["auth_user"] = username
        st.query_params["auth_role"] = role
        st.query_params["auth_iat"] = issued_at
        st.query_params["auth_sig"] = signature
    except Exception:
        pass


def _restore_login_from_query_params() -> bool:
    """Restore login state after a browser refresh, if the signed URL session is valid."""
    if st.session_state.get("authenticated"):
        return True
    try:
        username = str(st.query_params.get("auth_user", "")).strip()
        role = str(st.query_params.get("auth_role", "viewer")).strip().lower()
        issued_at = str(st.query_params.get("auth_iat", "")).strip()
        signature = str(st.query_params.get("auth_sig", "")).strip()
    except Exception:
        return False

    if not username or not issued_at or not signature:
        return False

    # Optional safety window: keep browser-refresh sessions valid for 12 hours.
    try:
        if int(time.time()) - int(issued_at) > 12 * 60 * 60:
            _clear_login_query_params()
            return False
    except Exception:
        _clear_login_query_params()
        return False

    users = get_users()
    if username not in users:
        _clear_login_query_params()
        return False

    expected_role = str(users[username].get("role", role)).lower()
    password = str(users[username].get("password", ""))
    expected = _make_session_signature(username, expected_role, password, issued_at)

    if hmac.compare_digest(signature, expected):
        st.session_state["authenticated"] = True
        st.session_state["username"] = username
        st.session_state["role"] = expected_role or "user"
        st.session_state["last_login"] = _format_login_time(issued_at)
        return True

    _clear_login_query_params()
    return False



def _excel_bool(value, default: bool = False) -> bool:
    text = str(value).strip().lower()
    if text in ["yes", "y", "true", "1", "on"]:
        return True
    if text in ["no", "n", "false", "0", "off"]:
        return False
    return default


def _excel_clean(value) -> str:
    if value is None:
        return ""
    text = str(value).strip()
    if text.lower() in ["nan", "none", "null"]:
        return ""
    return text


def _normalize_permission_sheet(raw_df: pd.DataFrame, required_cols: List[str]) -> pd.DataFrame:
    """Normalize Excel sheets that have a title row before the real header."""
    if raw_df is None or raw_df.empty:
        return pd.DataFrame(columns=required_cols)
    raw = raw_df.fillna("").copy()
    required_lower = {str(c).strip().lower() for c in required_cols}
    header_idx = None
    for idx, row in raw.iterrows():
        lower_values = {str(v).strip().lower() for v in row.tolist() if str(v).strip()}
        if required_lower.issubset(lower_values):
            header_idx = idx
            break
    if header_idx is None:
        out = raw.copy()
        out.columns = [str(c).strip() for c in out.columns]
        return out
    headers = [str(v).strip() if str(v).strip() else f"Unnamed_{i}" for i, v in enumerate(raw.loc[header_idx].tolist())]
    out = raw.iloc[header_idx + 1:].copy()
    out.columns = headers
    out = out.reset_index(drop=True)
    if not out.empty:
        out = out.loc[~out.apply(lambda r: all(str(x).strip() == "" for x in r.tolist()), axis=1)].reset_index(drop=True)
    return out



def _permissions_mtime() -> float:
    """Return mtime for permissions.xlsx; used to detect GitHub/Admin uploads."""
    try:
        return PERMISSIONS_XLSX_PATH.stat().st_mtime if PERMISSIONS_XLSX_PATH.exists() else 0.0
    except Exception:
        return 0.0


def _permissions_signature() -> str:
    """Stable signature for the currently mounted permission workbook."""
    try:
        if not PERMISSIONS_XLSX_PATH.exists():
            return "missing"
        stat = PERMISSIONS_XLSX_PATH.stat()
        return f"{int(stat.st_mtime)}:{stat.st_size}"
    except Exception:
        return str(time.time())


def _permissions_last_modified_text() -> str:
    try:
        if not PERMISSIONS_XLSX_PATH.exists():
            return "Missing"
        return datetime.fromtimestamp(PERMISSIONS_XLSX_PATH.stat().st_mtime).strftime("%Y-%m-%d %H:%M:%S")
    except Exception:
        return "Unknown"


def _read_permissions_excel() -> Dict[str, pd.DataFrame]:
    """Read Excel permission matrix from data/permissions.xlsx.

    V43.3: intentionally reads from disk every rerun and uses file mtime/signature
    to avoid stale permission behavior after GitHub updates or Admin uploads.
    Robust against title rows / merged headers in the workbook template.
    """
    if not PERMISSIONS_XLSX_PATH.exists():
        return {}
    try:
        raw_sheets = pd.read_excel(PERMISSIONS_XLSX_PATH, sheet_name=None, dtype=str, header=None).copy()
        specs = {
            "Users": ["Username", "Password"],
            "User_Page_Access": ["Username", "Executive Overview"],
            "User_Component_Access": ["Username", "Page", "Component / Table", "Show"],
        }
        sheets: Dict[str, pd.DataFrame] = {}
        for name, df in raw_sheets.items():
            sheets[name] = _normalize_permission_sheet(df, specs.get(name, [])) if name in specs else df.fillna("")
        st.session_state["permissions_last_signature"] = _permissions_signature()
        st.session_state["permissions_last_loaded"] = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
        return sheets
    except Exception as exc:
        st.warning(f"Unable to read permissions.xlsx. Falling back to Secrets/default permissions. Details: {exc}")
        return {}



def _write_permissions_excel_sheets(sheets: Dict[str, pd.DataFrame]) -> None:
    """Persist permission sheets to data/permissions.xlsx so Admin Board edits become live immediately.
    The app reads this workbook on each normal rerun; users pick up changes after browser refresh or Logout/Login.
    """
    DATA_DIR.mkdir(exist_ok=True)
    clean_sheets: Dict[str, pd.DataFrame] = {}
    for name, df in (sheets or {}).items():
        if isinstance(df, pd.DataFrame):
            out = df.copy().fillna("")
            out.columns = [str(c).strip() for c in out.columns]
            clean_sheets[str(name)] = out
    if not clean_sheets:
        clean_sheets = {
            "Users": pd.DataFrame(columns=["Username", "Password", "Department / Display Role", "Active", "Full Name / Department"]),
            "User_Page_Access": pd.DataFrame(columns=["Username"]),
            "User_Component_Access": pd.DataFrame(columns=["Username", "Page", "Component / Table", "Show", "Export Excel", "Export PDF", "Export PPT", "Notes"]),
        }
    with pd.ExcelWriter(PERMISSIONS_XLSX_PATH, engine="openpyxl") as writer:
        for name, df in clean_sheets.items():
            safe_name = str(name)[:31] or "Sheet1"
            df.to_excel(writer, sheet_name=safe_name, index=False)
    st.session_state["permission_runtime_signature"] = _permissions_signature()
    st.session_state["permissions_last_loaded"] = ksa_now().strftime("%Y-%m-%d %H:%M:%S")


def _update_permission_sheet_from_editor(sheet_name: str, edited_df: pd.DataFrame, selected_user: str = "All Users") -> None:
    sheets = _read_permissions_excel()
    edited = edited_df.copy().fillna("")
    if selected_user and selected_user != "All Users" and sheet_name in sheets and "Username" in edited.columns:
        original = sheets[sheet_name].copy().fillna("")
        if "Username" in original.columns:
            keep = original[original["Username"].astype(str).str.strip().str.lower() != selected_user.strip().lower()]
            sheets[sheet_name] = pd.concat([keep, edited], ignore_index=True)
        else:
            sheets[sheet_name] = edited
    else:
        sheets[sheet_name] = edited
    _write_permissions_excel_sheets(sheets)


def render_permission_auto_refresh(interval_seconds: int = 12) -> None:
    """Disabled by design.

    Permissions are reloaded only on a normal Streamlit rerun, browser refresh,
    Logout/Login, or the manual Reload Permissions button. This avoids forcing
    page reloads while users are working.
    """
    return

def _excel_permission_available() -> bool:
    return PERMISSIONS_XLSX_PATH.exists()


def get_excel_user_records() -> Tuple[Dict[str, Dict[str, str]], List[str]]:
    """V43: Return active users from permissions.xlsx.
    The Role/Department column is display-only and is no longer used for permissions.
    Active=No removes the same username from Secrets/fallback users.
    """
    sheets = _read_permissions_excel()
    if "Users" not in sheets:
        return {}, []
    df = sheets["Users"].fillna("")
    required = {"Username", "Password"}
    if not required.issubset(set(df.columns)):
        return {}, []

    active_users: Dict[str, Dict[str, str]] = {}
    inactive_users: List[str] = []
    for _, row in df.iterrows():
        username = _excel_clean(row.get("Username"))
        password = _excel_clean(row.get("Password"))
        role = _excel_clean(row.get("Department / Display Role")) or _excel_clean(row.get("Department")) or _excel_clean(row.get("Role")) or "user"
        role_key = str(role).lower()
        active = _excel_clean(row.get("Active")).lower() or "yes"
        if not username:
            continue
        if active in ["no", "false", "0", "inactive", "disabled"]:
            inactive_users.append(username)
            continue
        if not password:
            continue
        # Keep the value for display only. If it is not a built-in role, current_role will safely fallback.
        active_users[username] = {"password": password, "role": role_key, "department": role, "full_name": _excel_clean(row.get("Full Name / Department")) or role}
    return active_users, inactive_users


def get_excel_users() -> Dict[str, Dict[str, str]]:
    users, _inactive = get_excel_user_records()
    return users


def _page_name_to_key(page_name: str) -> str:
    text = _excel_clean(page_name).lower()
    mapping = {
        "executive overview": "dashboard",
        "tables & exports": "dashboard",
        "pmo audit": "dashboard",
        "kpi performance": "dashboard",
        "performance explanation": "dashboard",
        "pmo report assistant": "dashboard",
        "executive reports": "reports",
        "executive ppt builder": "ppt_builder",
        "document upload center": "documents",
        "upload csv": "upload",
        "smart alerts": "alerts",
        "ai executive assistant": "assistant",
        "admin board": "admin",
        "project updates center": "project_updates",
        "data update agent": "data_update_agent",
        "notification center": "notification_center",
        "notification center 🔔": "notification_center",
        "executive daily digest": "executive_daily_digest",
        "whatsapp agent": "whatsapp_agent",
    }
    return mapping.get(text, "")


def _page_name_to_dashboard_tab(page_name: str) -> str:
    text = _excel_clean(page_name).lower()
    mapping = {
        "executive overview": "overview",
        "tables & exports": "tables",
        "pmo audit": "pmo",
        "kpi performance": "performance",
        "performance explanation": "perf-explanation",
        "pmo report assistant": "decision",
        "executive reports": "reports",
    }
    return mapping.get(text, "")


def _page_display_name(page_name: str) -> str:
    text = _excel_clean(page_name)
    if text == "Executive PPT Builder":
        return "📊 Executive PPT Builder"
    if text == "Document Upload Center":
        return "📤 Document Upload Center"
    return text



def get_user_based_policy_from_excel(username: str) -> Dict:
    """V43 User-Based Permission Engine.
    Permissions are taken from:
      - Users: login + Active only
      - User_Page_Access: page/tab access by Username
      - User_Component_Access: component/table visibility and export buttons by Username
    Role sheets are ignored for permissions.
    """
    sheets = _read_permissions_excel()
    username = str(username or "").strip()
    if not username:
        return {}

    policy: Dict[str, Any] = {
        "dashboard": False,
        "assistant": False,
        "alerts": False,
        "reports": False,
        "admin": False,
        "upload": False,
        "email": False,
        "documents": False,
        "ppt_builder": False,
        "project_updates": False,
        "export": False,
        "export_excel": False,
        "export_pdf": False,
        "export_ppt": False,
        "pages": [],
        "dashboard_tabs": [],
        "hide_buttons": [],
        "hide_tables": [],
        "show_tables": [],
        "hide_excel_components": [],
        "allowed_excel_components": [],
        "hide_pdf_components": [],
        "allowed_pdf_components": [],
        "hide_ppt_components": [],
    }

    pages: List[str] = []
    tabs: List[str] = []

    # Page access by username
    if "User_Page_Access" in sheets:
        df = sheets["User_Page_Access"].fillna("")
        if "Username" in df.columns:
            rows = df[df["Username"].astype(str).str.strip().str.lower() == username.lower()]
            if not rows.empty:
                row = rows.iloc[0]
                for col in df.columns:
                    if col == "Username":
                        continue
                    enabled = _excel_bool(row.get(col), False)
                    display = _page_display_name(col)
                    tab = _page_name_to_dashboard_tab(col)
                    key = _page_name_to_key(col)
                    if enabled:
                        if display and display not in pages:
                            pages.append(display)
                        if tab and tab not in tabs:
                            tabs.append(tab)
                        if key:
                            policy[key] = True
                    else:
                        if key:
                            policy[key] = False

    # If any dashboard tab is enabled, include Dashboard container page.
    if tabs and "Dashboard" not in pages:
        pages.insert(0, "Dashboard")

    # Component/table access by username
    show_tables: List[str] = []
    hide_tables: List[str] = []
    hide_excel_components: List[str] = []
    hide_pdf_components: List[str] = []
    allowed_pdf_components: List[str] = []
    hide_ppt_components: List[str] = []
    any_excel = False
    any_pdf = False
    any_ppt = False
    global_pdf_report = False
    global_pdf_names = {
        "global pdf report",
        "export pdf report",
        "dashboard pdf report",
        "full dashboard pdf report",
        "dashboard full pdf report",
    }

    if "User_Component_Access" in sheets:
        df = sheets["User_Component_Access"].fillna("")
        if "Username" in df.columns:
            # V43.3 security rule: default deny.
            # Any component listed in the workbook but not explicitly Show=Yes for this username is hidden.
            all_components = [
                _excel_clean(x) for x in df.get("Component / Table", pd.Series(dtype=str)).tolist()
                if _excel_clean(x)
            ]
            rows = df[df["Username"].astype(str).str.strip().str.lower() == username.lower()]
            for _, row in rows.iterrows():
                component = _excel_clean(row.get("Component / Table"))
                if not component:
                    continue
                show = _excel_bool(row.get("Show"), False)
                ex_excel = _excel_bool(row.get("Export Excel"), False)
                ex_pdf = _excel_bool(row.get("Export PDF"), False)
                ex_ppt = _excel_bool(row.get("Export PPT"), False)

                # Global PDF Report is a permission for the TOP dashboard PDF button only.
                # It must NOT be counted as a normal table/component and must NOT grant
                # permission to export all tables.
                if component.strip().lower() in global_pdf_names:
                    global_pdf_report = bool(show or ex_pdf)
                    continue

                if show:
                    show_tables.append(component)
                    if ex_excel:
                        any_excel = True
                    else:
                        hide_excel_components.append(component)
                    if ex_pdf:
                        any_pdf = True
                        allowed_pdf_components.append(component)
                    else:
                        hide_pdf_components.append(component)
                    if ex_ppt:
                        any_ppt = True
                    else:
                        hide_ppt_components.append(component)
                else:
                    hide_tables.append(component)
                    hide_excel_components.append(component)
                    hide_pdf_components.append(component)
                    hide_ppt_components.append(component)

            # Hide every workbook-known component that is not explicitly allowed.
            for component in all_components:
                if component and component not in show_tables and component not in hide_tables:
                    hide_tables.append(component)
                    hide_excel_components.append(component)
                    hide_pdf_components.append(component)
                    hide_ppt_components.append(component)

    # Same component names can exist in more than one page. If the component is
    # explicitly Show=Yes / Export=Yes anywhere for this user, that positive
    # permission must win. This fixes cases like Executive Reports=Yes while the
    # same component name is No in Executive Overview, which previously hid it.
    show_set = {str(x).strip().lower() for x in show_tables}
    pdf_set = {str(x).strip().lower() for x in allowed_pdf_components}
    excel_allowed_components = [x for x in show_tables if str(x).strip().lower() not in {str(h).strip().lower() for h in hide_excel_components}]
    excel_allowed_set = {str(x).strip().lower() for x in excel_allowed_components}
    ppt_allowed_set = {str(x).strip().lower() for x in show_tables if str(x).strip().lower() not in {str(h).strip().lower() for h in hide_ppt_components}}
    hide_tables = [x for x in hide_tables if str(x).strip().lower() not in show_set]
    hide_pdf_components = [x for x in hide_pdf_components if str(x).strip().lower() not in pdf_set]
    hide_excel_components = [x for x in hide_excel_components if str(x).strip().lower() not in excel_allowed_set]
    hide_ppt_components = [x for x in hide_ppt_components if str(x).strip().lower() not in ppt_allowed_set]

    # Canonical order and deduplication
    canonical_pages = ["Dashboard", "Project Updates Center", "Data Update Agent", "Notification Center 🔔", "Executive Daily Digest", "WhatsApp Agent", "AI Executive Assistant", "Smart Alerts", "Executive Reports", "📊 Executive PPT Builder", "Upload CSV", "📤 Document Upload Center", "Admin Board"]
    pages = [p for p in canonical_pages if p in list(dict.fromkeys(pages))]
    canonical_tabs = ["overview", "tables", "pmo", "performance", "perf-explanation", "decision", "reports"]
    tabs = [t for t in canonical_tabs if t in list(dict.fromkeys(tabs))]
    policy["pages"] = pages
    policy["dashboard_tabs"] = tabs
    policy["show_tables"] = show_tables
    policy["hide_tables"] = hide_tables
    policy["hide_excel_components"] = list(dict.fromkeys(hide_excel_components))
    policy["allowed_excel_components"] = list(dict.fromkeys(excel_allowed_components))
    policy["hide_pdf_components"] = list(dict.fromkeys(hide_pdf_components))
    policy["allowed_pdf_components"] = list(dict.fromkeys(allowed_pdf_components))
    policy["hide_ppt_components"] = list(dict.fromkeys(hide_ppt_components))
    policy["export_excel"] = any_excel
    policy["export_pdf"] = any_pdf
    # Dedicated control for the top dashboard button: Export PDF Report.
    # This is intentionally separated from table-level Export PDF permissions.
    policy["global_pdf_report"] = bool(global_pdf_report)
    policy["export_ppt"] = any_ppt  # Export permission only; does NOT grant PPT Builder page access
    policy["export"] = bool(policy.get("export_excel") or policy.get("export_pdf") or policy.get("export_ppt") or policy.get("global_pdf_report"))
    return policy


def get_excel_role_policy(role: str) -> Dict:
    """Disabled: the portal uses User-Based Permissions Only."""
    return {}


def get_excel_user_override(username: str) -> Dict:
    """Disabled: all exceptions must be entered directly per user in User_Page_Access/User_Component_Access."""
    return {}


def get_users() -> Dict[str, Dict[str, str]]:
    """Stable user loader with Excel permissions as the final authority.
    Priority: fallback -> Secrets -> Excel. Excel Active=No removes user.
    """
    users = _default_users()
    try:
        raw_users = st.secrets.get("users", {})
    except Exception:
        raw_users = {}
    if raw_users:
        try:
            items = raw_users.items() if hasattr(raw_users, "items") else []
            for raw_username, data in items:
                username = str(raw_username).strip()
                if not username:
                    continue
                if isinstance(data, str):
                    password = data.strip()
                    role = "viewer"
                else:
                    password = _secret_get(data, "password", "").strip()
                    role = _secret_get(data, "role", "viewer").strip().lower()
                if not password:
                    continue
                users[username] = {"password": password, "role": role or "user"}
        except Exception:
            pass
    excel_users, inactive_excel_users = get_excel_user_records()
    users.update(excel_users)
    for inactive in inactive_excel_users:
        users.pop(inactive, None)
    return users


def _as_bool(value, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    return str(value).strip().lower() in ["1", "true", "yes", "y", "on"]

def _as_list(value) -> List[str]:
    if value is None:
        return []
    if isinstance(value, (list, tuple, set)):
        return [str(x).strip() for x in value if str(x).strip()]
    text = str(value).strip()
    if not text:
        return []
    return [x.strip() for x in text.split(",") if x.strip()]

def _secret_to_plain_dict(obj) -> Dict:
    out = {}
    try:
        items = obj.items() if hasattr(obj, "items") else []
        for k, v in items:
            key = str(k)
            if hasattr(v, "items"):
                out[key] = _secret_to_plain_dict(v)
            elif isinstance(v, (list, tuple)):
                out[key] = [str(x) if not isinstance(x, bool) else x for x in v]
            else:
                out[key] = v
    except Exception:
        return {}
    return out

def _merge_policy(base: Dict, override: Dict) -> Dict:
    merged = dict(base or {})
    for key, value in (override or {}).items():
        if key in ["password", "role"]:
            continue
        if key in ["pages", "dashboard_tabs", "hide_buttons", "hide_tables", "show_tables", "components"]:
            merged[key] = _as_list(value)
        elif isinstance(value, bool):
            merged[key] = value
        elif str(value).strip().lower() in ["true", "false", "yes", "no", "1", "0", "on", "off"]:
            merged[key] = _as_bool(value)
        else:
            merged[key] = value
    return merged

def get_roles_override_from_secrets(role: str) -> Dict:
    """Disabled: roles in Secrets are not used for permissions."""
    return {}

def get_user_override_from_secrets(username: str) -> Dict:
    """Disabled by design.

    V3.4 uses permissions.xlsx as the ONLY authority for page, tab, component,
    and export permissions. Streamlit Secrets may still contain old role/page
    values from V41/V42, but those values must not override Excel.
    This fixes the issue where Active=Yes/No changed correctly while page and
    component permission changes from GitHub did not appear on the dashboard.
    """
    return {}

def user_policy(username: str | None = None, role: str | None = None) -> Dict:
    """V43: user-only permission policy.
    Role is used for display only. No role-based page/component permissions are applied.
    """
    username = username or str(st.session_state.get("username", "")).strip()
    if username:
        policy = get_user_based_policy_from_excel(username)

        # V3 hard rule requested by Ahmed:
        # Admin Board is user-based and visible only for username ahmedfekry,
        # regardless of any legacy Role/Department value or old workbook cache.
        if _is_admin_board_owner(username):
            policy["admin"] = True
            pages = _as_list(policy.get("pages"))
            if "Admin Board" not in pages:
                pages.append("Admin Board")
            policy["pages"] = pages
        else:
            policy["admin"] = False
            policy["pages"] = [p for p in _as_list(policy.get("pages")) if p != "Admin Board"]

        # Optional emergency user-specific Secret override only.
        policy = _merge_policy(policy, get_user_override_from_secrets(username))

        # Re-apply the hard gate after any emergency override.
        if not _is_admin_board_owner(username):
            policy["admin"] = False
            policy["pages"] = [p for p in _as_list(policy.get("pages")) if p != "Admin Board"]
        elif "Admin Board" not in _as_list(policy.get("pages")):
            policy["pages"] = _as_list(policy.get("pages")) + ["Admin Board"]

        return policy

    # Before login, keep the minimum safe policy.
    return {
        "dashboard": True,
        "pages": ["Dashboard"],
        "dashboard_tabs": ["overview"],
        "export": False,
        "export_excel": False,
        "export_pdf": False,
        "export_ppt": False,
        "hide_buttons": ["Export", "Upload", "Delete", "Import"],
        "hide_tables": [],
    }


def allowed_pages_for_current_user() -> List[str]:
    policy = user_policy()
    pages = _as_list(policy.get("pages"))
    if pages:
        return pages
    out = []
    if policy.get("dashboard"): out.append("Dashboard")
    if policy.get("assistant"): out.append("AI Executive Assistant")
    if policy.get("alerts"): out.append("Smart Alerts")
    if policy.get("reports"): out.append("Executive Reports")
    if policy.get("project_updates"): out.append("Project Updates Center")
    if policy.get("data_update_agent") or policy.get("admin"): out.append("Data Update Agent")
    if policy.get("notification_center"): out.append("Notification Center 🔔")
    if policy.get("executive_daily_digest"): out.append("Executive Daily Digest")
    if policy.get("whatsapp_agent"): out.append("WhatsApp Agent")
    if policy.get("ppt_builder"): out.append("📊 Executive PPT Builder")
    if policy.get("upload"): out.append("Upload CSV")
    if policy.get("documents"): out.append("📤 Document Upload Center")
    if policy.get("admin"): out.append("Admin Board")
    return out

def current_role() -> str:
    # V43: Role/Department is display-only.
    return str(st.session_state.get("role", "user")).lower() or "user"


def role_policy(role: str | None = None) -> Dict:
    # Legacy compatibility only. Permissions are user-based in V43.
    return user_policy()


def can(permission: str) -> bool:
    return bool(user_policy().get(permission, False))


def _is_admin_board_owner(username: str | None = None) -> bool:
    """Hard security gate: Admin Board is visible only for ahmedfekry."""
    username = username or str(st.session_state.get("username", "")).strip()
    return str(username or "").strip().lower() == "ahmedfekry"



def validate_current_authenticated_user() -> bool:
    """Ensure current session still exists in permissions.xlsx after file updates.

    If Active becomes No or the user is removed, force logout immediately.
    If Department/Password changed, update the session and signed URL.
    """
    if not st.session_state.get("authenticated"):
        return False
    username = str(st.session_state.get("username", "")).strip()
    if not username:
        _clear_login_query_params()
        st.session_state.clear()
        return False

    current_sig = _permissions_signature()
    if st.session_state.get("permission_runtime_signature") != current_sig:
        st.session_state["permission_runtime_signature"] = current_sig
        st.session_state["permission_runtime_checked_at"] = ksa_now().strftime("%Y-%m-%d %H:%M:%S")

    users = get_users()
    if username not in users:
        _clear_login_query_params()
        st.session_state.clear()
        st.error("Your account access has been changed or disabled. Please contact the PMO System Administrator.")
        st.rerun()
        return False

    new_role = str(users[username].get("role", "user")).lower() or "user"
    if st.session_state.get("role") != new_role:
        st.session_state["role"] = new_role
        try:
            _persist_login(username, new_role, str(users[username].get("password", "")))
        except Exception:
            pass
    return True


def allowed_dashboard_tabs(role: str | None = None) -> List[str]:
    policy = user_policy(role=role) if role else user_policy()
    tabs = _as_list(policy.get("dashboard_tabs"))
    return tabs


def login_page() -> bool:
    if st.session_state.get("authenticated") or _restore_login_from_query_params():
        return True

    dawiyat_logo = image_to_base64(DAWIYAT_LOGO_PATH)
    met_logo = image_to_base64(MET_LOGO_PATH)

    st.markdown(
        f"""
        <div class="portal-login">
            <div class="portal-login-top"></div>
            <div class="portal-login-inner">
                <div class="logo-row">
                    <div class="logo-card dawiyat-logo">
                        <img src="data:image/jpeg;base64,{dawiyat_logo}" alt="Dawiyat Logo" />
                    </div>
                    <div class="logo-card met-logo">
                        <img src="data:image/jpeg;base64,{met_logo}" alt="Middle Sea Telecom Logo" />
                    </div>
                </div>
                <div class="portal-badge">Executive PMO Governance Portal</div>
                <div class="portal-title">Dawiyat PMO Executive Portal</div>
                <div class="portal-subtitle">
                    Secure executive portal for PMO Dashboard, Arabic AI Assistant, Smart Alerts,
                    PDF Reports, CSV Governance, and Executive Decision Support.
                </div>
                <div class="portal-feature-row">
                    <div class="portal-feature">🔐 User-Based Access</div>
                    <div class="portal-feature">📊 Live PMO Dashboard</div>
                    <div class="portal-feature">🤖 Arabic AI Assistant</div>
                </div>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    c1, c2, c3 = st.columns([0.42, 2.15, 0.42])
    with c2:
        st.markdown(
            """
            <div class="account-access-row">
                <div class="account-access-pill">
                    <div>For account access, please contact</div>
                    <span class="admin-contact-name">Eng./Ahmed Fekry (PMO System Administrator)</span>
                </div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        with st.form("login_form", clear_on_submit=False):
            username = st.text_input("Username")
            password = st.text_input("Password", type="password")
            login = st.form_submit_button("Login", use_container_width=True)

        if login:
            users = get_users()
            if username in users and password == users[username]["password"]:
                role = users[username].get("role", "viewer").lower()
                role = role or "user"
                issued_at = str(int(time.time()))
                st.session_state["authenticated"] = True
                st.session_state["username"] = username
                st.session_state["role"] = role
                st.session_state["last_login"] = _format_login_time(issued_at)
                _persist_login(username, role, users[username]["password"], issued_at=issued_at)
                st.rerun()
            else:
                st.error("Invalid username or password.")

    return False


@st.cache_data(show_spinner=False)
def read_csv_cached(path_str: str, mtime: float) -> pd.DataFrame:
    path = Path(path_str)
    if not path.exists() or path.stat().st_size == 0:
        return pd.DataFrame()
    encodings = ["utf-8-sig", "utf-8", "cp1256", "latin1"]
    last_error = None
    for enc in encodings:
        try:
            return pd.read_csv(path, dtype=str, keep_default_na=False, encoding=enc)
        except Exception as exc:
            last_error = exc
    st.warning(f"Could not read {path.name}: {last_error}")
    return pd.DataFrame()


def safe_read_csv(path: Path) -> pd.DataFrame:
    mtime = path.stat().st_mtime if path.exists() else 0
    return read_csv_cached(str(path), mtime)


@st.cache_data(show_spinner=False)
def load_master_workorders(path_str: str, mtime: float) -> pd.DataFrame:
    """Single source-of-truth loader for u_osp_work_order.csv.

    V4 performance rule:
    - Read the workorder file once through Streamlit cache.
    - Strip accidental spaces from column names.
    - Keep City/District inside the master file when available, reducing dependence on District.csv.
    """
    df = read_csv_cached(path_str, mtime).copy()
    if df.empty:
        return df
    df.columns = [str(c).strip() for c in df.columns]
    for col in ["Region", "City", "District", "Validation", "PM last Update Date"]:
        if col not in df.columns:
            df[col] = ""
    return df.fillna("").astype(str)


def master_workorders_df() -> pd.DataFrame:
    mtime = WO_PATH.stat().st_mtime if WO_PATH.exists() else 0
    return load_master_workorders(str(WO_PATH), mtime)


def derive_district_records_from_workorders(df: pd.DataFrame) -> pd.DataFrame:
    """Create lightweight location records from u_osp_work_order.csv.
    This replaces District.csv as the dashboard location source when City/District exist in the master file.
    """
    needed = ["Link Code", "Work Order", "Region", "City", "District"]
    if df.empty or not all(c in df.columns for c in needed):
        return pd.DataFrame(columns=needed)
    out = df[needed].copy()
    out = out[(out["Link Code"].astype(str).str.strip() != "") | (out["Work Order"].astype(str).str.strip() != "")]
    return out.drop_duplicates().fillna("").astype(str)


def df_to_records(df: pd.DataFrame) -> List[dict]:
    if df.empty:
        return []
    return df.fillna("").astype(str).to_dict(orient="records")





# -----------------------------
# V5.9 Governance / PMO Intelligence helpers
# -----------------------------
def _normalize_text(v: Any) -> str:
    return str(v if v is not None else "").strip()


def _normalize_status(v: Any) -> str:
    return _normalize_text(v).lower().replace("_", " ").replace("-", " ").strip()


POSITIVE_STATUSES = {
    "approved", "accepted", "submitted", "completed", "closed", "created", "done", "issued",
    "ready", "rfs approved", "under clearing remarks", "scheduled",
}
NEGATIVE_STATUSES = {
    "rejected", "failed", "not approved", "returned", "missing documents", "not started", "not start",
    "sor not create", "not create", "cancelled", "canceled",
}
NEUTRAL_STATUSES = {"in progress", "under preparation", "under review", "requested", "scheduled"}


def classify_change_impact(field: str, old_value: Any, new_value: Any) -> str:
    """Simple rule engine used by Change Detection Agent and Digest."""
    old_s = _normalize_status(old_value)
    new_s = _normalize_status(new_value)
    if not new_s or old_s == new_s:
        return "Neutral"
    if new_s in POSITIVE_STATUSES and old_s not in POSITIVE_STATUSES:
        return "Positive"
    if new_s in NEGATIVE_STATUSES and old_s not in NEGATIVE_STATUSES:
        return "Negative"
    if new_s in NEUTRAL_STATUSES:
        return "Neutral"
    if any(x in new_s for x in ["approved", "accepted", "submitted", "completed", "closed", "created"]):
        return "Positive"
    if any(x in new_s for x in ["reject", "fail", "missing", "not start", "not create"]):
        return "Negative"
    return "Neutral"


def category_for_change(field: str, impact: str = "") -> str:
    f = _normalize_status(field)
    if "invoice" in f or "sor" in f or "amount" in f:
        return "Financial"
    if "handover" in f:
        return "Handover"
    if "pat" in f:
        return "PAT"
    if "asbuilt" in f or "as built" in f or "rfs" in f or "redline" in f or "boq" in f or "dcr" in f:
        return "Documents"
    if "civil" in f or "fiber" in f or "full wo" in f:
        return "Implementation"
    if str(impact).lower() == "negative":
        return "Critical"
    return "Project Update"


def notification_access_enabled(username: str, category: str) -> bool:
    """Admin-controlled notification distribution matrix.
    Supports rows like Username=ahmedfekry, Category=ALL, Enabled=Yes.
    """
    ensure_governance_files()
    df = safe_read_csv(NOTIFICATION_ACCESS_PATH)
    if df.empty or "Username" not in df.columns:
        return True
    u = str(username).strip().lower()
    cat = str(category).strip().lower()
    rows = df[df["Username"].astype(str).str.strip().str.lower().isin([u, "all", "*"])].copy()
    if rows.empty:
        return True
    if "Enabled" not in rows.columns:
        return True
    if "Category" not in rows.columns:
        return rows["Enabled"].astype(str).str.lower().isin(["yes", "true", "1", "on"]).any()
    exact = rows[rows["Category"].astype(str).str.strip().str.lower().isin([cat, "all", "*"])]
    if exact.empty:
        return False
    return exact["Enabled"].astype(str).str.lower().isin(["yes", "true", "1", "on"]).any()


def build_effective_operational_data() -> pd.DataFrame:
    """Single source of truth used for snapshots and downloads: Master + latest project updates + derived fields."""
    base = master_workorders_df().copy()
    if base.empty:
        return base
    return apply_derived_billing_fields(apply_project_updates_to_workorders(base)).copy()


def write_master_operational_snapshot(reason: str = "manual") -> Path | None:
    """Persist the current effective master and a dated snapshot for version control.
    Note: on Streamlit Cloud this persists while the app filesystem persists; download snapshots regularly or connect a GitHub write token for permanent storage.
    """
    ensure_governance_files()
    SNAPSHOT_DIR.mkdir(parents=True, exist_ok=True)
    df = build_effective_operational_data()
    if df.empty:
        return None
    df.to_csv(MASTER_OPERATIONAL_PATH, index=False, encoding="utf-8-sig")
    stamp = ksa_now().strftime("%Y%m%d_%H%M%S")
    snap = SNAPSHOT_DIR / f"master_operational_data_{stamp}_{reason}.csv"
    df.to_csv(snap, index=False, encoding="utf-8-sig")
    return snap


def append_change_impacts(change_rows: List[dict]) -> None:
    if not change_rows:
        return
    now = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
    impact_rows = []
    kpi_rows = []
    for ch in change_rows:
        impact = classify_change_impact(ch.get("Field", ""), ch.get("Old Value", ""), ch.get("New Value", ""))
        category = category_for_change(ch.get("Field", ""), impact)
        impact_rows.append({
            "Impact ID": hashlib.md5(f"{ch.get('Change ID','')}-{impact}-{now}".encode()).hexdigest()[:12],
            "Created At": now,
            "Change ID": ch.get("Change ID", ""),
            "Updated By": ch.get("Updated By", ""),
            "Link Code": ch.get("Link Code", ""),
            "Work Order": ch.get("Work Order", ""),
            "Field": ch.get("Field", ""),
            "Old Value": ch.get("Old Value", ""),
            "New Value": ch.get("New Value", ""),
            "Impact": impact,
            "Category": category,
            "Impact Note": f"{category} change classified as {impact}.",
        })
        if impact in ["Positive", "Negative"]:
            kpi_rows.append({
                "KPI Impact ID": hashlib.md5(f"kpi-{ch.get('Change ID','')}-{now}".encode()).hexdigest()[:12],
                "Created At": now,
                "Updated By": ch.get("Updated By", ""),
                "Link Code": ch.get("Link Code", ""),
                "Work Order": ch.get("Work Order", ""),
                "Field": ch.get("Field", ""),
                "Impact": impact,
                "KPI Area": category,
                "Message": f"{category} KPI impact: {ch.get('Field','')} changed from {ch.get('Old Value','')} to {ch.get('New Value','')}.",
            })
    _append_csv_rows(CHANGE_IMPACT_PATH, impact_rows)
    _append_csv_rows(KPI_IMPACT_PATH, kpi_rows)


def create_digest_for_date(date_text: str | None = None, audience: str = "Executive") -> str:
    ensure_governance_files()
    date_text = date_text or ksa_now().strftime("%Y-%m-%d")
    log = safe_read_csv(CHANGE_LOG_PATH)
    impacts = safe_read_csv(CHANGE_IMPACT_PATH)
    if not log.empty and "Updated At" in log.columns:
        daily_log = log[log["Updated At"].astype(str).str.startswith(date_text)].copy()
    else:
        daily_log = pd.DataFrame()
    if not impacts.empty and "Created At" in impacts.columns:
        daily_impacts = impacts[impacts["Created At"].astype(str).str.startswith(date_text)].copy()
    else:
        daily_impacts = pd.DataFrame()
    lines = [f"Executive Daily Digest - {date_text}", "", f"Total changes: {len(daily_log):,}"]
    if not daily_log.empty:
        lines.append(f"Updated Link Codes: {daily_log.get('Link Code', pd.Series(dtype=str)).astype(str).replace('', pd.NA).dropna().nunique():,}")
        lines.append(f"Updated Work Orders: {daily_log.get('Work Order', pd.Series(dtype=str)).astype(str).replace('', pd.NA).dropna().nunique():,}")
        top_users = daily_log.groupby("Updated By").size().sort_values(ascending=False).head(5)
        if not top_users.empty:
            lines += ["", "Top Contributors:"] + [f"- {u}: {c} changes" for u, c in top_users.items()]
        top_fields = daily_log.groupby("Field").size().sort_values(ascending=False).head(8)
        if not top_fields.empty:
            lines += ["", "Most Updated Fields:"] + [f"- {f}: {c}" for f, c in top_fields.items()]
    if not daily_impacts.empty and "Impact" in daily_impacts.columns:
        impact_counts = daily_impacts.groupby("Impact").size().to_dict()
        lines += ["", "Impact Summary:"] + [f"- {k}: {v}" for k, v in impact_counts.items()]
        negatives = daily_impacts[daily_impacts["Impact"].astype(str).str.lower() == "negative"].head(10)
        if not negatives.empty:
            lines += ["", "Critical / Negative Changes:"]
            for _, r in negatives.iterrows():
                lines.append(f"- {r.get('Link Code','')} | {r.get('Work Order','')} | {r.get('Field','')}: {r.get('Old Value','')} → {r.get('New Value','')}")
    return "\n".join(lines)


def create_system_notification(to_user: str, title: str, message: str, category: str = "System", link_code: str = "", work_order: str = "") -> None:
    if not notification_access_enabled(to_user, category):
        return
    now = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
    _append_csv_rows(NOTIFICATIONS_PATH, [{
        "Notification ID": hashlib.md5(f"{now}-{to_user}-{category}-{title}".encode()).hexdigest()[:12],
        "Created At": now, "To User": to_user, "Category": category, "Title": title,
        "Message": message, "Is Read": "No", "Related Link Code": link_code, "Related Work Order": work_order,
    }])


def ensure_governance_files() -> None:
    DATA_DIR.mkdir(exist_ok=True)
    defaults = {
        PROJECT_UPDATES_PATH: ["Link Code", "Work Order", *PROJECT_UPDATE_EDITABLE_COLUMNS, "Updated By", "Updated At"],
        CHANGE_LOG_PATH: ["Change ID", "Updated At", "Updated By", "Link Code", "Work Order", "Field", "Old Value", "New Value", "Source"],
        NOTIFICATIONS_PATH: ["Notification ID", "Created At", "To User", "Category", "Title", "Message", "Is Read", "Related Link Code", "Related Work Order"],
        NOTIFICATION_ACCESS_PATH: ["Username", "Category", "Enabled", "WhatsApp", "Daily Digest", "Critical Only"],
        CHANGE_IMPACT_PATH: ["Impact ID", "Created At", "Change ID", "Updated By", "Link Code", "Work Order", "Field", "Old Value", "New Value", "Impact", "Category", "Impact Note"],
        KPI_IMPACT_PATH: ["KPI Impact ID", "Created At", "Updated By", "Link Code", "Work Order", "Field", "Impact", "KPI Area", "Message"],
        DAILY_DIGEST_PATH: ["Digest ID", "Created At", "Created By", "Audience", "Digest Text"],
        WHATSAPP_OUTBOX_PATH: ["Message ID", "Created At", "To User", "To WhatsApp", "Message", "Status", "Sent At"],
        MASTER_OPERATIONAL_PATH: [],
    }
    for path, cols in defaults.items():
        if not path.exists():
            pd.DataFrame(columns=cols).to_csv(path, index=False, encoding="utf-8-sig")


def _append_csv_rows(path: Path, rows: List[dict]) -> None:
    if not rows:
        return
    ensure_governance_files()
    old = safe_read_csv(path) if path.exists() else pd.DataFrame()
    new_df = pd.DataFrame(rows).fillna("")
    all_cols = list(dict.fromkeys(list(old.columns) + list(new_df.columns)))
    old = old.reindex(columns=all_cols, fill_value="")
    new_df = new_df.reindex(columns=all_cols, fill_value="")
    pd.concat([old, new_df], ignore_index=True).to_csv(path, index=False, encoding="utf-8-sig")


def apply_project_updates_to_workorders(wo: pd.DataFrame) -> pd.DataFrame:
    """Layer 2 overlay: keep u_osp_work_order as master, apply latest user updates from project_updates.csv."""
    if wo is None or wo.empty:
        return wo
    ensure_governance_files()
    updates = safe_read_csv(PROJECT_UPDATES_PATH).copy()
    if updates.empty:
        return wo
    out = wo.copy()
    for col in PROJECT_UPDATE_EDITABLE_COLUMNS:
        if col not in out.columns:
            out[col] = ""
    link_col = first_existing_col(out, ["Link Code"])
    wo_col = first_existing_col(out, ["Work Order", "WO", "WO ID", "Workorder"])
    # Apply by Work Order first, then Link Code fallback. Keep latest row per key.
    if "Updated At" in updates.columns:
        updates = updates.sort_values("Updated At")
    if wo_col and "Work Order" in updates.columns:
        for col in PROJECT_UPDATE_EDITABLE_COLUMNS:
            if col in updates.columns:
                sub = updates[(updates["Work Order"].astype(str).str.strip() != "") & (updates[col].astype(str).str.strip() != "")]
                if not sub.empty:
                    u_by_wo = sub.drop_duplicates("Work Order", keep="last").set_index("Work Order")
                    mapped = out[wo_col].astype(str).map(u_by_wo[col].astype(str).to_dict())
                    out[col] = out[col].where(mapped.isna() | (mapped.astype(str).str.strip() == ""), mapped)
    if link_col and "Link Code" in updates.columns:
        for col in PROJECT_UPDATE_EDITABLE_COLUMNS:
            if col in updates.columns:
                sub = updates[(updates["Link Code"].astype(str).str.strip() != "") & (updates[col].astype(str).str.strip() != "")]
                if not sub.empty:
                    u_by_link = sub.drop_duplicates("Link Code", keep="last").set_index("Link Code")
                    mapped = out[link_col].astype(str).map(u_by_link[col].astype(str).to_dict())
                    out[col] = out[col].where(mapped.isna() | (mapped.astype(str).str.strip() == ""), mapped)
    return out


def create_update_notifications(changes: List[dict]) -> None:
    if not changes:
        return
    users = list(get_users().keys())
    now = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
    rows = []
    for ch in changes:
        impact = classify_change_impact(ch.get("Field", ""), ch.get("Old Value", ""), ch.get("New Value", ""))
        category = category_for_change(ch.get("Field", ""), impact)
        title = f"{category}: {ch.get('Field','')} updated"
        msg = f"{ch.get('Updated By','')} changed {ch.get('Field','')} for Link Code {ch.get('Link Code','')} / WO {ch.get('Work Order','')} from '{ch.get('Old Value','')}' to '{ch.get('New Value','')}'. Impact: {impact}."
        for user in users:
            if str(user).strip().lower() == str(ch.get('Updated By','')).strip().lower():
                continue
            if not notification_access_enabled(user, category):
                continue
            rows.append({
                "Notification ID": hashlib.md5(f"{now}-{user}-{title}-{len(rows)}".encode()).hexdigest()[:12],
                "Created At": now,
                "To User": user,
                "Category": category,
                "Title": title,
                "Message": msg,
                "Is Read": "No",
                "Related Link Code": ch.get("Link Code", ""),
                "Related Work Order": ch.get("Work Order", ""),
            })
    append_change_impacts(changes)
    _append_csv_rows(NOTIFICATIONS_PATH, rows)


def unread_notifications_count(username: str) -> int:
    ensure_governance_files()
    df = safe_read_csv(NOTIFICATIONS_PATH)
    if df.empty or "To User" not in df.columns:
        return 0
    mask = df["To User"].astype(str).str.lower().eq(str(username).lower())
    if "Is Read" in df.columns:
        mask &= df["Is Read"].astype(str).str.lower().ne("yes")
    return int(mask.sum())

def read_cached_document_status_records() -> List[dict]:
    """Read last Google Drive document scan status for dashboard preview.

    The HTML dashboard cannot call Google Drive API directly. The Streamlit
    Document Upload Center scans Google Drive and writes this cache so the
    dashboard preview reflects real files inside each Link Code subfolder.
    """
    try:
        if "document_status_df" in st.session_state and isinstance(st.session_state["document_status_df"], pd.DataFrame):
            df = st.session_state["document_status_df"].copy()
            if not df.empty:
                return df_to_records(df)
    except Exception:
        pass
    try:
        if DOC_STATUS_CACHE_PATH.exists():
            df = safe_read_csv(DOC_STATUS_CACHE_PATH).fillna("")
            return df_to_records(df)
    except Exception:
        return []
    return []


# --- V5.8.12 derived billing/handover fields ---
def _excel_column_by_letter(df: pd.DataFrame, letter: str) -> str:
    """Return DataFrame column name by Excel-style letter (A=1)."""
    n = 0
    for ch in str(letter).strip().upper():
        if "A" <= ch <= "Z":
            n = n * 26 + (ord(ch) - 64)
    idx = n - 1
    return df.columns[idx] if 0 <= idx < len(df.columns) else ""


def apply_derived_billing_fields(wo: pd.DataFrame) -> pd.DataFrame:
    """Apply row-by-row derived fields required for WO Billing & Handover report.

    Excel mapping requested:
    - implementation update uses BH Fiber Status, BI Civil Status, BJ FULL WO STATUS logic per row.
    - SOR Status.1 = BL
    - First 50% status = BM
    - Second 50% status = BP
    """
    if wo is None or wo.empty:
        return wo
    out = wo.copy()

    fiber_col = first_existing_col(out, ["Fiber Status"]) or _excel_column_by_letter(out, "BH")
    civil_col = first_existing_col(out, ["Civil Status"]) or _excel_column_by_letter(out, "BI")
    full_col = first_existing_col(out, ["FULL WO STATUS", "Full WO Status"]) or _excel_column_by_letter(out, "BJ")

    def _txt(v: Any) -> str:
        return str(v if v is not None else "").strip()

    def _impl(row: pd.Series) -> str:
        bj = _txt(row.get(full_col, "")) if full_col else ""
        bi = _txt(row.get(civil_col, "")) if civil_col else ""
        bh = _txt(row.get(fiber_col, "")) if fiber_col else ""
        if bj == "Completed":
            return "Completed"
        if bj == "Not Start":
            return "Civil Not Start"
        if bj == "In Progress":
            if bi == "In Progress":
                return "Civil In Progress"
            if bi == "Completed":
                return "Fiber In Progress"
            if bh == "In Progress":
                return "Fiber In Progress"
        return ""

    out["implementation update"] = out.apply(_impl, axis=1)

    # Preserve source columns and write display/report output columns row-by-row.
    sor_source = _excel_column_by_letter(out, "BL") or first_existing_col(out, ["SOR Status"])
    first50_source = _excel_column_by_letter(out, "BM") or first_existing_col(out, ["1st 50 Invoice Status", "First 50 Invoice Status"])
    second50_source = _excel_column_by_letter(out, "BP") or first_existing_col(out, ["Second 50% status", "2nd 50 Invoice Status", "Second 50 Invoice Status"])

    if sor_source:
        out["SOR Status.1"] = out[sor_source].astype(str)
    if first50_source:
        out["First 50% status"] = out[first50_source].astype(str)
    if second50_source:
        out["Second 50% status"] = out[second50_source].astype(str)

    # 1st 50 Invoice Cost Amount formula:
    # If Scope Target is Implementation, amount = 50% of WO Cost; if WO Cost is blank/zero, use Cost.
    # Otherwise the value is left blank, matching the intended Excel formula behaviour.
    scope_col = first_existing_col(out, ["Scope Target", "Scope_Target"]) or _excel_column_by_letter(out, "BW")
    wo_cost_col = first_existing_col(out, ["WO Cost", "WO COST", "Work Order Cost"]) or _excel_column_by_letter(out, "H")
    cost_col = first_existing_col(out, ["Cost", "COST"]) or _excel_column_by_letter(out, "I")

    def _num(v: Any) -> float:
        s = str(v if v is not None else "").strip()
        if not s or s.lower() in {"nan", "none", "n/a", "na", "-"}:
            return 0.0
        s = s.replace(",", "").replace("SAR", "").replace("ر.س", "").replace("ر.س.", "").strip()
        try:
            return float(s)
        except Exception:
            return 0.0

    def _fmt_amount(x: float) -> str:
        if not x:
            return ""
        if abs(x - round(x)) < 0.005:
            return str(int(round(x)))
        return f"{x:.2f}"

    if scope_col and (wo_cost_col or cost_col):
        def _calc_first50(row: pd.Series) -> str:
            scope = _txt(row.get(scope_col, "")).lower()
            if scope != "implementation":
                return ""
            wo_val = _num(row.get(wo_cost_col, "")) if wo_cost_col else 0.0
            cost_val = _num(row.get(cost_col, "")) if cost_col else 0.0
            base = wo_val if wo_val else cost_val
            return _fmt_amount(base * 0.5)
        out["1st 50 Invoice Cost Amount"] = out.apply(_calc_first50, axis=1)
    elif "1st 50 Invoice Cost Amount" not in out.columns:
        out["1st 50 Invoice Cost Amount"] = ""

    if "Asbuilt Final Amount" not in out.columns:
        out["Asbuilt Final Amount"] = ""
    return out


def build_initial_raw() -> Dict[str, List[dict]]:
    # V4 Phase 1: u_osp_work_order.csv is the main source for operational + location data.
    # District.csv is no longer required for dashboard Region/City/District when those
    # columns are already present in the master workorder file.
    master_wo = master_workorders_df()
    merged_workorders = apply_derived_billing_fields(apply_project_updates_to_workorders(master_wo))
    district_records = derive_district_records_from_workorders(merged_workorders)
    if district_records.empty and DISTRICT_PATH.exists():
        district_records = safe_read_csv(DISTRICT_PATH)
    return {
        "workorders": df_to_records(merged_workorders),
        "penalties": df_to_records(safe_read_csv(PENALTIES_PATH)),
        "districts": df_to_records(district_records),
        "document_status": read_cached_document_status_records(),
    }


def make_safe_js_initial_raw(raw_data: Dict[str, List[dict]]) -> str:
    payload = json.dumps(raw_data, ensure_ascii=False, separators=(",", ":")).encode("utf-8")
    b64 = base64.b64encode(payload).decode("ascii")
    return (
        'const INITIAL_RAW = JSON.parse('
        'new TextDecoder().decode('
        f'Uint8Array.from(atob("{b64}"), c => c.charCodeAt(0))'
        '));\n\n'
    )


def inject_data_into_dashboard(html: str, raw_data: Dict[str, List[dict]]) -> str:
    replacement = make_safe_js_initial_raw(raw_data)

    patterns = [
        r"const\s+INITIAL_RAW\s*=\s*.*?;\s*(?=const\s+FILTERS\s*=)",
        r"const\s+INITIAL_RAW\s*=\s*.*?;\s*(?=const\s+state\s*=)",
    ]

    updated = html
    replaced = 0
    for pattern in patterns:
        updated, replaced = re.subn(pattern, replacement, updated, count=1, flags=re.S)
        if replaced:
            break

    if not replaced:
        st.error("CSV injection failed: INITIAL_RAW block was not found inside dashboard.html.")
        return html

    # Replace the native full-page PDF print listener with the permission-aware selective PDF handler.
    updated = updated.replace(
        "document.getElementById('export-pdf').addEventListener('click', () => window.print());",
        "document.getElementById('export-pdf').addEventListener('click', () => { if (window.DAWIYAT_handleGlobalPdfExport) { window.DAWIYAT_handleGlobalPdfExport(); } else { window.print(); } });"
    )

    role = current_role()
    policy = user_policy()
    allowed_tabs = json.dumps(allowed_dashboard_tabs())
    hide_excel = "true" if not policy.get("export_excel", False) else "false"
    hide_pdf = "true" if not policy.get("export_pdf", False) else "false"
    hide_global_pdf = "true" if not policy.get("global_pdf_report", False) else "false"
    hide_all_exports = "true" if not policy.get("export", False) else "false"
    role_label = role.title()
    hide_buttons = json.dumps(_as_list(policy.get("hide_buttons")))
    hide_tables = json.dumps(_as_list(policy.get("hide_tables")))
    hide_excel_components = json.dumps(_as_list(policy.get("hide_excel_components")))
    allowed_excel_components = json.dumps(_as_list(policy.get("allowed_excel_components")))
    hide_pdf_components = json.dumps(_as_list(policy.get("hide_pdf_components")))
    allowed_pdf_components = json.dumps(_as_list(policy.get("allowed_pdf_components")))
    hide_ppt_components = json.dumps(_as_list(policy.get("hide_ppt_components")))
    smart_bulk_filter = json.dumps(_current_smart_bulk_filter_payload(), ensure_ascii=False)
    all_dashboard_tabs = ["overview", "tables", "pmo", "performance", "perf-explanation", "decision", "reports"]
    denied_tabs = [t for t in all_dashboard_tabs if t not in allowed_dashboard_tabs()]
    deny_tab_css = "\n".join([f'.tab[data-tab="{t}"], .report-tab[data-tab="{t}"], [data-tab="{t}"], #tab-{t} {{ display: none !important; visibility: hidden !important; }}' for t in denied_tabs])
    export_button_css = ""
    if hide_global_pdf == "true" or hide_all_exports == "true":
        export_button_css += "#export-pdf, button#export-pdf, .btn#export-pdf { display:none !important; visibility:hidden !important; pointer-events:none !important; }\n"
    if hide_excel == "true" or hide_all_exports == "true":
        export_button_css += "#export-excel, button#export-excel, [id*=\"export-excel\"] { display:none !important; visibility:hidden !important; pointer-events:none !important; }\n"

    portal_patch = f"""
<style>
{deny_tab_css}
{export_button_css}
#dawiyat-selective-pdf-root {{ display:none; }}
@media print {{
  body.dawiyat-selective-pdf-active > *:not(#dawiyat-selective-pdf-root) {{ display:none !important; }}
  body.dawiyat-selective-pdf-active #dawiyat-selective-pdf-root {{ display:block !important; padding: 0 !important; margin: 0 !important; }}
  #dawiyat-selective-pdf-root .dawiyat-pdf-cover {{ border:1px solid #d9e3ef; border-radius:16px; padding:14px 18px; margin-bottom:12px; break-inside:avoid; }}
  #dawiyat-selective-pdf-root .dawiyat-pdf-cover h1 {{ margin:0 0 6px; font-size:20px; color:#10223a; }}
  #dawiyat-selective-pdf-root .dawiyat-pdf-cover p {{ margin:0; font-size:11px; color:#64748b; }}
  #dawiyat-selective-pdf-root .dawiyat-pdf-block {{ display:block !important; visibility:visible !important; break-inside:avoid-page; page-break-inside:avoid; margin-bottom:12px; }}
  #dawiyat-selective-pdf-root .hidden {{ display:block !important; visibility:visible !important; }}
  #dawiyat-selective-pdf-root .btn, #dawiyat-selective-pdf-root button, #dawiyat-selective-pdf-root input, #dawiyat-selective-pdf-root select, #dawiyat-selective-pdf-root .action-row {{ display:none !important; }}
  #dawiyat-selective-pdf-root .table-wrap, #dawiyat-selective-pdf-root .pmo-table-wrap, #dawiyat-selective-pdf-root .assist-table-wrap {{ max-height:none !important; overflow:visible !important; }}
}}
.file-label, #apply-imports {{ display: none !important; }}
.header-actions::after {{
    content: "Data linked directly from Version 2 Executive Portal";
    display: block;
    color: #64748b;
    font-size: 12px;
    text-align: right;
    margin-top: 4px;
}}
body.role-viewer .doc-repo-panel,
body.role-board .doc-repo-panel,
body.role-finance .doc-repo-panel {{ display:none !important; }}
body.hide-excel button,
body.hide-excel .btn {{}}
</style>
<script>
window.DAWIYAT_SMART_BULK_FILTER = {smart_bulk_filter};
window.DAWIYAT_RBAC = {{
  role: {json.dumps(role)},
  permissionMode: "user-based-only",
  roleLabel: {json.dumps(role_label)},
  allowedTabs: {allowed_tabs},
  hideExcel: {hide_excel},
  hidePdf: {hide_pdf},
  hideGlobalPdf: {hide_global_pdf},
  hideAllExports: {hide_all_exports},
  hideButtons: {hide_buttons},
  hideTables: {hide_tables},
  hideExcelComponents: {hide_excel_components},
  excelAllowedComponents: {allowed_excel_components},
  hidePdfComponents: {hide_pdf_components},
  pdfAllowedComponents: {allowed_pdf_components},
  hidePptComponents: {hide_ppt_components}
}};
(function applySmartBulkFilterFromStreamlit() {{
  function apply() {{
    try {{
      const f = window.DAWIYAT_SMART_BULK_FILTER || {{}};
      if (typeof state === 'undefined' || typeof renderAll !== 'function') return;
      state.uploadedSiteFilter = f.active ? f : {{active:false, fileName:'', linkCodes:[], workOrders:[], linkColumn:'', woColumn:''}};
      if (f.active) {{
        if (Array.isArray(f.linkCodes) && f.linkCodes.length) state.filters.linkCode = f.linkCodes;
        if (Array.isArray(f.workOrders) && f.workOrders.length) state.filters.workOrder = f.workOrders;
        if (state.pmo && state.pmo.filters) {{
          if (Array.isArray(f.linkCodes) && f.linkCodes.length) state.pmo.filters.linkCode = f.linkCodes;
          if (Array.isArray(f.workOrders) && f.workOrders.length) state.pmo.filters.workOrder = f.workOrders;
        }}
      }}
      renderAll();
    }} catch(e) {{ console.warn('Smart bulk filter apply failed', e); }}
  }}
  window.addEventListener('load', () => setTimeout(apply, 250));
}})();
(function applyDawiyatRBAC() {{
  let rbacApplying = false;
  function norm(t) {{ return (t || '').replace(/\s+/g,' ').trim().toLowerCase(); }}
  function softNorm(t) {{
    return norm(t)
      .replace(/executive/g, '')
      .replace(/[–—-]/g, ' ')
      .replace(/&/g, ' and ')
      .replace(/[^a-z0-9؀-ۿ ]/g, ' ')
      .replace(/\s+/g,' ')
      .trim();
  }}
  function titleMatches(title, needle) {{
    const a = softNorm(title), b = softNorm(needle);
    if (!a || !b) return false;
    return a.includes(b) || b.includes(a);
  }}
  function hideExportButtons() {{
    const cfg = window.DAWIYAT_RBAC || {{}};
    if (cfg.hideGlobalPdf || cfg.hideAllExports) {{
      const pdfBtn = document.getElementById('export-pdf');
      if (pdfBtn) {{
        pdfBtn.style.setProperty('display','none','important');
        pdfBtn.style.setProperty('visibility','hidden','important');
        pdfBtn.style.setProperty('pointer-events','none','important');
      }}
    }}
    const buttons = Array.from(document.querySelectorAll('button, a, .btn'));
    buttons.forEach(el => {{
      if (el.dataset && el.dataset.forceExport) return;
      const txt = norm(el.textContent);
      if (cfg.hideAllExports && (txt.includes('export') || txt.includes('csv') || txt.includes('excel') || txt.includes('pdf'))) {{
        el.style.display = 'none';
        return;
      }}
      if (cfg.hideExcel && (txt.includes('export excel') || txt.includes('export csv') || txt === 'export')) {{
        el.style.display = 'none';
      }}
      if (cfg.hidePdf && txt.includes('export pdf') && el.id !== 'export-pdf') {{
        el.style.display = 'none';
      }}
      (cfg.hideButtons || []).forEach(needle => {{
        if (needle && txt.includes(String(needle).toLowerCase())) el.style.display = 'none';
      }});
    }});
  }}
  function blockTitle(el) {{
    const h = el.querySelector && el.querySelector('h1,h2,h3,.panel-head,.report-title');
    return h ? norm(h.textContent) : '';
  }}
  function hideTablesByText() {{
    const cfg = window.DAWIYAT_RBAC || {{}};
    const needles = (cfg.hideTables || []).map(x => String(x || '')).filter(Boolean);
    if (!needles.length) return;
    const blocks = Array.from(document.querySelectorAll('.panel, .kpi, .table-wrap, .report-card, .report-section')); // do not hide top-level tab sections
    blocks.forEach(el => {{
      const title = blockTitle(el);
      if (title && needles.some(n => titleMatches(title, n))) {{
        el.style.setProperty('display', 'none', 'important');
        el.style.setProperty('visibility', 'hidden', 'important');
      }}
    }});
  }}
  function hideComponentExportButtons() {{
    const cfg = window.DAWIYAT_RBAC || {{}};
    const blocks = Array.from(document.querySelectorAll('.panel, .report-card, .report-section-card, .pmo-compact-panel, .sor-exec-summary-panel, .stage-exec-summary-panel, .overview-exec-summary-panel, .overview-summary-table-card, .assist-card')); // scoped blocks only; avoids hiding table buttons from parent containers
    blocks.forEach(block => {{
      const title = blockTitle(block);
      if (!title) return;
      const hideExcel = (cfg.hideExcelComponents || []).some(n => titleMatches(title, n));
      const hidePdf = (cfg.hidePdfComponents || []).some(n => titleMatches(title, n));
      const hidePpt = (cfg.hidePptComponents || []).some(n => titleMatches(title, n));
      if (!hideExcel && !hidePdf && !hidePpt) return;
      Array.from(block.querySelectorAll('button, a, .btn')).forEach(el => {{
        if (el.dataset && el.dataset.forceExport) return;
        const txt = norm(el.textContent);
        if (hideExcel && (txt.includes('excel') || txt.includes('csv'))) el.style.display = 'none';
        if (hidePdf && txt.includes('pdf')) el.style.display = 'none';
        if (hidePpt && (txt.includes('ppt') || txt.includes('presentation') || txt.includes('powerpoint'))) el.style.display = 'none';
      }});
    }});
  }}
  function showAllowedExcelExportButtons() {{
    const cfg = window.DAWIYAT_RBAC || {{}};
    if (cfg.hideAllExports || cfg.hideExcel) return;
    const allowed = (cfg.excelAllowedComponents || []).map(x => String(x || '')).filter(Boolean);
    if (!allowed.length) return;
    const tablesTabAllowed = (cfg.allowedTabs || []).includes('tables');
    const blocks = Array.from(document.querySelectorAll('.panel, .report-card, .report-section-card, .pmo-compact-panel, .sor-exec-summary-panel, .stage-exec-summary-panel, .overview-exec-summary-panel, .overview-summary-table-card, .assist-card'));
    blocks.forEach(block => {{
      const title = blockTitle(block) || norm(block.textContent || '').slice(0, 220);
      const isAllowedBlock = allowed.some(n => titleMatches(title, n));
      if (!isAllowedBlock) return;
      block.style.removeProperty('display');
      block.style.removeProperty('visibility');
      Array.from(block.querySelectorAll('button, a, .btn')).forEach(el => {{
        const txt = norm(el.textContent);
        const call = String(el.getAttribute('onclick') || '').toLowerCase();
        const id = String(el.id || '').toLowerCase();
        const isExcelButton = txt.includes('export excel') || txt.includes('excel') || txt.includes('csv') || call.includes('exporttabletoexcel') || id.includes('export') && id.includes('excel');
        if (isExcelButton) {{
          el.style.removeProperty('display');
          el.style.removeProperty('visibility');
          el.style.removeProperty('pointer-events');
          el.style.setProperty('pointer-events','auto','important');
        }}
      }});
    }});
    // If the workbook contains a generic "Export Excel" component for Tables & Exports, keep all
    // table export buttons visible on that tab when the Tables tab itself is allowed.
    if (tablesTabAllowed && allowed.some(n => titleMatches(n, 'Export Excel'))) {{
      const tablesSection = document.getElementById('tab-tables');
      if (tablesSection) {{
        Array.from(tablesSection.querySelectorAll('button, a, .btn')).forEach(el => {{
          const txt = norm(el.textContent);
          const call = String(el.getAttribute('onclick') || '').toLowerCase();
          if (txt.includes('export excel') || call.includes('exporttabletoexcel')) {{
            el.style.removeProperty('display');
            el.style.removeProperty('visibility');
            el.style.removeProperty('pointer-events');
            el.style.setProperty('pointer-events','auto','important');
          }}
        }});
      }}
    }}
  }}
  function pdfPrintableBlocks() {{
    return Array.from(document.querySelectorAll('.panel, .report-card, .report-section-card, .pmo-compact-panel, .sor-exec-summary-panel, .stage-exec-summary-panel, .overview-exec-summary-panel, .overview-summary-table-card, .assist-card'));
  }}
  function findPdfComponentBlock(componentName) {{
    const name = String(componentName || '');
    const exactSelectorMap = [
      ['Link Code Summary Table', '#link-summary-table'],
      ['WO Billing & Handover Status Report', '#wo-status-report-panel'],
      ['WO Billing & Handover Status Report', '#reports-wo-billing-handover-table'],
      ['PMO Audit', '#pmo-audit-panel'],
      ['PMO Audit Table', '#pmo-audit-panel'],
      ['PMO Master Table', '#pmo-master-panel'],
      ['SOR Details', '#pmo-sor-details-panel'],
      ['SOR Summary & Revenue Exposure', '#reports-sor-summary'],
      ['Overall Projects Stages', '#pmo-cost-analysis-panel'],
      ['Overall Projects Stages — Cost Analysis', '#pmo-cost-analysis-panel'],
      ['Overall Stages Summary & Cost Exposure', '#reports-stage-summary'],
      ['Executive Overall Stages Summary & Cost Exposure', '#reports-stage-summary'],
      ['Executive KPI Cards', '#reports-kpi-grid'],
      ['Portfolio Summary & Cost Exposure', '#reports-portfolio-summary'],
      ['Executive Portfolio Summary & Cost Exposure', '#reports-portfolio-summary'],
      ['Executive Reports Library', '#reports-orange-dynamic'],
      ['Document Upload Center — Status Preview', '#doc-status-preview-panel']
    ];
    for (const pair of exactSelectorMap) {{
      if (titleMatches(pair[0], name)) {{
        const direct = document.querySelector(pair[1]);
        if (direct) return direct.closest('.panel, .report-card, .report-section-card, .pmo-compact-panel, .sor-exec-summary-panel, .stage-exec-summary-panel, .overview-exec-summary-panel, .overview-summary-table-card, .assist-card') || direct;
      }}
    }}
    const candidates = pdfPrintableBlocks();
    let best = null;
    let bestScore = 0;
    candidates.forEach(block => {{
      const title = blockTitle(block) || norm(block.textContent || '').slice(0, 220);
      if (!title) return;
      if (titleMatches(title, name)) {{
        const score = Math.min(softNorm(title).length, softNorm(name).length);
        if (score > bestScore) {{ best = block; bestScore = score; }}
      }}
    }});
    return best;
  }}
  function prepareSelectivePdfPrint() {{
    const cfg = window.DAWIYAT_RBAC || {{}};
    const allowed = (cfg.pdfAllowedComponents || []).map(x => String(x || '').trim()).filter(Boolean);
    if (!allowed.length) {{
      alert('No PDF-exportable tables/components are enabled for this user. Enable Export PDF for at least one Component / Table in Admin Board.');
      return false;
    }}
    restoreSelectivePdfPrint();
    const allowedTabs = new Set(cfg.allowedTabs || []);
    ['overview','performance','tables','decision','pmo','perf-explanation','reports'].forEach(tab => {{
      const sec = document.getElementById('tab-' + tab);
      if (sec && allowedTabs.has(tab)) {{
        sec.dataset.pdfWasHidden = sec.classList.contains('hidden') ? '1' : '0';
        sec.dataset.pdfOldDisplay = sec.style.display || '';
        sec.classList.remove('hidden');
        sec.style.setProperty('display','block','important');
        sec.style.setProperty('visibility','visible','important');
      }}
    }});
    if (typeof renderAll === 'function') {{ try {{ renderAll(); }} catch(e) {{}} }}
    if (typeof renderExecutiveReports === 'function') {{ try {{ renderExecutiveReports(); }} catch(e) {{}} }}
    if (typeof renderExecutiveOverviewSummary === 'function') {{ try {{ renderExecutiveOverviewSummary(getFilteredWorkorders ? getFilteredWorkorders() : []); }} catch(e) {{}} }}

    const root = document.createElement('div');
    root.id = 'dawiyat-selective-pdf-root';
    const cover = document.createElement('div');
    cover.className = 'dawiyat-pdf-cover';
    cover.innerHTML = '<h1>Dawiyat PMO Executive PDF Report</h1><p>Selective PDF export based on User_Component_Access → Export PDF permissions. Current filters and Smart Bulk Filter are applied before export.</p>';
    root.appendChild(cover);

    const missing = [];
    let matched = 0;
    allowed.forEach(componentName => {{
      const block = findPdfComponentBlock(componentName);
      if (!block) {{ missing.push(componentName); return; }}
      const clone = block.cloneNode(true);
      clone.classList.add('dawiyat-pdf-block');
      clone.style.removeProperty('display');
      clone.style.removeProperty('visibility');
      clone.querySelectorAll('[style]').forEach(el => {{
        if (String(el.style.display || '').includes('none')) el.style.removeProperty('display');
        el.style.removeProperty('visibility');
        el.style.removeProperty('pointer-events');
      }});
      root.appendChild(clone);
      matched += 1;
    }});
    if (!matched) {{
      restoreSelectivePdfPrint();
      alert('PDF permission is enabled, but no dashboard component matched the enabled Component / Table names: ' + allowed.join(', ') + '. Please use the exact component names available in Admin Board.');
      return false;
    }}
    if (missing.length) console.warn('PDF components not found:', missing);
    document.body.appendChild(root);
    document.body.classList.add('dawiyat-selective-pdf-active');
    return true;
  }}
  function restoreSelectivePdfPrint() {{
    document.body.classList.remove('dawiyat-selective-pdf-active');
    const root = document.getElementById('dawiyat-selective-pdf-root');
    if (root) root.remove();
    ['overview','performance','tables','decision','pmo','perf-explanation','reports'].forEach(tab => {{
      const sec = document.getElementById('tab-' + tab);
      if (!sec || sec.dataset.pdfWasHidden === undefined) return;
      if (sec.dataset.pdfWasHidden === '1') sec.classList.add('hidden');
      else sec.classList.remove('hidden');
      sec.style.display = sec.dataset.pdfOldDisplay || '';
      sec.style.removeProperty('visibility');
      delete sec.dataset.pdfWasHidden;
      delete sec.dataset.pdfOldDisplay;
    }});
  }}
  function handleGlobalPdfExport() {{
    if (prepareSelectivePdfPrint()) setTimeout(() => window.print(), 160);
  }}
  window.DAWIYAT_prepareSelectivePdfPrint = prepareSelectivePdfPrint;
  window.DAWIYAT_restoreSelectivePdfPrint = restoreSelectivePdfPrint;
  window.DAWIYAT_handleGlobalPdfExport = handleGlobalPdfExport;
  window.addEventListener('afterprint', restoreSelectivePdfPrint);
  document.addEventListener('click', function(e) {{
    const target = e.target && e.target.closest ? e.target.closest('#export-pdf') : null;
    if (!target) return;
    e.preventDefault();
    e.stopPropagation();
    if (e.stopImmediatePropagation) e.stopImmediatePropagation();
    handleGlobalPdfExport();
  }}, true);

  function applyTabs() {{
    if (rbacApplying) return;
    rbacApplying = true;
    const cfg = window.DAWIYAT_RBAC || {{ allowedTabs: ['overview'] }};
    const allowed = new Set(cfg.allowedTabs || ['overview']);
    document.body.classList.add('role-' + (cfg.role || 'viewer'));
    document.querySelectorAll('.tab[data-tab], .report-tab[data-tab], [data-tab]').forEach(btn => {{
      const ok = allowed.has(btn.dataset.tab);
      if (ok) {{ btn.style.removeProperty('display'); btn.style.removeProperty('visibility'); }}
      else {{ btn.style.setProperty('display','none','important'); btn.style.setProperty('visibility','hidden','important'); }}
    }});
    ['overview','performance','tables','decision','pmo','perf-explanation','reports'].forEach(tab => {{
      const sec = document.getElementById('tab-' + tab);
      if (sec && !allowed.has(tab)) {{ sec.classList.add('hidden'); sec.style.setProperty('display','none','important'); }}
      if (sec && allowed.has(tab)) {{ sec.style.removeProperty('display'); }}
    }});
    const active = document.querySelector('.tab.active[data-tab]');
    const activeAllowed = active && allowed.has(active.dataset.tab);
    if (!activeAllowed) {{
      const first = Array.from(allowed)[0] || 'overview';
      if (typeof window.setTab === 'function') window.setTab(first);
      else {{
        document.querySelectorAll('.tab[data-tab]').forEach(t => t.classList.toggle('active', t.dataset.tab === first));
        ['overview','performance','tables','decision','pmo','perf-explanation','reports'].forEach(tab => {{
          const sec = document.getElementById('tab-' + tab);
          if (sec) sec.classList.toggle('hidden', tab !== first);
        }});
      }}
    }}
    hideExportButtons();
    hideTablesByText();
    hideComponentExportButtons();
    showAllowedExcelExportButtons();
    rbacApplying = false;
  }}
  document.addEventListener('DOMContentLoaded', applyTabs);
  setTimeout(applyTabs, 250);
  setTimeout(applyTabs, 800);
  setTimeout(applyTabs, 1800);
  setTimeout(applyTabs, 3500);
  try {{
    const mo = new MutationObserver(() => setTimeout(applyTabs, 60));
    mo.observe(document.body, {{childList:true, subtree:true, attributes:true, attributeFilter:['class','style']}});
  }} catch(e) {{}}
}})();
</script>
"""
    if "</head>" in updated:
        updated = updated.replace("</head>", portal_patch + "\n</head>", 1)

    return updated


def parse_num(value) -> float:
    if value is None:
        return 0.0
    match = re.search(r"-?\d[\d,]*(?:\.\d+)?", str(value))
    if not match:
        return 0.0
    try:
        return float(match.group(0).replace(",", ""))
    except Exception:
        return 0.0



def effective_row_cost_from_wo(wo_df: pd.DataFrame) -> pd.Series:
    """Match dashboard effectiveCost: first positive value from WO Cost, then Cost, then related cost columns."""
    candidates = ["WO Cost", "Cost", "WO Cost / Cost", "Work Order Cost", "Total WO Cost", "Total Cost"]
    result = pd.Series([0.0] * len(wo_df), index=wo_df.index, dtype="float64")
    for col in candidates:
        if col in wo_df.columns:
            vals = wo_df[col].apply(parse_num)
            result = result.where(result > 0, vals)
    return result.fillna(0.0)

def first_existing_col(df: pd.DataFrame, candidates: List[str]) -> str:
    normalized = {re.sub(r"[^a-z0-9]", "", c.lower()): c for c in df.columns}
    for c in candidates:
        key = re.sub(r"[^a-z0-9]", "", c.lower())
        if key in normalized:
            return normalized[key]
    return ""


def load_workorders() -> pd.DataFrame:
    wo = apply_derived_billing_fields(apply_project_updates_to_workorders(master_workorders_df()))
    if wo.empty:
        return wo
    link_col = first_existing_col(wo, ["Link Code"])
    wo_col = first_existing_col(wo, ["Work Order"])
    cost_col = first_existing_col(wo, ["WO Cost", "Cost"])
    progress_col = first_existing_col(wo, ["Percentage of Completion"])
    stage_col = first_existing_col(wo, ["Stage", "Invoice Status", "FULL WO STATUS"])
    sor_col = first_existing_col(wo, ["SOR Status"])
    status_col = first_existing_col(wo, ["Work Order Status"])
    updated_col = first_existing_col(wo, ["Updated"])

    wo["_link"] = wo[link_col].astype(str) if link_col else ""
    wo["_wo"] = wo[wo_col].astype(str) if wo_col else ""
    wo["_cost"] = wo[cost_col].apply(parse_num) if cost_col else 0
    wo["_progress"] = wo[progress_col].apply(parse_num) if progress_col else 0
    wo["_stage"] = wo[stage_col].astype(str) if stage_col else ""
    wo["_sor"] = wo[sor_col].astype(str) if sor_col else ""
    wo["_status"] = wo[status_col].astype(str) if status_col else ""
    wo["_updated"] = wo[updated_col].astype(str) if updated_col else ""
    return wo


def portfolio_metrics() -> Dict[str, float]:
    wo = load_workorders()
    if wo.empty:
        return {"links": 0, "wos": 0, "cost": 0, "progress": 0}
    return {
        "links": wo["_link"].replace("", pd.NA).dropna().nunique(),
        "wos": len(wo),
        "cost": wo["_cost"].sum(),
        "progress": wo["_progress"].mean(),
    }



# ---------------- Smart Bulk Site Filter ----------------
LINK_CODE_HEADER_ALIASES = [
    "link code", "linkcode", "link_code", "link", "site code", "site id", "site", "area code",
]
WORK_ORDER_HEADER_ALIASES = [
    "work order", "workorder", "work_order", "wo", "wo id", "wo_id", "work order id", "order id",
]


def _norm_header_name(value: Any) -> str:
    return re.sub(r"[^a-z0-9]", "", str(value or "").strip().lower())


def _find_smart_filter_column(df: pd.DataFrame, aliases: List[str]) -> str:
    alias_keys = {_norm_header_name(x) for x in aliases}
    for col in df.columns:
        key = _norm_header_name(col)
        if key in alias_keys:
            return str(col)
    # soft contains matching, useful for files like "Dawiyat Link Code" or "WO Number"
    for col in df.columns:
        key = _norm_header_name(col)
        if any(a and (a in key or key in a) for a in alias_keys):
            return str(col)
    return ""


def _clean_smart_filter_values(values: Any) -> List[str]:
    out: List[str] = []
    seen = set()
    for value in list(values or []):
        if pd.isna(value):
            continue
        # Allow pasted cells containing several IDs separated by commas/new lines.
        parts = re.split(r"[\n,;\t]+", str(value))
        for part in parts:
            text = str(part or "").strip()
            if not text or text.lower() in {"nan", "none", "null"}:
                continue
            key = re.sub(r"\s+", "", text).upper()
            if key not in seen:
                seen.add(key)
                out.append(text)
    return out





def _smart_norm_id(value: Any) -> str:
    return re.sub(r"\s+", "", str(value or "").strip().upper())


def _smart_row_link(row: dict) -> str:
    """Return Link Code from either Streamlit-normalized dataframe rows or dashboard-style dict rows."""
    return str(
        row.get("linkCode", "")
        or row.get("_link", "")
        or row.get("Link Code", "")
        or row.get("LINK CODE", "")
        or ""
    ).strip()


def _smart_row_wo(row: dict) -> str:
    """Return Work Order from either Streamlit-normalized dataframe rows or dashboard-style dict rows."""
    return str(
        row.get("workOrder", "")
        or row.get("_wo", "")
        or row.get("Work Order", "")
        or row.get("WORK ORDER", "")
        or row.get("WO", "")
        or row.get("WO ID", "")
        or ""
    ).strip()


def _smart_bulk_scope_summary(workorders: List[dict]) -> Dict[str, Any]:
    """Validate uploaded Smart Bulk IDs against dashboard data.

    Final required logic:
    - Link Codes and Work Orders are checked independently against u_osp_work_order.
    - A missing Link Code is shown only in the Missing Link Codes table.
    - A missing Work Order is shown only in the Missing Work Orders table.
    - If a Work Order exists but the uploaded Link Code label does not exist, the WO is still matched;
      the missing Link Code is reported for data-quality review only.
    - Filtering stays inclusive and operational; validation is only for transparency.
    """
    uploaded_links = _clean_smart_filter_values(st.session_state.get("smart_bulk_link_codes", []))
    uploaded_wos = _clean_smart_filter_values(st.session_state.get("smart_bulk_work_orders", []))
    uploaded_pairs = st.session_state.get("smart_bulk_pairs", []) or (st.session_state.get("smart_bulk_uploaded", {}) or {}).get("pairs", []) or []

    data_links_map: Dict[str, List[dict]] = {}
    data_wos_map: Dict[str, dict] = {}
    for r in (workorders or []):
        ln_raw = _smart_row_link(r)
        wn_raw = _smart_row_wo(r)
        ln = _smart_norm_id(ln_raw)
        wn = _smart_norm_id(wn_raw)
        if ln:
            data_links_map.setdefault(ln, []).append(r)
        if wn:
            data_wos_map[wn] = r

    data_links = set(data_links_map.keys())
    data_wos = set(data_wos_map.keys())

    def _dedupe_display(values: List[str]) -> List[str]:
        seen = set(); out = []
        for v in values or []:
            n = _smart_norm_id(v)
            if n and n not in seen:
                seen.add(n); out.append(str(v).strip())
        return out

    uploaded_links = _dedupe_display(uploaded_links)
    uploaded_wos = _dedupe_display(uploaded_wos)

    matched_links = [x for x in uploaded_links if _smart_norm_id(x) in data_links]
    missing_links = [x for x in uploaded_links if _smart_norm_id(x) and _smart_norm_id(x) not in data_links]
    matched_wos = [x for x in uploaded_wos if _smart_norm_id(x) in data_wos]
    missing_wos = [x for x in uploaded_wos if _smart_norm_id(x) and _smart_norm_id(x) not in data_wos]

    # Context table for missing Link Codes: show related WOs from the uploaded rows and whether those WOs exist.
    missing_link_rows: List[Dict[str, str]] = []
    seen_missing_links = set()
    for pair in uploaded_pairs or []:
        link_original = str(pair.get("link_code", "") or "").strip()
        wo_original = str(pair.get("work_order", "") or "").strip()
        link_norm = _smart_norm_id(link_original)
        wo_norm = _smart_norm_id(wo_original)
        if not link_norm or link_norm in data_links or link_norm in seen_missing_links:
            continue
        seen_missing_links.add(link_norm)
        wo_exists = bool(wo_norm and wo_norm in data_wos)
        actual_link_from_wo = _smart_row_link(data_wos_map[wo_norm]) if wo_exists else ""
        missing_link_rows.append({
            "Missing Link Code": link_original,
            "Uploaded Work Order in same row": wo_original or "—",
            "Work Order Found?": "Yes" if wo_exists else "No",
            "Actual Link Code from Work Order": actual_link_from_wo or "—",
        })
    # Add any manual/multiselect missing links that were not part of pairs.
    for link_original in missing_links:
        link_norm = _smart_norm_id(link_original)
        if link_norm not in seen_missing_links:
            seen_missing_links.add(link_norm)
            missing_link_rows.append({
                "Missing Link Code": link_original,
                "Uploaded Work Order in same row": "—",
                "Work Order Found?": "—",
                "Actual Link Code from Work Order": "—",
            })

    # Context table for missing Work Orders: show related Link Code and whether that Link Code exists.
    missing_wo_rows: List[Dict[str, str]] = []
    seen_missing_wos = set()
    for pair in uploaded_pairs or []:
        link_original = str(pair.get("link_code", "") or "").strip()
        wo_original = str(pair.get("work_order", "") or "").strip()
        link_norm = _smart_norm_id(link_original)
        wo_norm = _smart_norm_id(wo_original)
        if not wo_norm or wo_norm in data_wos or wo_norm in seen_missing_wos:
            continue
        seen_missing_wos.add(wo_norm)
        link_exists = bool(link_norm and link_norm in data_links)
        missing_wo_rows.append({
            "Missing Work Order": wo_original,
            "Uploaded Link Code in same row": link_original or "—",
            "Link Code Found?": "Yes" if link_exists else "No",
        })
    for wo_original in missing_wos:
        wo_norm = _smart_norm_id(wo_original)
        if wo_norm not in seen_missing_wos:
            seen_missing_wos.add(wo_norm)
            missing_wo_rows.append({
                "Missing Work Order": wo_original,
                "Uploaded Link Code in same row": "—",
                "Link Code Found?": "—",
            })

    # Operational matched rows are still based on Work Orders first, plus Link-Code fallback when no WOs are supplied.
    # This keeps dashboard row counts stable while still surfacing all missing IDs in the validation tables.
    matched_rows_by_index: Dict[int, dict] = {}
    if uploaded_wos:
        for r in (workorders or []):
            if _smart_norm_id(_smart_row_wo(r)) in {_smart_norm_id(x) for x in uploaded_wos}:
                matched_rows_by_index[id(r)] = r
    elif uploaded_links:
        for r in (workorders or []):
            if _smart_norm_id(_smart_row_link(r)) in {_smart_norm_id(x) for x in uploaded_links}:
                matched_rows_by_index[id(r)] = r
    matched_rows = list(matched_rows_by_index.values())

    return {
        "uploaded_links": uploaded_links,
        "uploaded_wos": uploaded_wos,
        "uploaded_pairs": uploaded_pairs,
        "matched_links_count": len(matched_links),
        "matched_wos_count": len(matched_wos),
        "missing_link_labels": missing_links,
        "missing_wos": missing_wos,
        "missing_link_rows": missing_link_rows,
        "missing_wo_rows": missing_wo_rows,
        "matched_rows_count": len(matched_rows),
        "mode": "Independent Link Code + Work Order validation; Work Order controls filtering when available",
    }

def _smart_bulk_options_from_workorders(workorders: List[dict]) -> Tuple[List[str], List[str]]:
    link_options = sorted({str(r.get("linkCode", "")).strip() for r in workorders if str(r.get("linkCode", "")).strip()}, key=lambda x: x.upper())
    wo_options = sorted({str(r.get("workOrder", "")).strip() for r in workorders if str(r.get("workOrder", "")).strip()}, key=lambda x: x.upper())
    # Keep uploaded/manual/search values in the option list so Streamlit multiselect never receives a default not in options.
    link_options = sorted(set(link_options) | set(_clean_smart_filter_values(st.session_state.get("smart_bulk_link_codes", []))), key=lambda x: x.upper())
    wo_options = sorted(set(wo_options) | set(_clean_smart_filter_values(st.session_state.get("smart_bulk_work_orders", []))), key=lambda x: x.upper())
    return link_options, wo_options

def _parse_smart_bulk_upload(uploaded_file) -> Dict[str, Any]:
    if uploaded_file is None:
        return {"link_codes": [], "work_orders": [], "pairs": [], "link_column": "", "wo_column": "", "file_name": ""}
    name = getattr(uploaded_file, "name", "uploaded_file")
    try:
        if str(name).lower().endswith(".csv"):
            df = pd.read_csv(uploaded_file, dtype=str, encoding_errors="ignore")
        else:
            df = pd.read_excel(uploaded_file, dtype=str)
    except Exception as exc:
        st.error(f"Could not read uploaded Smart Bulk Filter file: {exc}")
        return {"link_codes": [], "work_orders": [], "pairs": [], "link_column": "", "wo_column": "", "file_name": name}

    df.columns = [str(c).strip() for c in df.columns]
    link_col = _find_smart_filter_column(df, LINK_CODE_HEADER_ALIASES)
    wo_col = _find_smart_filter_column(df, WORK_ORDER_HEADER_ALIASES)
    link_codes = _clean_smart_filter_values(df[link_col].dropna().tolist()) if link_col else []
    work_orders = _clean_smart_filter_values(df[wo_col].dropna().tolist()) if wo_col else []

    # Keep row-level pairs so Missing Scope is based on the uploaded row relationship:
    # Link Code + Work Order together. This prevents reporting all independent Link Codes
    # as missing when the Work Order exists and should be the actual matching key.
    pairs: List[Dict[str, str]] = []
    if link_col or wo_col:
        for _, row in df.iterrows():
            link_vals = _clean_smart_filter_values([row.get(link_col, "")]) if link_col else []
            wo_vals = _clean_smart_filter_values([row.get(wo_col, "")]) if wo_col else []
            max_len = max(len(link_vals), len(wo_vals), 0)
            for i in range(max_len):
                link = link_vals[i] if i < len(link_vals) else (link_vals[0] if link_vals else "")
                wo = wo_vals[i] if i < len(wo_vals) else (wo_vals[0] if wo_vals else "")
                if str(link or "").strip() or str(wo or "").strip():
                    pairs.append({"link_code": str(link or "").strip(), "work_order": str(wo or "").strip()})

    # de-duplicate pairs while preserving order
    seen_pairs = set()
    unique_pairs = []
    for pair in pairs:
        key = (_smart_norm_id(pair.get("link_code", "")), _smart_norm_id(pair.get("work_order", "")))
        if key not in seen_pairs:
            seen_pairs.add(key)
            unique_pairs.append(pair)

    if not link_codes and not work_orders:
        st.warning("No Link Code or Work Order column was detected. The file can use names such as Link Code, LinkCode, Site Code, WO, WO ID, or Work Order.")
    return {
        "link_codes": link_codes,
        "work_orders": work_orders,
        "pairs": unique_pairs,
        "link_column": link_col,
        "wo_column": wo_col,
        "file_name": name,
    }


def _current_smart_bulk_filter_payload() -> Dict[str, Any]:
    link_codes = _clean_smart_filter_values(st.session_state.get("smart_bulk_link_codes", []))
    work_orders = _clean_smart_filter_values(st.session_state.get("smart_bulk_work_orders", []))
    uploaded = st.session_state.get("smart_bulk_uploaded", {}) or {}
    pairs = st.session_state.get("smart_bulk_pairs", []) or uploaded.get("pairs", []) or []
    return {
        "active": bool(link_codes or work_orders),
        "fileName": uploaded.get("file_name", "Manual / Search Selection") if uploaded else "Manual / Search Selection",
        "linkCodes": link_codes,
        "workOrders": work_orders,
        "pairs": pairs,
        "linkColumn": uploaded.get("link_column", ""),
        "woColumn": uploaded.get("wo_column", ""),
        "filterMode": "WORK_ORDER_PRIMARY" if work_orders else ("LINK_CODE" if link_codes else "NONE"),
    }


def render_smart_bulk_filter_panel(raw: Dict[str, List[dict]]) -> None:
    """Streamlit-native Smart Bulk Filter panel above the dashboard iframe.
    It supports Excel/CSV upload, searchable multi-select chips for Link Code and Work Order,
    and manual paste for large WO lists. Reset Filters in dashboard clears the iframe state;
    the Clear button here clears the Streamlit selection before rendering the dashboard.
    """
    workorders = raw.get("workorders", []) or []

    is_active = bool(st.session_state.get("smart_bulk_link_codes") or st.session_state.get("smart_bulk_work_orders"))
    show_filter = st.session_state.get("show_smart_bulk_filter", False)

    top_cols = st.columns([1.2, 3, 1])
    with top_cols[0]:
        if st.button(("🙈 Hide Smart Bulk Filter" if show_filter else "🎯 Show Smart Bulk Filter"), use_container_width=True, key="toggle_smart_bulk_filter"):
            st.session_state["show_smart_bulk_filter"] = not show_filter
            st.rerun()
    with top_cols[1]:
        if is_active:
            st.success(f"Smart Bulk Filter active: {len(st.session_state.get('smart_bulk_link_codes', []))} Link Codes, {len(st.session_state.get('smart_bulk_work_orders', []))} Work Orders.")
    with top_cols[2]:
        if is_active and st.button("🧹 Clear", use_container_width=True, key="smart_bulk_quick_clear"):
            for k in [
                "smart_bulk_uploaded", "smart_bulk_link_codes", "smart_bulk_work_orders", "smart_bulk_pairs",
                "smart_bulk_link_codes_multiselect", "smart_bulk_work_orders_multiselect", "smart_bulk_manual_wo_text",
            ]:
                st.session_state.pop(k, None)
            st.rerun()

    if not st.session_state.get("show_smart_bulk_filter", False):
        return

    with st.container(border=True):
        st.markdown("### 🎯 Smart Bulk Filter")
        st.caption("Upload any Excel/CSV containing Link Code and/or Work Order, or search directly. The selection is applied to the full dashboard and PMO Audit.")

        uploaded = st.file_uploader(
            "Upload Excel / CSV containing Link Code and/or Work Order",
            type=["xlsx", "xls", "csv"],
            key="smart_bulk_upload_file",
        )
        if uploaded is not None:
            parsed = _parse_smart_bulk_upload(uploaded)
            st.session_state["smart_bulk_uploaded"] = parsed
            if parsed.get("link_codes"):
                st.session_state["smart_bulk_link_codes"] = parsed["link_codes"]
            if parsed.get("work_orders"):
                st.session_state["smart_bulk_work_orders"] = parsed["work_orders"]
            if parsed.get("pairs"):
                st.session_state["smart_bulk_pairs"] = parsed["pairs"]
            if parsed.get("link_codes") or parsed.get("work_orders"):
                st.success(
                    f"Loaded from {parsed.get('file_name','file')}: "
                    f"{len(parsed.get('link_codes', []))} Link Codes, "
                    f"{len(parsed.get('work_orders', []))} Work Orders."
                )

        # IMPORTANT: recompute after upload parsing so newly uploaded values become valid options
        # before st.multiselect is created. This prevents StreamlitAPIException when the
        # uploaded Excel contains IDs not already present in the dashboard dataset.
        link_options, wo_options = _smart_bulk_options_from_workorders(workorders)

        col1, col2, col3, col4 = st.columns([1.15, 1.15, 1.15, .8])
        with col1:
            selected_links = st.multiselect(
                "Scan Link Codes",
                options=link_options,
                default=[x for x in st.session_state.get("smart_bulk_link_codes", []) if x in set(link_options)],
                key="smart_bulk_link_codes_multiselect",
                placeholder="Search / select Link Codes",
            )
        with col2:
            selected_wos = st.multiselect(
                "Scan Work Orders",
                options=wo_options,
                default=[x for x in st.session_state.get("smart_bulk_work_orders", []) if x in set(wo_options)],
                key="smart_bulk_work_orders_multiselect",
                placeholder="Search / select Work Orders",
            )
        with col3:
            manual_wo = st.text_area(
                "Add Work Orders",
                value=st.session_state.get("smart_bulk_manual_wo_text", ""),
                key="smart_bulk_manual_wo_text",
                placeholder="Paste WO IDs separated by new lines or commas",
                height=80,
            )
        with col4:
            st.write("")
            st.write("")
            if st.button("🧹 Clear Smart Filter", use_container_width=True, key="smart_bulk_clear"):
                for k in [
                    "smart_bulk_uploaded", "smart_bulk_link_codes", "smart_bulk_work_orders", "smart_bulk_pairs",
                    "smart_bulk_link_codes_multiselect", "smart_bulk_work_orders_multiselect", "smart_bulk_manual_wo_text",
                ]:
                    st.session_state.pop(k, None)
                st.rerun()

        manual_wos = _clean_smart_filter_values([manual_wo])
        final_links = _clean_smart_filter_values(selected_links or st.session_state.get("smart_bulk_link_codes", []))
        final_wos = _clean_smart_filter_values((selected_wos or st.session_state.get("smart_bulk_work_orders", [])) + manual_wos)
        st.session_state["smart_bulk_link_codes"] = final_links
        st.session_state["smart_bulk_work_orders"] = final_wos

        scope = _smart_bulk_scope_summary(workorders)
        st.markdown("#### Smart Bulk Scope Check")
        c1, c2, c3 = st.columns(3)
        with c1:
            st.metric("Uploaded Scope", f"{len(scope['uploaded_links'])} Links / {len(scope['uploaded_wos'])} WOs")
            st.caption("Values read directly from the uploaded Excel/CSV file.")
        with c2:
            st.metric("Matched Scope", f"{scope['matched_links_count']} Links / {scope['matched_wos_count']} WOs")
            st.caption(scope.get("mode", "Independent validation"))
        with c3:
            missing_link_count = len(scope.get('missing_link_labels', []) or [])
            missing_wo_count = len(scope.get('missing_wos', []) or [])
            st.metric("Missing Scope", f"{missing_link_count} Links / {missing_wo_count} WOs")
            st.caption("Missing Link Codes and Missing Work Orders are checked separately against u_osp_work_order.")

        missing_link_rows = scope.get('missing_link_rows', []) or []
        missing_wo_rows = scope.get('missing_wo_rows', []) or []
        if missing_link_rows or missing_wo_rows:
            with st.expander("Show Missing Link Codes / Work Orders", expanded=True):
                if missing_link_rows:
                    st.markdown("##### Missing Link Codes")
                    st.dataframe(pd.DataFrame(missing_link_rows), use_container_width=True, hide_index=True)
                else:
                    st.success("No Missing Link Codes.")
                if missing_wo_rows:
                    st.markdown("##### Missing Work Orders")
                    st.dataframe(pd.DataFrame(missing_wo_rows), use_container_width=True, hide_index=True)
                else:
                    st.success("No Missing Work Orders.")
        else:
            st.success("No missing Link Codes or Work Orders were found in u_osp_work_order.")

@st.cache_data(show_spinner=False)
def read_dashboard_html_cached(path_str: str, mtime: float) -> str:
    path = Path(path_str)
    if not path.exists():
        return ""
    return path.read_text(encoding="utf-8", errors="ignore")


def render_dashboard() -> None:
    if not DASHBOARD_PATH.exists():
        st.error("Dashboard HTML file is missing: dashboard/dashboard.html")
        return

    raw = build_initial_raw()

    with st.sidebar.expander("Data check", expanded=False):
        st.write(f"Work Orders: {len(raw['workorders']):,}")
        st.write(f"Penalties: {len(raw['penalties']):,}")
        st.write(f"Location records from u_osp_work_order: {len(raw['districts']):,}")

    render_smart_bulk_filter_panel(raw)

    # Hidden action/governance pages: shown as compact buttons on Dashboard according to permissions.
    all_allowed = allowed_pages_for_current_user()
    quick_actions = []
    if "Project Updates Center" in all_allowed:
        quick_actions.append(("📝 Open Project Updates Center", "Project Updates Center", "secondary"))
    if "Data Update Agent" in all_allowed:
        quick_actions.append(("🧠 Open Data Update Agent", "Data Update Agent", "secondary"))
    if "Notification Center 🔔" in all_allowed:
        quick_actions.append((f"🔔 Open Notification Center ({unread_notifications_count(st.session_state.get('username',''))})", "Notification Center 🔔", "secondary"))
    if "Executive Daily Digest" in all_allowed:
        quick_actions.append(("📩 Open Executive Daily Digest", "Executive Daily Digest", "secondary"))
    if "WhatsApp Agent" in all_allowed:
        quick_actions.append(("🟢 Open WhatsApp Agent Outbox", "WhatsApp Agent", "secondary"))
    if "📤 Document Upload Center" in all_allowed:
        quick_actions.append(("📤 Open Document Upload Center", "📤 Document Upload Center", "secondary"))
    if "📊 Executive PPT Builder" in all_allowed:
        quick_actions.append(("📊 Open Executive PPT Builder", "📊 Executive PPT Builder", "secondary"))
    if "Admin Board" in all_allowed and _is_admin_board_owner():
        quick_actions.append(("⚙️ Open Admin Board", "Admin Board", "primary"))

    if quick_actions:
        st.markdown(
            """
            <div class="quick-actions-panel">
                <div class="quick-actions-title">Quick Actions & Governance Agents</div>
                <div class="quick-actions-subtitle">Data Update Agent, Notification Center, Daily Digest, WhatsApp Agent, Document Center, PPT Builder, and Admin Board open only from here according to user permissions.</div>
            </div>
            """,
            unsafe_allow_html=True,
        )
        for i in range(0, len(quick_actions), 4):
            action_cols = st.columns(min(4, len(quick_actions) - i))
            for col, (label, target_page, btn_type) in zip(action_cols, quick_actions[i:i+4]):
                with col:
                    if st.button(label, use_container_width=True, type=btn_type, key=f"open_hidden_{target_page}"):
                        st.session_state["active_hidden_page"] = target_page
                        st.rerun()

    dashboard_html = read_dashboard_html_cached(str(DASHBOARD_PATH), DASHBOARD_PATH.stat().st_mtime)
    dashboard_html = inject_data_into_dashboard(dashboard_html, raw)

    components.html(dashboard_html, height=1800, scrolling=True)



def render_back_to_dashboard_button(key_suffix: str = "") -> None:
    """Back button for hidden action/governance pages opened from Dashboard Quick Actions."""
    safe_key = f"back_to_dashboard_{key_suffix}" if key_suffix else "back_to_dashboard"
    back_cols = st.columns([1.2, 4])
    with back_cols[0]:
        if st.button("← Back to Dashboard", use_container_width=True, key=safe_key):
            st.session_state.pop("active_hidden_page", None)
            st.session_state["force_dashboard"] = True
            st.rerun()


def _unique_keep_order(items: List[str]) -> List[str]:
    out: List[str] = []
    for item in items:
        if item and item not in out:
            out.append(item)
    return out


def _component_enabled_for_current_user(component_name: str) -> bool:
    """Return True when the component is explicitly Show=Yes in permissions.xlsx.

    This keeps Project Updates controlled by the User_Component_Access sheet instead
    of only the display Role text.
    """
    target = str(component_name or "").strip().lower()
    if not target:
        return False
    try:
        policy = user_policy()
        shown = [str(x).strip().lower() for x in _as_list(policy.get("show_tables"))]
        return target in shown
    except Exception:
        return False


def _explicit_editable_columns_from_permissions() -> List[str]:
    """Optional future-proof column-level permissions.

    If Admin Board contains rows like:
      Page = Project Updates Center
      Component / Table = Editable Column: SOR Status
      Show = Yes
    then only those explicit columns are added in addition to role defaults.
    """
    explicit: List[str] = []
    try:
        policy = user_policy()
        for comp in _as_list(policy.get("show_tables")):
            txt = str(comp or "").strip()
            low = txt.lower()
            if low.startswith("editable column:"):
                col = txt.split(":", 1)[1].strip()
                if col in PROJECT_UPDATE_EDITABLE_COLUMNS:
                    explicit.append(col)
            elif low in {"save updates", "bulk update", "editable columns"}:
                # Control components; not columns themselves. Kept for permission UI clarity.
                continue
    except Exception:
        pass
    return _unique_keep_order(explicit)


def current_user_update_columns() -> List[str]:
    username = str(st.session_state.get("username", "")).strip().lower()
    role = str(st.session_state.get("role", "")).strip().lower()

    # Page + component gate. A user must have the Project Updates page and the
    # Project Updates Editable Grid component enabled. Ahmed remains full admin.
    if username == "ahmedfekry" or role in {"admin", "pmo"}:
        return PROJECT_UPDATE_EDITABLE_COLUMNS.copy()

    if not can("project_updates"):
        return []

    has_grid = _component_enabled_for_current_user("Project Updates Editable Grid")
    if not has_grid:
        return []

    explicit_cols = _explicit_editable_columns_from_permissions()

    # Finance / invoicing / document controller users.
    # Fixes mohamed_syed: his display role is "Invoicing & Document Controller",
    # so checking only for the word "finance" made him view-only.
    if (
        username in {"mohamed_syed", "mohamed_sayed"}
        or any(k in role for k in ["finance", "invoice", "invoicing", "document", "commercial", "billing"])
    ):
        return _unique_keep_order(
            PROJECT_UPDATE_FINANCE_COLUMNS
            + ["DCR_Status", "RFS Certificate", "As-built BOQ", "Redline", "Handover O&M _Status", "Handover Consultant _Status", "Asbuilt Final Amount"]
            + explicit_cols
        )

    # Project users can update implementation/closure fields. SOR Status is included
    # because the PMO workflow requires project managers to update the same row-level
    # status set without being blocked on this specific column.
    if username == "adham_ismail" or any(k in role for k in ["project", "manager", "operation", "operations", "pm", "coordinator"]):
        return _unique_keep_order(PROJECT_UPDATE_PM_COLUMNS + ["SOR Status"] + explicit_cols)

    # Fallback: if Admin grants explicit editable columns, use them even if the role
    # label is not recognized. This prevents valid users from becoming view-only just
    # because the display Role text changed.
    if explicit_cols:
        return explicit_cols

    return []

def project_updates_center_page() -> None:
    render_back_to_dashboard_button("project_updates")
    st.title("📝 Project Updates Center")
    st.caption("Authorized users update controlled status columns only. Master data remains protected; changes are saved in data/project_updates.csv and overlaid on the dashboard after refresh.")
    if not can("project_updates") and not can("admin"):
        st.error("You do not have permission to access Project Updates Center.")
        return

    df = master_workorders_df().copy()
    if df.empty:
        st.warning("u_osp_work_order.csv is empty or missing.")
        return
    # Ensure required editable columns exist in current data file.
    for col in PROJECT_UPDATE_EDITABLE_COLUMNS:
        if col not in df.columns:
            df[col] = ""

    st.markdown("#### 🎯 Smart Bulk Filter for Project Updates")
    st.caption("Use the same Smart Bulk Filter to narrow the editable grid to a large selected list of Link Codes and/or Work Orders before making updates.")
    render_smart_bulk_filter_panel({"workorders": df.to_dict("records")})

    link_col = first_existing_col(df, ["Link Code"])
    wo_col = first_existing_col(df, ["Work Order"])
    region_col = first_existing_col(df, ["Region"])
    project_col = first_existing_col(df, ["Project"])
    stage_col = first_existing_col(df, ["Stage"])

    st.markdown("#### 🎛️ Dashboard Filters for Project Updates")
    st.caption("Use the same dashboard filters used by the Executive PPT Builder. Filters are applied first, then Smart Bulk / Link Code / Work Order search narrows the editable grid further.")

    def _pu_options(source_df: pd.DataFrame, aliases: List[str]) -> List[str]:
        col = first_existing_col(source_df, aliases)
        if not col or col not in source_df.columns:
            return []
        vals = []
        for v in source_df[col].dropna().astype(str).tolist():
            t = str(v).strip()
            if t and t.lower() not in {"nan", "none"}:
                vals.append(t)
        return sorted(list(dict.fromkeys(vals)))

    def _pu_apply(source_df: pd.DataFrame, aliases: List[str], selected: List[str]) -> pd.DataFrame:
        col = first_existing_col(source_df, aliases)
        selected = [str(x).strip() for x in (selected or []) if str(x).strip()]
        if not col or not selected:
            return source_df
        return source_df[source_df[col].astype(str).str.strip().isin(selected)]

    with st.expander("🎛️ Advanced Dashboard Filters", expanded=False):
        f1, f2, f3 = st.columns(3)
        with f1:
            pu_region = st.multiselect("Region", _pu_options(df, ["Region", "Updated Region"]), key="pu_filter_region")
            city_base = _pu_apply(df, ["Region", "Updated Region"], pu_region)
            pu_city = st.multiselect("City", _pu_options(city_base, ["City", "Updated City"]), key="pu_filter_city")
            district_base = _pu_apply(city_base, ["City", "Updated City"], pu_city)
            pu_district = st.multiselect("District", _pu_options(district_base, ["District", "WO Districts", "Updated District"]), key="pu_filter_district")
            pu_project = st.multiselect("Project", _pu_options(df, ["Project"]), key="pu_filter_project")
        with f2:
            pu_stage = st.multiselect("Stage", _pu_options(df, ["Stage"]), key="pu_filter_stage")
            pu_subclass = st.multiselect("Subclass", _pu_options(df, ["Subclass", "Sub Class"]), key="pu_filter_subclass")
            pu_year = st.multiselect("Year", _pu_options(df, ["Year"]), key="pu_filter_year")
            pu_wo_status = st.multiselect("Work Order Status", _pu_options(df, ["Work Order Status", "FULL WO STATUS"]), key="pu_filter_wo_status")
            pu_second50 = st.multiselect("Second 50% status", _pu_options(df, ["Second 50% status", "2nd 50 Invoice Status", "Second 50 Invoice Status"]), key="pu_filter_second50")
        with f3:
            pu_sor = st.multiselect("SOR Status", _pu_options(df, ["SOR Status", "SOR Status.1"]), key="pu_filter_sor")
            pu_sor_ref = st.multiselect("SOR Reference Number", _pu_options(df, ["SOR Reference Number", "SOR Ref Number", "SOR Reference"]), key="pu_filter_sor_ref")
            pu_type = st.multiselect("Type", _pu_options(df, ["Type"]), key="pu_filter_type")
            pu_class = st.multiselect("Class", _pu_options(df, ["Class"]), key="pu_filter_class")
            pu_scope = st.multiselect("Scope Target", _pu_options(df, ["Scope Target"]), key="pu_filter_scope")
            pu_invoice50 = st.multiselect("1st 50 Invoice Status", _pu_options(df, ["1st 50 Invoice Status", "First 50 Invoice Status"]), key="pu_filter_invoice50")

        if st.button("Clear Advanced Filters", use_container_width=True, key="pu_clear_advanced_filters"):
            for k in ["pu_filter_region","pu_filter_city","pu_filter_district","pu_filter_project","pu_filter_stage","pu_filter_subclass","pu_filter_year","pu_filter_wo_status","pu_filter_second50","pu_filter_sor","pu_filter_sor_ref","pu_filter_type","pu_filter_class","pu_filter_scope","pu_filter_invoice50"]:
                st.session_state[k] = []
            st.rerun()

    st.markdown("#### 🔎 Direct Search Filters")
    st.caption("Same search-and-add style used in the main dashboard filters. Select one or many Link Codes / Work Orders to narrow the editable grid.")
    c1, c2, c3 = st.columns([1.1, 1.1, 0.7])
    link_options = []
    wo_options = []
    if link_col and link_col in df.columns:
        link_options = sorted([str(x).strip() for x in df[link_col].dropna().astype(str).unique().tolist() if str(x).strip()])
    if wo_col and wo_col in df.columns:
        wo_options = sorted([str(x).strip() for x in df[wo_col].dropna().astype(str).unique().tolist() if str(x).strip()])
    with c1:
        selected_search_links = st.multiselect(
            "Link Code (search & add)",
            options=link_options,
            default=st.session_state.get("project_update_search_links", []),
            placeholder="Search / paste / select Link Codes",
            key="project_update_search_links",
        )
        st.caption("Select one or many Link Codes.")
    with c2:
        selected_search_wos = st.multiselect(
            "Work Order (search & add)",
            options=wo_options,
            default=st.session_state.get("project_update_search_wos", []),
            placeholder="Search / paste / select Work Orders",
            key="project_update_search_wos",
        )
        st.caption("Select one or many Work Orders.")
    with c3:
        max_rows = st.number_input("Max Rows", min_value=25, max_value=1000, value=200, step=25)
        if st.button("Clear Search", use_container_width=True, key="project_update_clear_search"):
            st.session_state["project_update_search_links"] = []
            st.session_state["project_update_search_wos"] = []
            st.rerun()

    view = df.copy()

    # Apply Project Updates advanced filters using the same concepts as Executive PPT Builder.
    # These filters use AND logic and remain linked with Smart Bulk and direct search filters below.
    _advanced_filters = [
        (["Region", "Updated Region"], st.session_state.get("pu_filter_region", [])),
        (["City", "Updated City"], st.session_state.get("pu_filter_city", [])),
        (["District", "WO Districts", "Updated District"], st.session_state.get("pu_filter_district", [])),
        (["Project"], st.session_state.get("pu_filter_project", [])),
        (["Stage"], st.session_state.get("pu_filter_stage", [])),
        (["Subclass", "Sub Class"], st.session_state.get("pu_filter_subclass", [])),
        (["Year"], st.session_state.get("pu_filter_year", [])),
        (["Work Order Status", "FULL WO STATUS"], st.session_state.get("pu_filter_wo_status", [])),
        (["Second 50% status", "2nd 50 Invoice Status", "Second 50 Invoice Status"], st.session_state.get("pu_filter_second50", [])),
        (["SOR Status", "SOR Status.1"], st.session_state.get("pu_filter_sor", [])),
        (["SOR Reference Number", "SOR Ref Number", "SOR Reference"], st.session_state.get("pu_filter_sor_ref", [])),
        (["Type"], st.session_state.get("pu_filter_type", [])),
        (["Class"], st.session_state.get("pu_filter_class", [])),
        (["Scope Target"], st.session_state.get("pu_filter_scope", [])),
        (["1st 50 Invoice Status", "First 50 Invoice Status"], st.session_state.get("pu_filter_invoice50", [])),
    ]
    for aliases, selected in _advanced_filters:
        view = _pu_apply(view, aliases, selected)

    smart_links = set(_clean_smart_filter_values(st.session_state.get("smart_bulk_link_codes", [])))
    smart_wos = set(_clean_smart_filter_values(st.session_state.get("smart_bulk_work_orders", [])))
    if (smart_links or smart_wos) and (link_col or wo_col):
        # Same source-of-truth rule as the dashboard:
        # If Work Orders exist, filter exact WOs. Otherwise use Link Code.
        if smart_wos and wo_col:
            view = view[view[wo_col].astype(str).str.strip().str.upper().isin({x.upper() for x in smart_wos})]
            st.info(f"Smart Bulk Filter applied to Project Updates by exact Work Order list: {len(smart_wos)} Work Orders.")
        elif smart_links and link_col:
            view = view[view[link_col].astype(str).str.strip().str.upper().isin({x.upper() for x in smart_links})]
            st.info(f"Smart Bulk Filter applied to Project Updates by Link Code: {len(smart_links)} Link Codes.")

    search_links = set(_clean_smart_filter_values(selected_search_links))
    search_wos = set(_clean_smart_filter_values(selected_search_wos))
    if (search_links or search_wos) and (link_col or wo_col):
        mask = pd.Series(False, index=view.index)
        if search_links and link_col:
            mask = mask | view[link_col].astype(str).str.strip().str.upper().isin({x.upper() for x in search_links})
        if search_wos and wo_col:
            mask = mask | view[wo_col].astype(str).str.strip().str.upper().isin({x.upper() for x in search_wos})
        view = view[mask]
        st.info(f"Direct Search applied: {len(search_links)} Link Codes / {len(search_wos)} Work Orders.")
    view = view.head(int(max_rows)).copy()
    view.insert(0, "_row_id", view.index.astype(int))

    editable_cols = current_user_update_columns()
    id_cols = [c for c in ["_row_id", link_col, wo_col, region_col, project_col, stage_col] if c and c in view.columns]
    display_cols = id_cols + [c for c in PROJECT_UPDATE_EDITABLE_COLUMNS if c in view.columns]
    view = view[display_cols]

    if not editable_cols:
        st.info("Your role is view-only for this update center.")

    col_config = {"_row_id": st.column_config.NumberColumn("Row ID", disabled=True)}
    for col, options in PROJECT_UPDATE_STATUS_OPTIONS.items():
        if col in view.columns:
            existing_values = [str(v).strip() for v in view[col].dropna().astype(str).unique().tolist() if str(v).strip()]
            values = list(dict.fromkeys(["", *options, *existing_values]))
            col_config[col] = st.column_config.SelectboxColumn(col, options=values, required=False)
    if "1st 50 Invoice Cost Amount" in view.columns:
        col_config["1st 50 Invoice Cost Amount"] = st.column_config.TextColumn("1st 50 Invoice Cost Amount")
    if "Asbuilt Final Amount" in view.columns:
        col_config["Asbuilt Final Amount"] = st.column_config.TextColumn("Asbuilt Final Amount")
    if "SOR Reference Number" in view.columns:
        col_config["SOR Reference Number"] = st.column_config.TextColumn("SOR Reference Number")

    disabled_cols = [c for c in view.columns if c == "_row_id" or c not in editable_cols]
    edited = st.data_editor(
        view,
        use_container_width=True,
        hide_index=True,
        disabled=disabled_cols,
        column_config=col_config,
        key="project_updates_editor",
    )
    st.caption(f"Showing {len(view):,} rows. Editable columns for this user: {', '.join(editable_cols) if editable_cols else 'None'}")

    if st.button("💾 Save Project Updates", type="primary", use_container_width=True):
        if not editable_cols:
            st.error("No editable columns are assigned to your role.")
            return
        base_full = master_workorders_df().copy()
        effective_full = apply_project_updates_to_workorders(base_full).copy()
        for col in PROJECT_UPDATE_EDITABLE_COLUMNS:
            if col not in base_full.columns:
                base_full[col] = ""
            if col not in effective_full.columns:
                effective_full[col] = ""
        username = st.session_state.get('username','')
        now = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
        update_rows = []
        change_rows = []
        for _, row in edited.iterrows():
            try:
                idx = int(row.get("_row_id"))
            except Exception:
                continue
            if idx not in effective_full.index:
                continue
            link_val = str(effective_full.at[idx, link_col]) if link_col else str(row.get("Link Code", ""))
            wo_val = str(effective_full.at[idx, wo_col]) if wo_col else str(row.get("Work Order", ""))
            any_changed = False
            wide_row = {"Link Code": link_val, "Work Order": wo_val, "Updated By": username, "Updated At": now}
            for col in editable_cols:
                if col in edited.columns:
                    old_val = str(effective_full.at[idx, col]) if col in effective_full.columns else ""
                    new_val = str(row.get(col, ""))
                    wide_row[col] = new_val
                    if old_val != new_val:
                        any_changed = True
                        change_rows.append({
                            "Change ID": hashlib.md5(f"{now}-{username}-{wo_val}-{col}-{old_val}-{new_val}".encode()).hexdigest()[:12],
                            "Updated At": now, "Updated By": username, "Link Code": link_val, "Work Order": wo_val,
                            "Field": col, "Old Value": old_val, "New Value": new_val, "Source": "Project Updates Center",
                        })
            if any_changed:
                update_rows.append(wide_row)
        if update_rows:
            _append_csv_rows(PROJECT_UPDATES_PATH, update_rows)
            _append_csv_rows(CHANGE_LOG_PATH, change_rows)
            create_update_notifications(change_rows)
            # Persist the effective operational layer back to the local master CSV as well.
            # This makes Project Updates survive a normal Streamlit app process reboot and keeps
            # dashboard.html injection aligned for every user after refresh/login.
            try:
                effective_after_save = apply_derived_billing_fields(apply_project_updates_to_workorders(base_full)).copy()
                backup_file(WO_PATH)
                effective_after_save.to_csv(WO_PATH, index=False, encoding="utf-8-sig")
            except Exception as e:
                st.warning(f"Updates were saved, but master CSV persistence failed: {e}")
            snapshot_path = write_master_operational_snapshot("project_update")
        else:
            snapshot_path = None
        st.cache_data.clear()
        if snapshot_path:
            st.success(f"Saved {len(change_rows):,} changed cells by {username}. Master operational data, project_updates.csv, and snapshot were updated. Snapshot: {snapshot_path.name}")
        else:
            st.info("No data changes detected.")


def backup_file(path: Path) -> Path:
    stamp = ksa_now().strftime("%Y%m%d_%H%M%S")
    dest = BACKUP_DIR / f"{path.stem}_{stamp}{path.suffix}"
    if path.exists():
        shutil.copy(path, dest)
    return dest



def data_update_agent_page() -> None:
    render_back_to_dashboard_button("data_update_agent")
    st.title("🧠 Data Update Agent")
    st.caption("Governed data layers, update persistence, snapshots, notification access, change impact, and digest controls.")
    ensure_governance_files()
    tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
        "Data Layers", "Master Data Admin", "User Notification Access", "Change Detection Agent", "Data Version Center", "KPI Impact Engine"
    ])
    with tab1:
        st.subheader("Dashboard Data Layers")
        st.markdown("""
        **Layer 1 — Protected Master:** `u_osp_work_order.csv`, `District.csv`, `Penalties.csv`  
        **Layer 2 — Incremental Updates:** `project_updates.csv`  
        **Layer 3 — Operational Master:** `master_operational_data.csv` = Master + latest updates + derived fields  
        **Layer 4 — Audit & Intelligence:** `change_log.csv`, `change_impact.csv`, `kpi_impact.csv`, `notifications.csv`
        """)
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Project Updates", len(safe_read_csv(PROJECT_UPDATES_PATH)))
        c2.metric("Change Log Rows", len(safe_read_csv(CHANGE_LOG_PATH)))
        c3.metric("Unread Notifications", unread_notifications_count(st.session_state.get("username", "")))
        op = safe_read_csv(MASTER_OPERATIONAL_PATH)
        c4.metric("Operational Rows", len(op))
        if st.button("🔄 Rebuild Operational Master + Snapshot", use_container_width=True):
            snap = write_master_operational_snapshot("manual_rebuild")
            st.cache_data.clear()
            st.success(f"Operational master rebuilt. Snapshot: {snap.name if snap else 'N/A'}")
    with tab2:
        st.subheader("Admin-only Master Data")
        if not _is_admin_board_owner():
            st.info("District.csv and Penalties.csv are protected master files. Only ahmedfekry can replace them.")
        else:
            for label, path in PROJECT_MASTER_ADMIN_FILES.items():
                st.markdown(f"### {label}")
                df = safe_read_csv(path)
                st.caption(f"Current rows: {len(df):,} | Path: data/{path.name}")
                uploaded = st.file_uploader(f"Replace {label}", type=["csv"], key=f"master_replace_{label}")
                if uploaded is not None and st.button(f"💾 Save {label}", key=f"save_master_{label}"):
                    backup_file(path)
                    path.write_bytes(uploaded.getvalue())
                    write_master_operational_snapshot(f"replace_{path.stem}")
                    st.cache_data.clear()
                    st.success(f"{label} replaced successfully and operational snapshot refreshed.")
    with tab3:
        st.subheader("User Notification Access")
        st.caption("Controls who receives Project Update, Financial, Documents, PAT, Handover, Implementation, Critical, Daily Digest, and WhatsApp queue notifications.")
        users = list(get_users().keys())
        categories = ["ALL", "Project Update", "Financial", "Documents", "PAT", "Handover", "Implementation", "Critical", "Daily Digest", "KPI Impact", "System"]
        access = safe_read_csv(NOTIFICATION_ACCESS_PATH)
        if access.empty or "Username" not in access.columns:
            rows = []
            for u in users:
                for cat in categories:
                    rows.append({"Username": u, "Category": cat, "Enabled": "Yes" if u.lower()=="ahmedfekry" or cat not in ["WhatsApp"] else "No", "WhatsApp": "No", "Daily Digest": "Yes", "Critical Only": "No"})
            access = pd.DataFrame(rows)
        edited = st.data_editor(access, use_container_width=True, hide_index=True, num_rows="dynamic", key="notification_access_editor")
        if st.button("💾 Save Notification Access", use_container_width=True):
            edited.fillna("").to_csv(NOTIFICATION_ACCESS_PATH, index=False, encoding="utf-8-sig")
            st.cache_data.clear()
            st.success("Notification access saved. New notifications will follow this matrix.")
    with tab4:
        st.subheader("Change Detection Agent")
        log = safe_read_csv(CHANGE_LOG_PATH)
        impacts = safe_read_csv(CHANGE_IMPACT_PATH)
        c1, c2, c3 = st.columns(3)
        c1.metric("Raw Changes", len(log))
        c2.metric("Impact Rows", len(impacts))
        neg = impacts[impacts.get("Impact", "").astype(str).str.lower()=="negative"] if not impacts.empty and "Impact" in impacts.columns else pd.DataFrame()
        c3.metric("Negative/Critical", len(neg))
        if impacts.empty:
            st.info("No change impacts recorded yet.")
        else:
            st.dataframe(impacts.sort_values("Created At", ascending=False), use_container_width=True, hide_index=True)
            st.download_button("⬇️ Download Change Impact", impacts.to_csv(index=False).encode("utf-8-sig"), "change_impact.csv", "text/csv")
        if not log.empty:
            with st.expander("Raw Change Log"):
                st.dataframe(log.sort_values("Updated At", ascending=False), use_container_width=True, hide_index=True)
                st.download_button("⬇️ Download Change Log", log.to_csv(index=False).encode("utf-8-sig"), "change_log.csv", "text/csv")
    with tab5:
        st.subheader("Data Version Center")
        st.caption("Download the latest operational master or any saved snapshot. Use this before/after major updates.")
        if MASTER_OPERATIONAL_PATH.exists():
            st.download_button("⬇️ Download Current Operational Master", MASTER_OPERATIONAL_PATH.read_bytes(), "master_operational_data.csv", "text/csv", use_container_width=True)
        else:
            st.info("No operational master yet. Click Rebuild Operational Master in Data Layers.")
        SNAPSHOT_DIR.mkdir(parents=True, exist_ok=True)
        snapshots = sorted(SNAPSHOT_DIR.glob("master_operational_data_*.csv"), reverse=True)
        st.metric("Snapshots", len(snapshots))
        if snapshots:
            snap_df = pd.DataFrame([{"Snapshot": x.name, "Modified": datetime.fromtimestamp(x.stat().st_mtime, KSA_TZ).strftime("%Y-%m-%d %H:%M:%S"), "Size KB": round(x.stat().st_size/1024, 1)} for x in snapshots])
            st.dataframe(snap_df, use_container_width=True, hide_index=True)
            selected = st.selectbox("Select snapshot to download", [x.name for x in snapshots])
            snap_path = next((x for x in snapshots if x.name == selected), None)
            if snap_path:
                st.download_button("⬇️ Download Selected Snapshot", snap_path.read_bytes(), snap_path.name, "text/csv", use_container_width=True)
    with tab6:
        st.subheader("KPI Impact Engine")
        kpi = safe_read_csv(KPI_IMPACT_PATH)
        if kpi.empty:
            st.info("No KPI impacts yet. Positive/Negative changes will appear here after project updates.")
        else:
            st.dataframe(kpi.sort_values("Created At", ascending=False), use_container_width=True, hide_index=True)
            st.download_button("⬇️ Download KPI Impact", kpi.to_csv(index=False).encode("utf-8-sig"), "kpi_impact.csv", "text/csv")


def notification_center_page() -> None:
    render_back_to_dashboard_button("notification_center")
    st.title("🔔 Notification Center")
    ensure_governance_files()
    username = st.session_state.get("username", "")
    df = safe_read_csv(NOTIFICATIONS_PATH)
    if df.empty or "To User" not in df.columns:
        st.info("No notifications yet.")
        return
    user_df = df[df["To User"].astype(str).str.lower() == str(username).lower()].copy()
    if user_df.empty:
        st.info("No notifications assigned to you.")
        return
    unread = user_df[user_df.get("Is Read", "No").astype(str).str.lower() != "yes"].copy()
    c1, c2, c3 = st.columns(3)
    c1.metric("Unread Notifications", len(unread))
    c2.metric("Total Notifications", len(user_df))
    cats = sorted([x for x in user_df.get("Category", pd.Series(dtype=str)).astype(str).unique() if x])
    selected_cat = c3.selectbox("Category", ["All"] + cats)
    if selected_cat != "All":
        user_df = user_df[user_df["Category"].astype(str).eq(selected_cat)]
    f1, f2 = st.columns(2)
    with f1:
        unread_only = st.checkbox("Unread only", value=False)
    with f2:
        if st.button("Mark all as read", use_container_width=True):
            df.loc[df["To User"].astype(str).str.lower() == str(username).lower(), "Is Read"] = "Yes"
            df.to_csv(NOTIFICATIONS_PATH, index=False, encoding="utf-8-sig")
            st.cache_data.clear()
            st.rerun()
    if unread_only:
        user_df = user_df[user_df.get("Is Read", "No").astype(str).str.lower() != "yes"]
    st.dataframe(user_df.sort_values("Created At", ascending=False), use_container_width=True, hide_index=True)


def executive_daily_digest_page() -> None:
    render_back_to_dashboard_button("executive_daily_digest")
    st.title("📩 Executive Daily Digest")
    ensure_governance_files()
    date_text = st.date_input("Digest Date", value=ksa_now().date()).strftime("%Y-%m-%d")
    digest = create_digest_for_date(date_text)
    st.text_area("Digest Preview", digest, height=420)
    c1, c2 = st.columns(2)
    with c1:
        if st.button("💾 Save Digest", use_container_width=True):
            row = {"Digest ID": hashlib.md5(f"{date_text}-{digest}".encode()).hexdigest()[:12], "Created At": ksa_now().strftime("%Y-%m-%d %H:%M:%S"), "Created By": st.session_state.get("username", ""), "Audience": "Executive", "Digest Text": digest}
            _append_csv_rows(DAILY_DIGEST_PATH, [row])
            for u in get_users().keys():
                if notification_access_enabled(u, "Daily Digest"):
                    create_system_notification(u, "Executive Daily Digest", digest[:900], "Daily Digest")
            st.success("Digest saved and notifications created for enabled users.")
    with c2:
        if st.button("🟢 Queue WhatsApp Digest", use_container_width=True):
            rows = []
            now = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
            access = safe_read_csv(NOTIFICATION_ACCESS_PATH)
            users = list(get_users().keys())
            for u in users:
                whats_enabled = False
                if not access.empty and "Username" in access.columns and "WhatsApp" in access.columns:
                    m = access[access["Username"].astype(str).str.lower().eq(str(u).lower())]
                    whats_enabled = m["WhatsApp"].astype(str).str.lower().isin(["yes","true","1","on"]).any()
                if whats_enabled or str(u).lower() == "ahmedfekry":
                    rows.append({"Message ID": hashlib.md5(f"{now}-{u}-digest".encode()).hexdigest()[:12], "Created At": now, "To User": u, "To WhatsApp": "", "Message": digest, "Status": "Queued", "Sent At": ""})
            _append_csv_rows(WHATSAPP_OUTBOX_PATH, rows)
            st.success(f"Queued {len(rows)} WhatsApp digest messages.")
    saved = safe_read_csv(DAILY_DIGEST_PATH)
    if not saved.empty:
        st.subheader("Saved Digests")
        st.dataframe(saved.sort_values("Created At", ascending=False), use_container_width=True, hide_index=True)


def whatsapp_agent_page() -> None:
    render_back_to_dashboard_button("whatsapp_agent")
    st.title("🟢 WhatsApp Agent")
    st.caption("Pilot outbox: messages are prepared here. Actual WhatsApp sending requires a connected WhatsApp Business API or approved integration.")
    ensure_governance_files()
    users = list(get_users().keys())
    to_user = st.selectbox("Recipient", users)
    message = st.text_area("Message", placeholder="Write or paste WhatsApp update message...")
    if st.button("Queue WhatsApp Message", use_container_width=True):
        row = {"Message ID": hashlib.md5(f"{ksa_now()}-{to_user}-{message}".encode()).hexdigest()[:12], "Created At": ksa_now().strftime("%Y-%m-%d %H:%M:%S"), "To User": to_user, "To WhatsApp": "", "Message": message, "Status": "Queued", "Sent At": ""}
        _append_csv_rows(WHATSAPP_OUTBOX_PATH, [row])
        st.success("Message queued in WhatsApp outbox.")
    outbox = safe_read_csv(WHATSAPP_OUTBOX_PATH)
    st.dataframe(outbox.sort_values("Created At", ascending=False) if not outbox.empty and "Created At" in outbox.columns else outbox, use_container_width=True, hide_index=True)


def upload_data_page() -> None:
    st.title("📤 Upload CSV Data")
    st.caption("Upload new CSV files directly from the browser. Existing files are backed up before replacement.")

    if not can("upload"):
        st.error("You do not have permission to upload data.")
        return

    for label, path in DATA_FILES.items():
        with st.container(border=True):
            st.subheader(label)
            current_size = path.stat().st_size / 1024 if path.exists() else 0
            st.caption(f"Current file: {path} | Size: {current_size:,.1f} KB")

            uploaded = st.file_uploader(f"Upload replacement for {label}", type=["csv"], key=f"upload_{label}")

            if uploaded is not None:
                if st.button(f"Apply update: {label}", key=f"apply_{label}", use_container_width=True):
                    backup_file(path)
                    path.write_bytes(uploaded.read())
                    st.cache_data.clear()
                    st.success(f"{label} updated successfully. Previous version was saved in backups.")
                    st.rerun()

    st.divider()
    st.subheader("Previous Versions / Backups")
    backups = sorted(BACKUP_DIR.glob("*"), reverse=True)
    if not backups:
        st.info("No backups yet.")
    else:
        rows = []
        for b in backups:
            rows.append({"File": b.name, "Size KB": round(b.stat().st_size / 1024, 1), "Modified": datetime.fromtimestamp(b.stat().st_mtime).strftime("%Y-%m-%d %H:%M:%S")})
        st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)


def smart_alerts_dataframe() -> pd.DataFrame:
    wo = load_workorders()
    if wo.empty:
        return pd.DataFrame()

    alerts = []

    for _, r in wo.iterrows():
        link = r.get("_link", "")
        wo_id = r.get("_wo", "")
        cost = float(r.get("_cost", 0) or 0)
        progress = float(r.get("_progress", 0) or 0)
        stage = str(r.get("_stage", ""))
        sor = str(r.get("_sor", ""))
        status = str(r.get("_status", ""))

        sor_missing = ("created" not in sor.lower()) or ("not" in sor.lower())
        if progress >= 80 and sor_missing:
            alerts.append({
                "Priority": "High",
                "Alert Type": "Revenue Leakage / SOR Missing",
                "Link Code": link,
                "Work Order": wo_id,
                "Cost": cost,
                "Progress": progress,
                "Reason": "Completion >= 80% but SOR is not created.",
                "Required Action": "Commercial / Back Office to create SOR and unblock billing.",
            })

        if "civil" in stage.lower() and ("not" in stage.lower() or "start" in stage.lower()):
            alerts.append({
                "Priority": "High" if cost >= 1_000_000 else "Medium",
                "Alert Type": "Civil Not Start",
                "Link Code": link,
                "Work Order": wo_id,
                "Cost": cost,
                "Progress": progress,
                "Reason": "Civil stage has not started.",
                "Required Action": "Project team to confirm permits, materials, and execution crew readiness.",
            })

        if progress >= 100 and status.lower() != "closed":
            alerts.append({
                "Priority": "Medium",
                "Alert Type": "Completed Not Closed",
                "Link Code": link,
                "Work Order": wo_id,
                "Cost": cost,
                "Progress": progress,
                "Reason": "Work order reached 100% but is not closed.",
                "Required Action": "PMO / Back Office to close administrative and documentation cycle.",
            })

    if not alerts:
        return pd.DataFrame()

    df = pd.DataFrame(alerts)
    order = {"High": 0, "Medium": 1, "Low": 2}
    df["_order"] = df["Priority"].map(order).fillna(9)
    return df.sort_values(["_order", "Cost"], ascending=[True, False]).drop(columns=["_order"])


def smart_alerts_page() -> None:
    if not can("alerts"):
        st.error("You do not have permission to access Smart Alerts.")
        return

    st.title("🚨 Smart Alerts Dashboard")
    st.caption("Early-warning executive module based on current uploaded CSV data.")

    df = smart_alerts_dataframe()
    if df.empty:
        st.success("No critical alerts under current data.")
        return

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Alerts", f"{len(df):,}")
    c2.metric("High Priority", f"{(df['Priority'] == 'High').sum():,}")
    c3.metric("Exposure Cost", f"{df['Cost'].sum():,.0f}")
    c4.metric("Affected Link Codes", f"{df['Link Code'].nunique():,}")

    priority_filter = st.multiselect("Priority", sorted(df["Priority"].unique()), default=sorted(df["Priority"].unique()))
    type_filter = st.multiselect("Alert Type", sorted(df["Alert Type"].unique()), default=sorted(df["Alert Type"].unique()))
    filtered = df[df["Priority"].isin(priority_filter) & df["Alert Type"].isin(type_filter)]

    st.dataframe(filtered, use_container_width=True, hide_index=True)

    csv = filtered.to_csv(index=False).encode("utf-8-sig")
    st.download_button("Download Smart Alerts CSV", csv, "smart_alerts.csv", "text/csv", use_container_width=True)

    if can("email"):
        st.divider()
        st.subheader("Email Alerts")
        recipients = st.text_input("Recipients", placeholder="name@company.com, name2@company.com")
        if st.button("Send Critical Alerts Email", use_container_width=True):
            ok, msg = send_alert_email(recipients, filtered)
            if ok:
                st.success(msg)
            else:
                st.error(msg)


def send_alert_email(recipients: str, df: pd.DataFrame) -> Tuple[bool, str]:
    try:
        smtp_host = st.secrets["email"]["smtp_host"]
        smtp_port = int(st.secrets["email"].get("smtp_port", 587))
        smtp_user = st.secrets["email"]["smtp_user"]
        smtp_password = st.secrets["email"]["smtp_password"]
        sender = st.secrets["email"].get("sender", smtp_user)
    except Exception:
        return False, "Email secrets are not configured. Add [email] settings in Streamlit Secrets."

    recips = [r.strip() for r in recipients.split(",") if r.strip()]
    if not recips:
        return False, "Please enter at least one recipient."

    high = df[df["Priority"] == "High"].head(20)
    body = "Dawiyat PMO Smart Alerts\n\n"
    body += f"Total Alerts: {len(df)}\nHigh Priority: {(df['Priority'] == 'High').sum()}\nExposure Cost: {df['Cost'].sum():,.0f} SAR\n\n"
    body += "Top Critical Alerts:\n"
    for _, r in high.iterrows():
        body += f"- {r['Link Code']} | {r['Alert Type']} | Cost {r['Cost']:,.0f} | Action: {r['Required Action']}\n"

    msg = EmailMessage()
    msg["Subject"] = "Dawiyat PMO Smart Alerts - Critical Sites"
    msg["From"] = sender
    msg["To"] = ", ".join(recips)
    msg.set_content(body)

    with smtplib.SMTP(smtp_host, smtp_port) as server:
        server.starttls()
        server.login(smtp_user, smtp_password)
        server.send_message(msg)

    return True, "Email sent successfully."


def ai_assistant_page() -> None:
    if not can("assistant"):
        st.error("You do not have permission to access AI Assistant.")
        return

    st.title("🤖 AI Executive Assistant بالعربي")
    st.caption("Connected to the latest uploaded CSV files.")

    wo = load_workorders()
    if wo.empty:
        st.warning("No work order data found.")
        return

    metrics = portfolio_metrics()
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Link Codes", f"{metrics['links']:,}")
    c2.metric("Work Orders", f"{metrics['wos']:,}")
    c3.metric("WO Cost", f"{metrics['cost']:,.0f}")
    c4.metric("Avg Progress", f"{metrics['progress']:.1f}%")

    q = st.selectbox(
        "اختر سؤال الإدارة التنفيذية",
        [
            "ما أعلى المواقع الجاهزة أو المتقدمة ولم يتم إنشاء SOR لها؟",
            "ما أعلى مواقع Civil Not Start من حيث التكلفة؟",
            "ما الأعمال التي وصلت 100% ولم تغلق؟",
            "ما أكبر عوائق الفوترة الحالية؟",
            "ما المواقع التي يجب تصعيدها اليوم؟",
            "ماذا يحدث إذا لم يتم اتخاذ إجراء خلال 30 يوم؟",
        ],
    )

    alerts = smart_alerts_dataframe()
    st.markdown('<div class="arabic-box">', unsafe_allow_html=True)

    if "SOR" in q:
        df = alerts[alerts["Alert Type"].eq("Revenue Leakage / SOR Missing")] if not alerts.empty else pd.DataFrame()
        st.subheader("Revenue Leakage / SOR Not Created")
        st.write("هذه المواقع متقدمة في الإنجاز ولكنها لا تدخل دورة الفوترة بسبب عدم إنشاء SOR.")
        st.dataframe(df.head(20), use_container_width=True, hide_index=True)
    elif "Civil Not Start" in q:
        df = alerts[alerts["Alert Type"].eq("Civil Not Start")] if not alerts.empty else pd.DataFrame()
        st.subheader("Civil Not Start Escalation")
        st.write("هذه المواقع تحتاج ضغط مباشر على إدارة المشروع لبدء التنفيذ وإزالة عوائق التصاريح أو الموارد أو المواد.")
        st.dataframe(df.head(20), use_container_width=True, hide_index=True)
    elif "100%" in q:
        df = alerts[alerts["Alert Type"].eq("Completed Not Closed")] if not alerts.empty else pd.DataFrame()
        st.subheader("Completed But Not Closed")
        st.write("هذه الأعمال وصلت 100% ولكنها تحتاج إغلاق إداري أو مستندي لتحرير دورة الفوترة.")
        st.dataframe(df.head(20), use_container_width=True, hide_index=True)
    elif "عوائق الفوترة" in q:
        st.subheader("Billing Bottlenecks")
        if not alerts.empty:
            summary = alerts.groupby("Alert Type").agg(Count=("Link Code", "count"), Exposure=("Cost", "sum")).reset_index().sort_values("Exposure", ascending=False)
            st.dataframe(summary, use_container_width=True, hide_index=True)
    elif "تصعيدها" in q:
        st.subheader("Executive Escalation Queue")
        st.write("تم ترتيب هذه القائمة حسب أولوية التعرض التجاري، عدم إنشاء SOR، Civil Not Start، والإغلاق.")
        st.dataframe(alerts.head(20), use_container_width=True, hide_index=True)
    else:
        exposure = alerts["Cost"].sum() if not alerts.empty else 0
        st.subheader("30 Days Executive Impact")
        st.write(f"إذا لم يتم اتخاذ إجراء خلال 30 يوم، ستظل قيمة تقريبية قدرها {exposure:,.0f} SAR معرضة للتجميد أو التأخير في الفوترة والإغلاق.")
        st.write("القرار المطلوب: War Room يومي للفوترة والإغلاق و Civil Not Start.")
    st.markdown("</div>", unsafe_allow_html=True)



def get_drive_service():
    """Create a Google Drive API service using Streamlit Secrets.

    Supports either [google_service_account] or [gcp_service_account]
    because Streamlit examples often use the second name.
    """
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
    except Exception as exc:
        raise RuntimeError("Google Drive libraries are not installed. Install requirements.txt first.") from exc

    info = None
    for secret_key in ["google_service_account", "gcp_service_account"]:
        try:
            candidate = st.secrets.get(secret_key, None)
            if candidate:
                info = dict(candidate)
                break
        except Exception:
            pass
    if not info:
        raise RuntimeError("Google service account secrets are not configured. Add [google_service_account] or [gcp_service_account] in Streamlit Secrets.")

    if "private_key" in info:
        # Streamlit Secrets uses TOML, while Google exports JSON. Users may paste the key
        # as triple-quoted TOML, or as a JSON-style value containing escaped \\n.
        # Normalize both formats before passing it to google-auth.
        pk = str(info["private_key"]).strip()
        pk = pk.strip('"').strip("'")
        pk = pk.replace("\\r\\n", "\n").replace("\\n", "\n").replace("\r\n", "\n")
        pk = pk.replace("-----BEGIN PRIVATE KEY----- ", "-----BEGIN PRIVATE KEY-----\n")
        pk = pk.replace(" -----END PRIVATE KEY-----", "\n-----END PRIVATE KEY-----")
        pk = re.sub(r"-----BEGIN PRIVATE KEY-----\s*", "-----BEGIN PRIVATE KEY-----\n", pk)
        pk = re.sub(r"\s*-----END PRIVATE KEY-----", "\n-----END PRIVATE KEY-----", pk)
        info["private_key"] = pk
    scopes = ["https://www.googleapis.com/auth/drive"]
    try:
        credentials = service_account.Credentials.from_service_account_info(info, scopes=scopes)
    except Exception as exc:
        raise RuntimeError(
            "Google Drive Service Account private_key is not valid. "
            "In Streamlit Secrets use TOML format. Example: "
            "private_key = \"\"\"-----BEGIN PRIVATE KEY-----\\n...\\n-----END PRIVATE KEY-----\\n\"\"\". "
            "Do not paste JSON with braces or commas."
        ) from exc
    return build("drive", "v3", credentials=credentials, cache_discovery=False)


def get_drive_root_folder_id() -> str:
    try:
        return str(st.secrets["google_drive"]["root_folder_id"]).strip()
    except Exception:
        return ""


def google_drive_root_config_message() -> str:
    return (
        "google_drive.root_folder_id is not configured in Streamlit Secrets. "
        "Add the Google Drive folder ID for the main 'Link Codes' folder under [google_drive]. "
        "The current version uses MANUAL Google Drive upload mode: users open the Link Code folder "
        "and upload files directly in Google Drive, while the dashboard scans document status."
    )


def drive_folder_url(folder_id: str) -> str:
    return f"https://drive.google.com/drive/folders/{folder_id}" if folder_id else ""


def extract_drive_folder_id(url_or_id: str) -> str:
    text = str(url_or_id or "").strip()
    if not text:
        return ""
    if re.fullmatch(r"[A-Za-z0-9_-]{20,}", text):
        return text
    for pattern in [r"/folders/([A-Za-z0-9_-]+)", r"[?&]id=([A-Za-z0-9_-]+)", r"/d/([A-Za-z0-9_-]+)"]:
        m = re.search(pattern, text)
        if m:
            return m.group(1)
    return ""


def drive_escape(value: str) -> str:
    return str(value or "").replace("'", "\'")


def find_drive_folder(service, parent_id: str, name: str) -> str:
    query = (
        f"'{drive_escape(parent_id)}' in parents and "
        f"mimeType = '{GOOGLE_DRIVE_FOLDER_MIME}' and "
        f"name = '{drive_escape(name)}' and trashed = false"
    )
    result = service.files().list(
        q=query,
        spaces="drive",
        fields="files(id,name)",
        pageSize=10,
        supportsAllDrives=True,
        includeItemsFromAllDrives=True,
    ).execute()
    files = result.get("files", [])
    return files[0]["id"] if files else ""


def create_drive_folder(service, parent_id: str, name: str) -> str:
    metadata = {"name": name, "mimeType": GOOGLE_DRIVE_FOLDER_MIME, "parents": [parent_id]}
    folder = service.files().create(body=metadata, fields="id", supportsAllDrives=True).execute()
    return folder.get("id", "")


def get_or_create_drive_folder(service, parent_id: str, name: str) -> str:
    existing = find_drive_folder(service, parent_id, name)
    return existing or create_drive_folder(service, parent_id, name)


def ensure_link_folder(service, link_code: str, existing_link: str = "") -> tuple:
    """Resolve the existing Link Code folder without uploading/creating files.

    Manual Upload Mode avoids the Service Account storage quota issue.
    It first uses Document_Link from the CSV. If missing, it searches the configured
    Link Codes root folder by Link Code name and stores the found URL back to CSV.
    It does not create folders or upload files.
    """
    folder_id = extract_drive_folder_id(existing_link)
    if folder_id:
        return folder_id, drive_folder_url(folder_id)
    root_id = get_drive_root_folder_id()
    if not root_id:
        raise RuntimeError(google_drive_root_config_message())
    folder_id = find_drive_folder(service, root_id, link_code)
    if folder_id:
        return folder_id, drive_folder_url(folder_id)
    return "", ""


def upload_file_to_drive(service, folder_id: str, uploaded_file, filename: str, mime_type: str = "application/octet-stream") -> str:
    try:
        from googleapiclient.http import MediaIoBaseUpload
    except Exception as exc:
        raise RuntimeError("google-api-python-client is not installed.") from exc
    uploaded_file.seek(0)
    media = MediaIoBaseUpload(uploaded_file, mimetype=mime_type, resumable=True)
    metadata = {"name": filename, "parents": [folder_id]}
    created = service.files().create(
        body=metadata,
        media_body=media,
        fields="id,webViewLink",
        supportsAllDrives=True,
    ).execute()
    return created.get("webViewLink", "")



def safe_drive_filename(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_.\-()\[\] ]+", "_", str(value or "")).strip()
    return cleaned or "uploaded_file"


def document_folder_name(doc_type: str) -> str:
    return DOCUMENT_FOLDER_MAP.get(str(doc_type), str(doc_type))


def ensure_document_subfolders(service, link_folder_id: str) -> Dict[str, str]:
    """Return existing standard Option A subfolders under a Link Code folder.

    In Manual Upload Mode the dashboard does not create folders with the
    Service Account because Service Accounts have no personal Drive storage quota.
    Users can create/upload directly inside Google Drive; this function only detects
    existing folders/files and updates the status preview.
    """
    folder_ids = {}
    for doc_type in DOCUMENT_TYPES:
        folder_ids[doc_type] = find_drive_folder(service, link_folder_id, document_folder_name(doc_type))
    return folder_ids


def list_drive_files(service, folder_id: str) -> List[Dict[str, str]]:
    query = f"'{drive_escape(folder_id)}' in parents and mimeType != '{GOOGLE_DRIVE_FOLDER_MIME}' and trashed = false"
    result = service.files().list(
        q=query,
        spaces="drive",
        fields="files(id,name,webViewLink,createdTime,modifiedTime,size,mimeType)",
        pageSize=100,
        supportsAllDrives=True,
        includeItemsFromAllDrives=True,
        orderBy="modifiedTime desc",
    ).execute()
    return result.get("files", [])


def document_status_for_link(service, link_folder_id: str) -> Dict[str, Dict[str, object]]:
    """Return status/count/latest link for the standard document subfolders."""
    status = {}
    for doc_type in DOCUMENT_TYPES:
        folder_name = document_folder_name(doc_type)
        folder_id = find_drive_folder(service, link_folder_id, folder_name)
        files = list_drive_files(service, folder_id) if folder_id else []
        status[doc_type] = {
            "folder_id": folder_id,
            "folder_url": drive_folder_url(folder_id),
            "count": len(files),
            "latest_file": files[0].get("name", "") if files else "",
            "latest_url": files[0].get("webViewLink", "") if files else "",
            "latest_created": files[0].get("createdTime", "") if files else "",
            "latest_modified": files[0].get("modifiedTime", "") if files else "",
            "state": "Uploaded" if files else "Missing",
        }
    return status


def status_badge_text(state: str, count: int = 0) -> str:
    if state == "Uploaded":
        return "🟢 Uploaded"
    if state == "Partial":
        return "🟡 Partial"
    return "🔴 Missing"


def get_document_link_for_link(wo: pd.DataFrame, link_code: str) -> str:
    if wo.empty or "Document_Link" not in wo.columns:
        return ""
    vals = wo.loc[wo["Link Code"].astype(str).eq(str(link_code)), "Document_Link"].astype(str)
    vals = [v.strip() for v in vals if str(v).strip()]
    return vals[0] if vals else ""


def update_document_link_in_csv(link_code: str, folder_url: str) -> None:
    wo = master_workorders_df()
    if wo.empty:
        raise RuntimeError("u_osp_work_order.csv is empty or missing.")
    if "Document_Link" not in wo.columns:
        wo["Document_Link"] = ""
    if "Link Code" not in wo.columns:
        raise RuntimeError("Link Code column is missing from u_osp_work_order.csv.")
    mask = wo["Link Code"].astype(str).eq(str(link_code))
    if not mask.any():
        raise RuntimeError(f"Link Code not found in CSV: {link_code}")
    backup_file(WO_PATH)
    wo.loc[mask, "Document_Link"] = folder_url
    wo.to_csv(WO_PATH, index=False, encoding="utf-8-sig")
    st.cache_data.clear()




def get_all_link_codes(wo: pd.DataFrame) -> List[str]:
    if wo.empty or "Link Code" not in wo.columns:
        return []
    return sorted([x for x in wo["Link Code"].astype(str).str.strip().unique() if x and x.lower() != "nan"])


def build_document_status_rows(service, wo: pd.DataFrame, link_codes: List[str], max_links: int = 50) -> Tuple[pd.DataFrame, Dict[str, str]]:
    """Scan Google Drive document folders and return one status row per Link Code.
    The scan is intentionally limited by max_links to avoid slow Google API calls.
    """
    rows = []
    folder_urls = {}
    for link_code in link_codes[:max_links]:
        current_folder_url = get_document_link_for_link(wo, link_code)
        try:
            link_folder_id, link_folder_url = ensure_link_folder(service, link_code, current_folder_url)
            if link_folder_url and not current_folder_url:
                update_document_link_in_csv(link_code, link_folder_url)
            if not link_folder_id:
                raise RuntimeError("No existing Google Drive folder found for this Link Code. Create the folder manually under the Link Codes root folder or add Document_Link to the CSV.")
            status = document_status_for_link(service, link_folder_id)
            folder_urls[link_code] = link_folder_url
            uploaded_count = sum(1 for d in DOCUMENT_TYPES if status.get(d, {}).get("state") == "Uploaded")
            row = {
                "Link Code": link_code,
                "Folder Link": link_folder_url,
                "Uploaded Types": uploaded_count,
                "Missing Types": len(DOCUMENT_TYPES) - uploaded_count,
                "Overall Status": "Uploaded" if uploaded_count == len(DOCUMENT_TYPES) else ("Partial" if uploaded_count > 0 else "Missing"),
                "Latest File": "",
                "Latest Created": "",
                "Latest Modified": "",
            }
            latest_candidates = []
            for doc_type in DOCUMENT_TYPES:
                info = status.get(doc_type, {})
                row[doc_type] = status_badge_text(str(info.get("state", "Missing")), int(info.get("count", 0) or 0))
                row[f"{document_folder_name(doc_type)} Uploaded / Created Date"] = str(info.get("latest_created", ""))
                row[f"{document_folder_name(doc_type)} Modified Date"] = str(info.get("latest_modified", ""))
                if info.get("latest_modified"):
                    latest_candidates.append((str(info.get("latest_modified", "")), str(info.get("latest_created", "")), str(info.get("latest_file", ""))))
            if latest_candidates:
                latest_candidates.sort(reverse=True)
                row["Latest Modified"], row["Latest Created"], row["Latest File"] = latest_candidates[0]
            rows.append(row)
        except Exception as exc:
            row = {
                "Link Code": link_code,
                "Folder Link": current_folder_url,
                "Uploaded Types": 0,
                "Missing Types": len(DOCUMENT_TYPES),
                "Overall Status": "Error",
                "Latest File": "",
                "Latest Created": "",
                "Latest Modified": "",
                "Error": str(exc),
            }
            for doc_type in DOCUMENT_TYPES:
                row[doc_type] = "❌ Missing"
                row[f"{document_folder_name(doc_type)} Uploaded / Created Date"] = ""
                row[f"{document_folder_name(doc_type)} Modified Date"] = ""
            rows.append(row)
    return pd.DataFrame(rows), folder_urls


def document_kpi_cards(status_df: pd.DataFrame) -> None:
    if status_df.empty:
        st.info("No document status rows scanned yet.")
        return
    total_links = len(status_df)
    complete_links = int((status_df["Overall Status"] == "Uploaded").sum()) if "Overall Status" in status_df.columns else 0
    partial_links = int((status_df["Overall Status"] == "Partial").sum()) if "Overall Status" in status_df.columns else 0
    missing_links = int((status_df["Overall Status"] == "Missing").sum()) if "Overall Status" in status_df.columns else 0
    total_uploaded_types = int(status_df.get("Uploaded Types", pd.Series(dtype=int)).sum())
    total_missing_types = int(status_df.get("Missing Types", pd.Series(dtype=int)).sum())

    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("Scanned Link Codes", f"{total_links:,}")
    c2.metric("Fully Uploaded", f"{complete_links:,}")
    c3.metric("Partial", f"{partial_links:,}")
    c4.metric("Missing", f"{missing_links:,}")
    c5.metric("Missing Doc Types", f"{total_missing_types:,}", help=f"Uploaded doc types: {total_uploaded_types:,}")


def upload_widget_for_document_type(service, link_code: str, link_folder_id: str, doc_type: str, doc_status: Dict[str, Dict[str, object]]) -> None:
    """Manual upload card: open the correct Google Drive subfolder and scan status.

    Direct upload from Streamlit is intentionally disabled to avoid the Google Drive
    Service Account storage-quota limitation on normal My Drive accounts.
    """
    folder_label = document_folder_name(doc_type)
    info = doc_status.get(doc_type, {})
    with st.container(border=True):
        left, right = st.columns([1, .55])
        with left:
            st.markdown(f"### {folder_label}")
            st.caption(status_badge_text(str(info.get("state", "Missing")), int(info.get("count", 0) or 0)))
            if info.get("latest_file"):
                st.caption(f"Latest: {info.get('latest_file')}")
                st.caption(f"Uploaded / Created: {info.get('latest_created', '')}")
                st.caption(f"Modified: {info.get('latest_modified', '')}")
        with right:
            if info.get("folder_url"):
                st.link_button("Open Subfolder", str(info.get("folder_url")), use_container_width=True)
            elif link_folder_id:
                st.info("Subfolder not found. Create it manually inside the Link Code folder.")
            if info.get("latest_url"):
                st.link_button("Open Latest File", str(info.get("latest_url")), use_container_width=True)



def document_upload_page() -> None:
    if not can("documents"):
        st.error("You do not have permission to access Document Upload Center.")
        return

    top_left, top_right = st.columns([1, .22])
    with top_left:
        st.title("📂 Document Upload Center")
        st.caption("Manual Google Drive upload workflow for every Link Code. Open the Link Code folder, upload files directly into: 01 Design / 02 Permit / 03 Photos / 04 PAT / 05 AsBuilt / 06 Handover / 07 Commercial, then refresh status.")
        st.info("Manual upload mode: upload files in Google Drive, then click Refresh Document Status. Created/Uploaded Date and Modified Date are scanned from Google Drive for all 7 stages.")
    with top_right:
        st.write("")
        st.write("")
        if st.button("⬅ Back to Dashboard", use_container_width=True):
            st.session_state.pop("active_hidden_page", None)
            st.session_state["force_dashboard"] = True
            st.rerun()

    if not can("documents"):
        st.error("You do not have permission to access documents.")
        return

    wo = master_workorders_df()
    if wo.empty:
        st.warning("u_osp_work_order.csv is missing or empty.")
        return
    if "Link Code" not in wo.columns:
        st.error("Link Code column is missing from u_osp_work_order.csv.")
        return

    links = get_all_link_codes(wo)
    if not links:
        st.warning("No Link Codes found.")
        return

    try:
        service = get_drive_service()
        drive_connected = True
    except Exception as exc:
        service = None
        drive_connected = False
        st.error(str(exc))
        st.info("Check Streamlit Secrets, enable Google Drive API, and share the Link Codes Google Drive folder with the service account email as Viewer/Editor so the dashboard can scan files.")

    st.markdown("### Executive Documents Dashboard")
    with st.container(border=True):
        scan_col1, scan_col2, scan_col3 = st.columns([1.2, .6, .6])
        with scan_col1:
            scan_scope = st.multiselect("Scan Link Codes", links, default=links[:min(10, len(links))], help="Scanning all Link Codes may take time because each scan checks Google Drive folders.")
        with scan_col2:
            max_scan = st.number_input("Max Scan", min_value=1, max_value=500, value=min(50, len(links)), step=10)
        with scan_col3:
            st.metric("Drive", "Connected" if drive_connected else "Not Connected")
        if drive_connected and st.button("Refresh Document Status", use_container_width=True, type="secondary"):
            with st.spinner("Scanning Google Drive folders..."):
                status_df, folder_urls = build_document_status_rows(service, wo, scan_scope or links, int(max_scan))
                st.session_state["document_status_df"] = status_df
                st.session_state["document_folder_urls"] = folder_urls
                try:
                    status_df.to_csv(DOC_STATUS_CACHE_PATH, index=False, encoding="utf-8-sig")
                    st.cache_data.clear()
                except Exception:
                    pass

        status_df = st.session_state.get("document_status_df", pd.DataFrame())
        document_kpi_cards(status_df)
        if not status_df.empty:
            date_cols = []
            for _doc_type in DOCUMENT_TYPES:
                date_cols.extend([f"{document_folder_name(_doc_type)} Uploaded / Created Date", f"{document_folder_name(_doc_type)} Modified Date"])
            show_cols = ["Link Code", "Overall Status", "Uploaded Types", "Missing Types"] + DOCUMENT_TYPES + date_cols + ["Latest File", "Latest Created", "Latest Modified", "Folder Link"]
            available_cols = [c for c in show_cols if c in status_df.columns]
            st.dataframe(status_df[available_cols], use_container_width=True, hide_index=True)
            st.download_button(
                "Export Document Status Excel/CSV",
                status_df.to_csv(index=False).encode("utf-8-sig"),
                "Dawiyat_Document_Status.csv",
                "text/csv",
                use_container_width=True,
            )

    st.divider()
    st.markdown("### Open Link Code Folder & Refresh Document Status")
    st.caption("Select one Link Code, open its Google Drive folder, upload files manually in Google Drive, then refresh the status scan. This avoids Service Account storage quota errors.")

    c1, c2, c3 = st.columns([1.35, .75, .75])
    with c1:
        # Keep selected Link Code stable during upload/page reruns.
        if st.session_state.get("doc_upload_link_code") not in links:
            st.session_state["doc_upload_link_code"] = links[0]
        link_code = st.selectbox("Link Code", links, key="doc_upload_link_code")
    with c2:
        st.metric("Document Types", len(DOCUMENT_TYPES))
    with c3:
        st.metric("Mode", "Manual Drive Upload")

    current_folder_url = get_document_link_for_link(wo, link_code)
    link_folder_id = ""
    link_folder_url = current_folder_url
    doc_status = {d: {"state": "Missing", "count": 0, "latest_file": "", "latest_url": "", "folder_url": ""} for d in DOCUMENT_TYPES}

    if drive_connected:
        try:
            link_folder_id, link_folder_url = ensure_link_folder(service, link_code, current_folder_url)
            if link_folder_url and not current_folder_url:
                update_document_link_in_csv(link_code, link_folder_url)
                current_folder_url = link_folder_url
                wo = master_workorders_df()
            if link_folder_id:
                doc_status = document_status_for_link(service, link_folder_id)
            else:
                st.warning("No existing Google Drive folder found for this Link Code. Create it manually under the Link Codes root folder, or add Document_Link to the CSV.")
        except Exception as exc:
            # This usually means root_folder_id is missing AND the selected Link Code has no Document_Link.
            st.warning(str(exc))
            if current_folder_url:
                st.info("This Link Code has an existing Document_Link, but the folder ID could not be read. Check the link format.")
            else:
                st.info("Select a Link Code that already has Documents = Open, or create a folder with the exact Link Code name under the Link Codes root folder. The dashboard will find it on the next refresh.")

    with st.container(border=True):
        h1, h2 = st.columns([1, .35])
        with h1:
            st.subheader(f"📁 {link_code}")
            st.caption("Manual upload mode: open this folder in Google Drive and upload files into the standard Option A subfolders.")
        with h2:
            if link_folder_url:
                st.link_button("Open Link Code Folder", link_folder_url, use_container_width=True)
            else:
                st.button("No Folder Link", disabled=True, use_container_width=True)

        metrics = st.columns(len(DOCUMENT_TYPES))
        for i, doc_type in enumerate(DOCUMENT_TYPES):
            info = doc_status.get(doc_type, {})
            metrics[i].metric(doc_type, status_badge_text(str(info.get("state", "Missing")), int(info.get("count", 0) or 0)))

    st.warning("Direct Streamlit upload is disabled in this version. Upload files manually inside Google Drive, then click Refresh Document Status.")

    for row_start in range(0, len(DOCUMENT_TYPES), 2):
        cols = st.columns(2)
        for idx, doc_type in enumerate(DOCUMENT_TYPES[row_start:row_start + 2]):
            with cols[idx]:
                upload_widget_for_document_type(service, link_code, link_folder_id, doc_type, doc_status)




def build_pdf_report() -> bytes:
    metrics = portfolio_metrics()
    alerts = smart_alerts_dataframe()

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=landscape(A4), rightMargin=24, leftMargin=24, topMargin=24, bottomMargin=24)
    styles = getSampleStyleSheet()
    story = []

    story.append(Paragraph("Dawiyat PMO Executive Report", styles["Title"]))
    story.append(Paragraph(f"Generated: {ksa_now().strftime('%Y-%m-%d %H:%M')}", styles["Normal"]))
    story.append(Spacer(1, 12))

    summary_data = [
        ["Metric", "Value"],
        ["Link Codes", f"{metrics['links']:,}"],
        ["Work Orders", f"{metrics['wos']:,}"],
        ["WO Cost", f"{metrics['cost']:,.0f} SAR"],
        ["Average Progress", f"{metrics['progress']:.1f}%"],
    ]

    t = Table(summary_data, hAlign="LEFT")
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#10223a")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#d9e3ef")),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("PADDING", (0,0), (-1,-1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 16))

    story.append(Paragraph("Top Smart Alerts", styles["Heading2"]))
    if alerts.empty:
        story.append(Paragraph("No critical alerts.", styles["Normal"]))
    else:
        top = alerts.head(12).copy()
        alert_data = [["Priority", "Alert Type", "Link Code", "Cost", "Progress", "Required Action"]]
        for _, r in top.iterrows():
            alert_data.append([
                str(r["Priority"]),
                str(r["Alert Type"])[:36],
                str(r["Link Code"])[:22],
                f"{float(r['Cost']):,.0f}",
                f"{float(r['Progress']):.1f}%",
                str(r["Required Action"])[:70],
            ])
        at = Table(alert_data, colWidths=[60, 160, 120, 80, 70, 300])
        at.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#1e3a5f")),
            ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#d9e3ef")),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
            ("PADDING", (0,0), (-1,-1), 6),
            ("FONTSIZE", (0,0), (-1,-1), 8),
        ]))
        story.append(at)

    story.append(Spacer(1, 14))
    story.append(Paragraph("Prepared by Eng/Ahmed Fekry - Quality & Performance Director (PMO)", styles["Normal"]))

    doc.build(story)
    return buffer.getvalue()



# -----------------------------------------------------------------------------
# Executive PPT Builder (Streamlit native page)
# This module is intentionally outside dashboard.html to keep the dashboard fast.
# -----------------------------------------------------------------------------
PPT_REPORT_OPTIONS = {
    "portfolio": "Portfolio Summary & Cost Exposure",
    "sor_summary": "SOR Summary & Revenue Exposure",
    "stage_summary": "Overall Stages Summary & Cost Exposure",
    "full_scope": "Dawiyat Project Full Scope",
    "regional": "Regional Performance Summary",
    "completion": "Sites Completion Analysis",
    "cost": "Sites Cost Analysis",
    "monthly": "Monthly Progress Trend",
    "financial": "Executive Financial Report",
    "readiness": "Executive Readiness Report",
    "pmo_audit": "PMO Audit Summary",
    "assistant_insights": "PMO Report Assistant Insights",
}


def _ppt_imports():
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    return Presentation, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, Inches, Pt, RGBColor


def _clean_text(v: Any, default: str = "N/A") -> str:
    s = str(v if v is not None else "").strip()
    return s if s else default


def _fmt_money(v: float) -> str:
    try:
        return f"{float(v):,.0f}"
    except Exception:
        return "0"


def _fmt_pct(v: float) -> str:
    try:
        return f"{float(v):.0f}%"
    except Exception:
        return "0%"


def _status_from_progress(v: float) -> str:
    try:
        x = float(v)
    except Exception:
        x = 0.0
    if x >= 100:
        return "Completed"
    if x > 0:
        return "In Progress"
    return "Not Start"


def _parse_date_any(v: Any):
    s = str(v or "").strip()
    if not s:
        return pd.NaT
    return pd.to_datetime(s, errors="coerce", dayfirst=True)


def load_ppt_workorders() -> pd.DataFrame:
    """Load the same CSV files as the dashboard, but keep it Streamlit-native.
    Includes fields needed for PPT Builder filters and reports.
    """
    wo = apply_derived_billing_fields(apply_project_updates_to_workorders(master_workorders_df())).copy()
    if wo.empty:
        return pd.DataFrame()

    dist = derive_district_records_from_workorders(master_workorders_df()).copy()

    link_col = first_existing_col(wo, ["Link Code"])
    wo_col = first_existing_col(wo, ["Work Order"])
    cost_col = first_existing_col(wo, ["WO Cost", "Cost"])
    progress_col = first_existing_col(wo, ["Percentage of Completion"])
    subclass_col = first_existing_col(wo, ["Subclass"])
    region_col = first_existing_col(wo, ["Region"])
    city_col = first_existing_col(wo, ["City"])
    district_col = first_existing_col(wo, ["District", "District "])
    project_col = first_existing_col(wo, ["Project"])
    stage_col = first_existing_col(wo, ["Stage"])
    sor_col = first_existing_col(wo, ["SOR Status"])
    sor_ref_col = first_existing_col(wo, ["SOR Reference Number"])
    updated_col = first_existing_col(wo, ["Updated", "Created"])
    closure_col = first_existing_col(wo, ["Work Order Status"])
    performance_col = first_existing_col(wo, ["Performance Status", "Performance"])

    out = pd.DataFrame({
        "Link Code": wo[link_col].astype(str) if link_col else "",
        "Work Order": wo[wo_col].astype(str) if wo_col else "",
        "Region": wo[region_col].astype(str) if region_col else "",
        "City": wo[city_col].astype(str) if city_col else "",
        "District": wo[district_col].astype(str) if district_col else "",
        "Project": wo[project_col].astype(str) if project_col else "",
        "Stage": wo[stage_col].astype(str) if stage_col else "",
        "SOR Status": wo[sor_col].astype(str) if sor_col else "",
        "SOR Reference Number": wo[sor_ref_col].astype(str) if sor_ref_col else "",
        "Cost": effective_row_cost_from_wo(wo),
        "Progress": wo[progress_col].apply(parse_num) if progress_col else 0.0,
        "Subclass": wo[subclass_col].astype(str) if subclass_col else "",
        "Updated": wo[updated_col].astype(str) if updated_col else "",
        "Closure Status": wo[closure_col].astype(str) if closure_col else "",
        "Performance": wo[performance_col].astype(str) if performance_col else "",
        "Year": wo[first_existing_col(wo, ["Year"])].astype(str) if first_existing_col(wo, ["Year"]) else "",
        "Work Order Status": wo[first_existing_col(wo, ["Work Order Status"])].astype(str) if first_existing_col(wo, ["Work Order Status"]) else "",
        "Type": wo[first_existing_col(wo, ["Type"])].astype(str) if first_existing_col(wo, ["Type"]) else "",
        "Class": wo[first_existing_col(wo, ["Class"])].astype(str) if first_existing_col(wo, ["Class"]) else "",
        "Scope Target": wo[first_existing_col(wo, ["Scope Target"])].astype(str) if first_existing_col(wo, ["Scope Target"]) else "",
        "1st 50 Invoice Status": wo[first_existing_col(wo, ["1st 50 Invoice Status"])].astype(str) if first_existing_col(wo, ["1st 50 Invoice Status"]) else "",
        "Second 50% status": wo[first_existing_col(wo, ["Second 50% status", "2nd 50 Invoice Status", "Second 50 Invoice Status"])].astype(str) if first_existing_col(wo, ["Second 50% status", "2nd 50 Invoice Status", "Second 50 Invoice Status"]) else "",
        "Missing MET Actual / PM Review": wo[first_existing_col(wo, ["Missing MET Actual / PM Review", "Missing MET Actual", "PM Review", "PM Review Status"])].astype(str) if first_existing_col(wo, ["Missing MET Actual / PM Review", "Missing MET Actual", "PM Review", "PM Review Status"]) else "",
    })

    if not dist.empty:
        d_link = first_existing_col(dist, ["Link Code"])
        d_wo = first_existing_col(dist, ["Work Order"])
        d_region = first_existing_col(dist, ["Region"])
        d_city = first_existing_col(dist, ["City"])
        d_district = first_existing_col(dist, ["District", "District "])
        cols = []
        for c in [d_link, d_wo, d_region, d_city, d_district]:
            if c and c not in cols:
                cols.append(c)
        if d_link and d_wo and cols:
            d = dist[cols].copy()
            rename = {d_link: "Link Code", d_wo: "Work Order"}
            if d_region: rename[d_region] = "Region_map"
            if d_city: rename[d_city] = "City_map"
            if d_district: rename[d_district] = "District_map"
            d = d.rename(columns=rename)
            out = out.merge(d, on=["Link Code", "Work Order"], how="left")
            if "Region_map" in out:
                out["Region"] = out["Region_map"].where(out["Region_map"].astype(str).str.strip().ne(""), out["Region"])
            if "City_map" in out:
                out["City"] = out["City_map"].where(out["City_map"].astype(str).str.strip().ne(""), out["City"])
            if "District_map" in out:
                out["District"] = out["District_map"].where(out["District_map"].astype(str).str.strip().ne(""), out["District"])
            out = out.drop(columns=[c for c in ["Region_map", "City_map", "District_map"] if c in out.columns])

    for c in ["Region", "City", "District", "Project", "Stage", "SOR Status", "SOR Reference Number", "Year", "Work Order Status", "Type", "Class", "Scope Target", "Subclass", "1st 50 Invoice Status", "Second 50% status", "Missing MET Actual / PM Review"]:
        out[c] = out[c].fillna("").astype(str).str.strip()
        out[c] = out[c].replace({"": "N/A", "nan": "N/A", "NaN": "N/A", "None": "N/A"})
    out["Status"] = out["Progress"].apply(_status_from_progress)
    out["Updated_dt"] = out["Updated"].apply(_parse_date_any)
    return out


def link_level_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    rows = _ensure_ppt_columns(rows)
    if rows.empty:
        return pd.DataFrame(columns=["Link Code", "Region", "City", "SOR Status", "Cost", "Progress", "Status"])
    grouped = rows.groupby("Link Code", dropna=False).agg(
        Region=("Region", lambda s: _clean_text(next((x for x in s if str(x).strip()), "N/A"))),
        City=("City", lambda s: _clean_text(next((x for x in s if str(x).strip()), "N/A"))),
        **{"SOR Status": ("SOR Status", lambda s: _clean_text(next((x for x in s if str(x).strip()), "N/A")))},
        Cost=("Cost", "sum"),
        Progress=("Progress", "mean"),
    ).reset_index()
    grouped["Status"] = grouped["Progress"].apply(_status_from_progress)
    return grouped


def city_summary_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    links = link_level_dataframe(rows)
    if links.empty:
        return pd.DataFrame(columns=["Region", "City", "No. of Link Codes", "Completed", "In Progress", "Not Start", "Completion %", "WO Amount"])
    summary = links.groupby(["Region", "City"], dropna=False).agg(
        **{"No. of Link Codes": ("Link Code", "nunique"), "WO Amount": ("Cost", "sum")}
    ).reset_index()
    for status in ["Completed", "In Progress", "Not Start"]:
        counts = links[links["Status"].eq(status)].groupby(["Region", "City"])["Link Code"].nunique()
        summary[status] = summary.set_index(["Region", "City"]).index.map(counts).fillna(0).astype(int)
    summary["Completion %"] = summary.apply(lambda r: (r["Completed"] / r["No. of Link Codes"] * 100) if r["No. of Link Codes"] else 0, axis=1)
    return summary.sort_values("WO Amount", ascending=False)


def status_cost_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    links = link_level_dataframe(rows)
    if links.empty:
        return pd.DataFrame(columns=["Status", "CIVIL", "WO Cost Civil", "FIBRE", "WO Cost Fibre", "Total WO Cost"])
    result = []
    for st_name in ["Completed", "In Progress", "Not Start"]:
        st_rows = rows[rows["Status"].eq(st_name)]
        civil = st_rows[st_rows["Subclass"].str.lower().str.contains("civil", na=False)]
        fiber = st_rows[st_rows["Subclass"].str.lower().str.contains("fiber", na=False)]
        result.append({
            "Status": st_name,
            "CIVIL": civil["Link Code"].nunique(),
            "WO Cost Civil": civil["Cost"].sum(),
            "FIBRE": fiber["Link Code"].nunique(),
            "WO Cost Fibre": fiber["Cost"].sum(),
            "Total WO Cost": st_rows["Cost"].sum(),
        })
    return pd.DataFrame(result)


def monthly_progress_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    if rows.empty or rows["Updated_dt"].isna().all():
        return pd.DataFrame(columns=["Month", "Civil %", "Fiber %", "Overall %"])
    temp = rows.dropna(subset=["Updated_dt"]).copy()
    temp["Month_dt"] = temp["Updated_dt"].dt.to_period("M").dt.to_timestamp()
    result = []
    for month, grp in temp.groupby("Month_dt"):
        civil = grp[grp["Subclass"].str.lower().str.contains("civil", na=False)]
        fiber = grp[grp["Subclass"].str.lower().str.contains("fiber", na=False)]
        result.append({
            "Month": month.strftime("%B %Y"),
            "Civil %": civil["Progress"].mean() if not civil.empty else 0,
            "Fiber %": fiber["Progress"].mean() if not fiber.empty else 0,
            "Overall %": grp["Progress"].mean() if not grp.empty else 0,
        })
    return pd.DataFrame(result).tail(9)


def penalty_total_filtered(rows: pd.DataFrame) -> float:
    pen = safe_read_csv(PENALTIES_PATH)
    if pen.empty:
        return 0.0
    link_col = first_existing_col(pen, ["Cluster Name", "Link Code"])
    amt_col = first_existing_col(pen, ["Penalties Amount", "Penalty Amount"])
    if not link_col or not amt_col:
        return 0.0
    valid_links = set(rows["Link Code"].dropna().astype(str)) if not rows.empty else set()
    p = pen.copy()
    p["_link"] = p[link_col].astype(str).str.strip()
    p["_amount"] = p[amt_col].apply(parse_num)
    p = p[p["_link"].ne("")]
    p = p[~p["_link"].str.lower().isin(["grand total", "total"])]
    if valid_links:
        p = p[p["_link"].isin(valid_links)]
    return float(p["_amount"].sum())


def readiness_summary_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    links = set(rows["Link Code"].dropna().astype(str)) if not rows.empty else set()
    status_df = pd.DataFrame(read_cached_document_status_records())
    out = []
    for doc in DOCUMENT_TYPES:
        uploaded = partial = missing = 0
        if status_df.empty or "Link Code" not in status_df.columns:
            missing = len(links)
        else:
            df = status_df[status_df["Link Code"].astype(str).isin(links)].copy() if links else status_df.copy()
            col = doc if doc in df.columns else first_existing_col(df, [doc])
            if not col:
                missing = len(links)
            else:
                s = df[col].astype(str).str.lower()
                uploaded = int(s.str.contains("uploaded|stage data available|accepted", regex=True).sum())
                partial = int(s.str.contains("partial|under|requested", regex=True).sum())
                missing = max(0, len(links) - uploaded - partial)
        out.append({"Document Type": doc, "Uploaded": uploaded, "Partial": partial, "Missing": missing})
    return pd.DataFrame(out)


def _slide_bg(slide, rgb=(253, 226, 204)):
    _, _, _, _, _, _, MSO_SHAPE, Inches, _, RGBColor = _ppt_imports()
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()
    shape.z_order if hasattr(shape, "z_order") else None


def _add_logo(slide, path: Path, x: float, y: float, w: float, h: float):
    _, _, _, _, _, _, _, Inches, _, _ = _ppt_imports()
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _add_header(slide, title: str):
    _, _, _, _, PP_ALIGN, _, _, Inches, Pt, RGBColor = _ppt_imports()
    _slide_bg(slide)
    _add_logo(slide, MET_LOGO_PATH, 0.25, 0.12, 1.95, 0.85)
    _add_logo(slide, DAWIYAT_LOGO_PATH, 10.55, 0.12, 2.55, 0.75)
    box = slide.shapes.add_textbox(Inches(2.1), Inches(0.52), Inches(9.1), Inches(0.42))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.bold = True; run.font.size = Pt(20); run.font.color.rgb = RGBColor(0, 0, 0)


def _add_footer(slide):
    """Fixed footer on all slides except the final Thanks slide."""
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(7.05), Inches(10.75), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0, 0, 128); bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.35), Inches(7.07), Inches(10.65), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    p.text = "Prepared by Eng/Ahmed Fekry - Quality & Performance Director (PMO)"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.bold = True; run.font.size = Pt(13); run.font.color.rgb = RGBColor(255, 221, 0)


def _add_table(slide, headers: List[str], rows: List[List[Any]], x: float, y: float, w: float, h: float, font_size: int = 11):
    _, _, _, _, PP_ALIGN, MSO_ANCHOR, _, Inches, Pt, RGBColor = _ppt_imports()
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for i in range(len(headers)):
        table.columns[i].width = Inches(w / len(headers))
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(head)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(244, 122, 42)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(font_size); r.font.color.rgb = RGBColor(0, 0, 0)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(252, 231, 215)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(font_size); r.font.color.rgb = RGBColor(0, 0, 0)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table


def _add_bar_chart(slide, title: str, labels: List[str], values: List[float], x: float, y: float, w: float, h: float):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    box = slide.shapes.add_textbox(Inches(x), Inches(y - 0.25), Inches(w), Inches(0.22))
    p = box.text_frame.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(12); p.runs[0].font.color.rgb = RGBColor(85, 85, 85)
    max_val = max([float(v or 0) for v in values] + [1])
    for i, (label, value) in enumerate(list(zip(labels, values))[:10]):
        yy = y + 0.08 + i * (h / 10)
        slide.shapes.add_textbox(Inches(x), Inches(yy), Inches(1.55), Inches(0.15)).text = str(label)[:18]
        bw = max(0.05, (float(value or 0) / max_val) * (w - 2.2))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 1.75), Inches(yy), Inches(bw), Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(79, 134, 217); bar.line.fill.background()
        slide.shapes.add_textbox(Inches(x + 1.8 + bw), Inches(yy - 0.02), Inches(1.1), Inches(0.17)).text = _fmt_money(value)


def _add_pie_chart(slide, title: str, labels: List[str], values: List[float], x: float, y: float, w: float, h: float):
    _, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION, _, _, _, Inches, _, _ = _ppt_imports()
    chart_data = ChartData()
    chart_data.categories = labels[:8] or ["N/A"]
    chart_data.add_series(title, values[:8] if values else [1])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False




def _ensure_ppt_columns(rows: pd.DataFrame) -> pd.DataFrame:
    """Ensure PPT report columns exist to prevent KeyError during PowerPoint generation."""
    required = {
        "Link Code": "", "Work Order": "", "Region": "N/A", "City": "N/A", "District": "N/A",
        "Project": "N/A", "Stage": "N/A", "SOR Status": "N/A", "SOR Reference Number": "N/A",
        "Cost": 0.0, "Progress": 0.0, "Subclass": "N/A", "Updated": "",
        "Closure Status": "N/A", "Performance": "N/A", "Year": "N/A",
        "Work Order Status": "N/A", "Type": "N/A", "Class": "N/A", "Scope Target": "N/A",
        "1st 50 Invoice Status": "N/A", "Missing MET Actual / PM Review": "N/A",
    }
    if rows is None or rows.empty:
        return pd.DataFrame(columns=list(required.keys()))
    out = rows.copy()
    for col, default in required.items():
        if col not in out.columns:
            out[col] = default
    out["Cost"] = pd.to_numeric(out["Cost"], errors="coerce").fillna(0.0)
    out["Progress"] = pd.to_numeric(out["Progress"], errors="coerce").fillna(0.0)
    for col, default in required.items():
        if col not in ["Cost", "Progress"]:
            out[col] = out[col].fillna(default).astype(str).replace({"": str(default), "nan": str(default), "None": str(default)})
    if "Status" not in out.columns:
        out["Status"] = out["Progress"].apply(_status_from_progress)
    if "Updated_dt" not in out.columns:
        out["Updated_dt"] = out["Updated"].apply(_parse_date_any)
    return out

def apply_ppt_filters(rows: pd.DataFrame, filters: Mapping[str, Any]) -> pd.DataFrame:
    """Apply PPT Builder filters. This replaces iframe-to-Streamlit filter sharing, which is not reliable in Streamlit components."""
    rows = _ensure_ppt_columns(rows)
    if rows.empty or not filters:
        return rows
    out = rows.copy()
    for col, val in filters.items():
        if col not in out.columns:
            continue
        if val is None:
            continue
        vals = val if isinstance(val, (list, tuple, set)) else [val]
        vals = [str(v).strip() for v in vals if str(v).strip() and str(v).strip() != "All"]
        if vals:
            out = out[out[col].astype(str).isin(vals)]
    return out


def _opt_values(df: pd.DataFrame, col: str) -> List[str]:
    if df.empty or col not in df.columns:
        return []
    vals = sorted([str(x) for x in df[col].dropna().astype(str).unique() if str(x).strip() and str(x).strip() not in ["nan", "None"]])
    return vals


def _summary_by_unique_link(rows: pd.DataFrame, dim_col: str, out_col: str) -> pd.DataFrame:
    rows = _ensure_ppt_columns(rows)
    columns = [out_col, "Link Codes", "WOs", "WO Cost", "Avg Progress", "Share"]
    if rows.empty:
        return pd.DataFrame(columns=columns)
    total_cost = float(pd.to_numeric(rows["Cost"], errors="coerce").fillna(0).sum()) or 1.0
    buckets: Dict[str, Dict[str, Any]] = {}
    # Rows without Link Code (for example Scope Target = Initiation) are included
    # in WO and cost, but they must not create artificial Link Code counts.
    for _, r in rows.iterrows():
        link = str(r.get("Link Code", "") or "").strip()
        wo = str(r.get("Work Order", "") or "").strip()
        has_link = bool(link) and link.lower() not in {"n/a", "nan", "none", "-"}
        group_key = f"LINK::{link}" if has_link else f"NO_LINK::{wo or _}"
        dim = str(r.get(dim_col, "N/A") or "N/A").strip() or "N/A"
        cost = float(pd.to_numeric(pd.Series([r.get("Cost", 0)]), errors="coerce").fillna(0).iloc[0])
        prog = float(pd.to_numeric(pd.Series([r.get("Progress", 0)]), errors="coerce").fillna(0).iloc[0])
        if group_key not in buckets:
            buckets[group_key] = {"link": link, "countable": has_link, "rows": [], "dims": {}}
        buckets[group_key]["rows"].append((dim, cost, prog))
        buckets[group_key]["dims"].setdefault(dim, {"cost": 0.0, "n": 0})
        buckets[group_key]["dims"][dim]["cost"] += cost
        buckets[group_key]["dims"][dim]["n"] += 1
    out: Dict[str, Dict[str, Any]] = {}
    for item in buckets.values():
        chosen = sorted(item["dims"].items(), key=lambda kv: (-kv[1]["cost"], -kv[1]["n"], str(kv[0])))[0][0]
        out.setdefault(chosen, {out_col: chosen, "_links": set(), "WOs": 0, "WO Cost": 0.0, "_progress": []})
        if item["countable"]:
            out[chosen]["_links"].add(item["link"])
        for _dim, cost, prog in item["rows"]:
            out[chosen]["WOs"] += 1
            out[chosen]["WO Cost"] += cost
            out[chosen]["_progress"].append(prog)
    records = []
    for item in out.values():
        cost = item["WO Cost"]
        records.append({
            out_col: item[out_col],
            "Link Codes": len(item["_links"]),
            "WOs": item["WOs"],
            "WO Cost": cost,
            "Avg Progress": sum(item["_progress"]) / len(item["_progress"]) if item["_progress"] else 0,
            "Share": cost / total_cost * 100,
        })
    return pd.DataFrame(records, columns=columns).sort_values(["WO Cost", "Link Codes"], ascending=[False, False])


def portfolio_summary_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    return _summary_by_unique_link(rows, "Project", "Project")


def stage_summary_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    return _summary_by_unique_link(rows, "Stage", "Stage")


def sor_summary_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    rows = _ensure_ppt_columns(rows)
    if rows.empty:
        return pd.DataFrame(columns=["SOR Status", "Link Codes", "WO Cost", "Share"])
    links = link_level_dataframe(rows)
    if "SOR Status" not in links.columns:
        links["SOR Status"] = "N/A"
    total_cost = links["Cost"].sum() or 1
    g = links.groupby("SOR Status", dropna=False).agg(
        **{"Link Codes": ("Link Code", "nunique"), "WO Cost": ("Cost", "sum")}
    ).reset_index()
    g["Share"] = g["WO Cost"] / total_cost * 100
    return g.sort_values("WO Cost", ascending=False)


def kpi_cards_dataframe(rows: pd.DataFrame) -> pd.DataFrame:
    rows = _ensure_ppt_columns(rows)
    penalties = penalty_total_filtered(rows)
    links = rows["Link Code"].nunique() if not rows.empty else 0
    cost = rows["Cost"].sum() if not rows.empty else 0
    civil = rows[rows["Subclass"].str.lower().str.contains("civil", na=False)]
    fiber = rows[rows["Subclass"].str.lower().str.contains("fiber", na=False)]
    closed = rows[rows["Closure Status"].astype(str).str.lower().str.contains("closed", na=False)]
    pending = rows[(rows["Progress"] >= 100) & (~rows.index.isin(closed.index))]
    data = [
        ["Total Link Codes", f"{links:,}"],
        ["Total WO Cost", _fmt_money(cost)],
        ["Civil Completion Rate", _fmt_pct(civil["Progress"].mean() if not civil.empty else 0)],
        ["Fiber Completion Rate", _fmt_pct(fiber["Progress"].mean() if not fiber.empty else 0)],
        ["Penalties", _fmt_money(penalties)],
        ["Closed Indicators", f"{closed['Link Code'].nunique():,}"],
        ["Pending Closure", f"{pending['Link Code'].nunique():,}"],
    ]
    return pd.DataFrame(data, columns=["Metric", "Value"])


def _add_cover_slide(prs, blank):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    slide = prs.slides.add_slide(blank)
    if PPT_COVER_PATH.exists():
        slide.shapes.add_picture(str(PPT_COVER_PATH), Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
    else:
        _slide_bg(slide, rgb=(16, 34, 58))
        _add_logo(slide, DAWIYAT_LOGO_PATH, 0.55, 0.35, 2.45, 0.9)
        _add_logo(slide, MET_LOGO_PATH, 10.45, 0.35, 2.15, 0.9)
        box = slide.shapes.add_textbox(Inches(1.0), Inches(3.05), Inches(11.3), Inches(0.75))
        p = box.text_frame.paragraphs[0]
        p.text = "Dawiyat Executive Dashboard Report"
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.bold = True; p.runs[0].font.size = Pt(31); p.runs[0].font.color.rgb = RGBColor(255,255,255)
    # Footer requested on first page only + Thanks page.
    _add_footer(slide)
    return slide


def _add_thanks_slide(prs, blank):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    slide = prs.slides.add_slide(blank)
    _slide_bg(slide)
    title = slide.shapes.add_textbox(Inches(3.3), Inches(2.55), Inches(6.8), Inches(0.9))
    p = title.text_frame.paragraphs[0]
    p.text = "Thanks"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(54); p.runs[0].font.color.rgb = RGBColor(0, 32, 96)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.45), Inches(5.35), Inches(10.45), Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0, 0, 128); bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(1.55), Inches(5.39), Inches(10.25), Inches(0.35))
    p = tx.text_frame.paragraphs[0]
    p.text = "Prepared by Eng/Ahmed Fekry - Quality & Performance Director (PMO)"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(22); p.runs[0].font.color.rgb = RGBColor(255, 221, 0)
    return slide




# ---------------- V38 Executive Presentation dashboard-style helpers ----------------
def _ppt_theme():
    _, _, _, _, _, _, _, _, _, RGBColor = _ppt_imports()
    return {
        "navy": RGBColor(16, 34, 58),
        "text": RGBColor(15, 31, 53),
        "muted": RGBColor(91, 117, 151),
        "line": RGBColor(217, 227, 239),
        "light": RGBColor(248, 251, 255),
        "card": RGBColor(255, 255, 255),
        "blue": RGBColor(37, 99, 235),
        "teal": RGBColor(20, 184, 166),
        "orange": RGBColor(245, 158, 11),
        "red": RGBColor(239, 68, 68),
        "green": RGBColor(34, 160, 107),
        "purple": RGBColor(139, 92, 246),
    }


def _ppt_set_text(shape, text, size=11, bold=False, color=None, align=None):
    _, _, _, _, PP_ALIGN, _, _, _, Pt, RGBColor = _ppt_imports()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text)
    if align is not None:
        p.alignment = align
    if not p.runs:
        p.add_run()
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Aptos"
        if color is not None:
            r.font.color.rgb = color


def _ppt_text(slide, text, x, y, w, h, size=11, bold=False, color=None, align=None):
    _, _, _, _, PP_ALIGN, _, _, Inches, Pt, RGBColor = _ppt_imports()
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _ppt_set_text(box, text, size=size, bold=bold, color=color, align=align)
    return box


def _ppt_header(slide, title: str, subtitle: str = ""):
    _, _, _, _, PP_ALIGN, _, _, Inches, Pt, RGBColor = _ppt_imports()
    theme = _ppt_theme()
    _slide_bg(slide, rgb=(248, 251, 255))
    _add_logo(slide, DAWIYAT_LOGO_PATH, 0.35, 0.18, 1.85, 0.58)
    _add_logo(slide, MET_LOGO_PATH, 11.0, 0.15, 1.7, 0.62)
    _ppt_text(slide, title, 2.15, 0.18, 9.0, 0.35, size=18, bold=True, color=theme["text"], align=PP_ALIGN.CENTER)
    if subtitle:
        _ppt_text(slide, subtitle, 0.35, 0.78, 12.6, 0.25, size=8.5, bold=False, color=theme["muted"])


def _ppt_round_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, radius=True):
    _, _, _, _, _, _, MSO_SHAPE, Inches, _, RGBColor = _ppt_imports()
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb or _ppt_theme()["card"]
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = 9144
    else:
        shp.line.fill.background()
    return shp


def _ppt_card(slide, label, value, meta="", x=0.3, y=1.15, w=2.3, h=1.05, accent=None, percent=None):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    theme = _ppt_theme(); accent = accent or theme["blue"]
    _ppt_round_rect(slide, x, y, w, h, fill_rgb=theme["card"], line_rgb=theme["line"])
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.08), Inches(y+0.06), Inches(w-0.16), Inches(0.045))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    _ppt_text(slide, str(label).upper(), x+0.12, y+0.18, w-0.24, 0.18, size=7.5, bold=True, color=theme["muted"])
    _ppt_text(slide, value, x+0.12, y+0.40, w-0.24, 0.30, size=18, bold=True, color=theme["text"])
    if meta:
        _ppt_text(slide, meta, x+0.12, y+0.73, w-0.24, 0.20, size=7.5, bold=False, color=theme["muted"])
    if percent is not None:
        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.12), Inches(y+h-0.18), Inches(w-0.24), Inches(0.07))
        track.fill.solid(); track.fill.fore_color.rgb = RGBColor(234, 240, 246); track.line.fill.background()
        fill_w = max(0.01, min(w-0.24, (w-0.24) * float(percent)/100))
        fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.12), Inches(y+h-0.18), Inches(fill_w), Inches(0.07))
        fill.fill.solid(); fill.fill.fore_color.rgb = accent; fill.line.fill.background()


def _ppt_chip(slide, text, x, y, w, h, fill, color=None, size=8):
    _, _, _, _, PP_ALIGN, _, _, Inches, Pt, RGBColor = _ppt_imports()
    theme = _ppt_theme(); color = color or theme["text"]
    _ppt_round_rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=None)
    _ppt_text(slide, text, x, y+0.03, w, h-0.04, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)


def _ppt_dashboard_table(slide, headers, rows, x, y, w, h, col_widths=None, font_size=7.2, max_rows=8):
    _, _, _, _, PP_ALIGN, MSO_ANCHOR, _, Inches, Pt, RGBColor = _ppt_imports()
    theme = _ppt_theme()
    rows = rows[:max_rows]
    table = slide.shapes.add_table(len(rows)+1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths[:len(headers)]):
            table.columns[i].width = Inches(cw)
    else:
        for i in range(len(headers)):
            table.columns[i].width = Inches(w/len(headers))
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(head).upper()
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(247, 250, 252)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(font_size); r.font.color.rgb = theme["muted"]
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = theme["card"]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(font_size); r.font.color.rgb = theme["text"]
                    if j in (0,1): r.font.bold = True
    return table


def _ppt_section_box(slide, title, subtitle, x, y, w, h):
    theme = _ppt_theme()
    _ppt_round_rect(slide, x, y, w, h, fill_rgb=theme["card"], line_rgb=theme["line"])
    _ppt_text(slide, title, x+0.12, y+0.10, w-0.24, 0.25, size=13, bold=True, color=theme["text"])
    if subtitle:
        _ppt_text(slide, subtitle, x+0.12, y+0.42, w-0.24, 0.20, size=7.3, color=theme["muted"])


def _ppt_bar_list(slide, labels, values, x, y, w, h, title="", accent=None, max_items=8):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    theme = _ppt_theme(); accent = accent or theme["teal"]
    if title:
        _ppt_text(slide, title, x, y-0.28, w, 0.22, size=9.5, bold=True, color=theme["text"])
    vals = [float(v or 0) for v in values[:max_items]]
    labs = [str(l) for l in labels[:max_items]]
    max_val = max(vals + [1])
    row_h = h / max(1, len(vals))
    for i, (lab, val) in enumerate(zip(labs, vals)):
        yy = y + i*row_h
        _ppt_text(slide, lab[:28], x, yy, w*0.32, row_h*0.55, size=6.8, bold=True, color=theme["text"])
        track_x = x + w*0.36
        track_w = w*0.48
        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(track_x), Inches(yy+0.04), Inches(track_w), Inches(0.08))
        track.fill.solid(); track.fill.fore_color.rgb = RGBColor(234, 240, 246); track.line.fill.background()
        fw = max(0.03, track_w * val / max_val)
        fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(track_x), Inches(yy+0.04), Inches(fw), Inches(0.08))
        fill.fill.solid(); fill.fill.fore_color.rgb = accent; fill.line.fill.background()
        _ppt_text(slide, _fmt_money(val), x+w*0.86, yy-0.01, w*0.14, row_h*0.6, size=6.8, bold=True, color=theme["text"])


def _pct(n, d):
    return (float(n) / float(d) * 100) if d else 0.0


def _critical_text_for_portfolio(rows, stages):
    if stages.empty:
        return "N/A", "No stage data"
    top = stages.iloc[0]
    return str(top["Stage"]), f'Top cost stage: {_fmt_pct(top["Share"])} | {_fmt_money(top["WO Cost"])} SAR | {int(top["Link Codes"]):,} Link Codes | {int(top["WOs"]):,} WOs'


def _sor_cards_data(sor_df, total_links):
    data = []
    for _, r in sor_df.head(5).iterrows():
        data.append((str(r["SOR Status"]), int(r["Link Codes"]), float(r["Share"]), float(r["WO Cost"])))
    return data


def _add_observation_card(slide, title, value, meta, x=10.2, y=0.90, w=2.7, h=0.85, tone="blue"):
    theme = _ppt_theme()
    tones = {"blue": (RGBColor(239,246,255), theme["blue"]), "red": (RGBColor(255,241,242), theme["red"]), "orange": (RGBColor(255,247,237), theme["orange"])}
    fill, accent = tones.get(tone, tones["blue"])
    _ppt_round_rect(slide, x, y, w, h, fill_rgb=fill, line_rgb=accent)
    _ppt_text(slide, title.upper(), x+0.12, y+0.10, w-0.24, 0.16, size=6.6, bold=True, color=accent)
    _ppt_text(slide, value, x+0.12, y+0.31, w-0.24, 0.25, size=13, bold=True, color=_ppt_theme()["text"])
    _ppt_text(slide, meta, x+0.12, y+0.60, w-0.24, 0.17, size=6.4, color=_ppt_theme()["muted"])


def _slide_portfolio(prs, blank, rows, portfolio, stages):
    theme = _ppt_theme(); slide = prs.slides.add_slide(blank)
    _ppt_header(slide, "Executive Portfolio Summary & Cost Exposure", "Project and Stage distribution is calculated from the current PPT Builder filters. Link Code counts and values are based on the same active filtered scope.")
    links = rows["Link Code"].nunique(); total_cost = rows["Cost"].sum(); wos = len(rows)
    top_project = portfolio.iloc[0] if not portfolio.empty else None
    top_stage = stages.iloc[0] if not stages.empty else None
    obs_title, obs_meta = _critical_text_for_portfolio(rows, stages)
    _add_observation_card(slide, "Critical Portfolio Observation", obs_title, obs_meta, tone="blue")
    _ppt_card(slide, "Total Link Codes", f"{links:,}", f"{wos:,} WOs in active filtered scope", 0.28, 1.22, 3.0, 1.05, theme["blue"], 100)
    _ppt_card(slide, "Total WO Cost", _fmt_money(total_cost), "Combined value for the active portfolio scope", 3.48, 1.22, 3.0, 1.05, theme["teal"], 100)
    _ppt_card(slide, "Top Project", _fmt_money(top_project["WO Cost"]) if top_project is not None else "0", (f'{top_project["Project"]} | {int(top_project["Link Codes"]):,} Link Codes | {int(top_project["WOs"]):,} WOs' if top_project is not None else "N/A"), 6.68, 1.22, 3.0, 1.05, theme["purple"], float(top_project["Share"]) if top_project is not None else 0)
    _ppt_card(slide, "Top Stage", _fmt_money(top_stage["WO Cost"]) if top_stage is not None else "0", (f'{top_stage["Stage"]} | {int(top_stage["Link Codes"]):,} Link Codes | {int(top_stage["WOs"]):,} WOs' if top_stage is not None else "N/A"), 9.88, 1.98, 3.0, 1.05, theme["orange"], float(top_stage["Share"]) if top_stage is not None else 0)
    # tables
    _ppt_section_box(slide, "Project Distribution — Link Codes & Cost", "Shows each Project with number of Link Codes, WOs, cost, and average completion.", 0.28, 2.55, 6.35, 4.45)
    _ppt_section_box(slide, "Stage Distribution — Link Codes & Cost", "Shows each Stage with number of Link Codes, WOs, cost, and average completion.", 6.78, 2.55, 6.25, 4.45)
    p_rows = [[r["Project"], f'{int(r["Link Codes"]):,}\n{int(r["WOs"]):,} WOs', _fmt_money(r["WO Cost"]), _fmt_pct(r["Avg Progress"]), _fmt_pct(r["Share"])] for _, r in portfolio.head(6).iterrows()]
    s_rows = [[r["Stage"], f'{int(r["Link Codes"]):,}\n{int(r["WOs"]):,} WOs', _fmt_money(r["WO Cost"]), _fmt_pct(r["Avg Progress"]), _fmt_pct(r["Share"])] for _, r in stages.head(6).iterrows()]
    _ppt_dashboard_table(slide, ["Project", "Link Codes / WOs", "WO Cost", "Avg Progress", "Share"], p_rows, 0.40, 3.10, 5.95, 3.55, [1.55,1.1,1.25,1.1,0.95], 6.0, 6)
    _ppt_dashboard_table(slide, ["Stage", "Link Codes / WOs", "WO Cost", "Avg Progress", "Share"], s_rows, 6.90, 3.10, 5.95, 3.55, [1.80,1.0,1.2,1.0,0.85], 5.7, 6)


def _slide_sor(prs, blank, rows, sor):
    theme = _ppt_theme(); slide = prs.slides.add_slide(blank)
    total_links = max(1, rows["Link Code"].nunique())
    _ppt_header(slide, "Executive SOR Summary & Revenue Exposure", f"Status distribution is calculated by Link Code. Percentages are based on {total_links:,} Link Codes under the current filters.")
    # cards across top
    colors = [theme["green"], theme["blue"], theme["orange"], theme["red"], RGBColor(148,163,184)]
    cards = _sor_cards_data(sor, total_links)
    x = 0.25
    for idx,(name,cnt,share,cost) in enumerate(cards[:5]):
        _ppt_card(slide, name, f"{cnt:,}", f"{_fmt_pct(share)} | Exposure: {_fmt_money(cost)}", x, 1.22, 2.45, 1.05, colors[idx%len(colors)], share)
        x += 2.55
    _ppt_section_box(slide, "SOR / Invoice Funnel", "Dynamic funnel based on SOR status and filtered cost exposure.", 0.28, 2.55, 4.1, 4.4)
    _ppt_bar_list(slide, sor["SOR Status"].astype(str).tolist(), sor["WO Cost"].tolist(), 0.55, 3.12, 3.55, 2.4, "", theme["teal"], 6)
    _ppt_section_box(slide, "SOR Exposure Link Codes", "Filtered by selected PPT Builder scope.", 4.55, 2.55, 8.45, 4.4)
    s_rows = [[r["SOR Status"], f'{int(r["Link Codes"]):,}', _fmt_money(r["WO Cost"]), _fmt_pct(r["Share"])] for _, r in sor.head(8).iterrows()]
    _ppt_dashboard_table(slide, ["SOR Status", "Link Codes", "Exposure Cost", "Share"], s_rows, 4.75, 3.05, 7.95, 3.55, [2.4,1.3,2.6,1.2], 7, 8)


def _slide_stage_summary(prs, blank, rows, stages):
    theme = _ppt_theme(); slide = prs.slides.add_slide(blank)
    _ppt_header(slide, "Executive Overall Stages Summary & Cost Exposure", "Stage distribution is calculated from the current PMO filters. SOR 50% Stage and Waiting SOR are shown as separate columns instead of being included under Other.")
    links = rows["Link Code"].nunique(); cost = rows["Cost"].sum(); top = stages.iloc[0] if not stages.empty else None
    active_districts = rows["District"].nunique() if "District" in rows else 0
    _add_observation_card(slide, "Critical Stage Observation", str(top["Stage"]) if top is not None else "N/A", (f'{_fmt_pct(top["Share"])} of filtered cost | {_fmt_money(top["WO Cost"])} SAR | {int(top["Link Codes"]):,} Link Codes | {int(top["WOs"]):,} WOs' if top is not None else "No stage data"), x=9.9, y=0.88, w=3.0, h=0.90, tone="orange")
    _ppt_card(slide, "Total Link Codes", f"{links:,}", f"{len(rows):,} WOs", 0.28, 1.22, 3.0, 1.05, theme["blue"], 100)
    _ppt_card(slide, "Total WO Cost", _fmt_money(cost), "Combined stage exposure", 3.48, 1.22, 3.0, 1.05, theme["teal"], 100)
    _ppt_card(slide, "Top Cost Stage", _fmt_money(top["WO Cost"]) if top is not None else "0", str(top["Stage"]) if top is not None else "N/A", 6.68, 1.22, 3.0, 1.05, theme["orange"], float(top["Share"]) if top is not None else 0)
    _ppt_card(slide, "Active Districts", f"{active_districts:,}", "Districts with active cost / WOs", 9.88, 1.90, 3.0, 1.05, theme["purple"], 80)
    _ppt_section_box(slide, "Stage Cost Distribution", "Top stages by exposure cost", 0.28, 2.55, 3.6, 4.4)
    _ppt_bar_list(slide, stages["Stage"].astype(str).tolist(), stages["WO Cost"].tolist(), 0.55, 3.12, 3.0, 2.9, "", theme["red"], 7)
    _ppt_section_box(slide, "Stage Cost Ranking", "Ranked cost exposure by stage", 4.05, 2.55, 4.1, 4.4)
    _ppt_bar_list(slide, stages["Stage"].astype(str).tolist(), stages["WO Cost"].tolist(), 4.30, 3.12, 3.55, 2.9, "", theme["purple"], 7)
    _ppt_section_box(slide, "Top 10 District / Stage Cost Drivers", "District and stage concentration", 8.32, 2.55, 4.65, 4.4)
    temp = rows.groupby(["District","Stage"], dropna=False).agg(**{"WO Cost":("Cost","sum"),"Link Codes":("Link Code","nunique"),"WOs":("Work Order","count")}).reset_index().sort_values("WO Cost", ascending=False).head(6)
    d_rows = [[r["District"], r["Stage"], _fmt_money(r["WO Cost"])] for _,r in temp.iterrows()]
    _ppt_dashboard_table(slide, ["District", "Stage", "WO Cost"], d_rows, 8.50, 3.05, 4.25, 3.55, [1.25,1.55,1.45], 6.2, 6)



# ---------------- V39 Executive Dashboard Replica Edition helpers ----------------
def _orange_theme():
    _, _, _, _, _, _, _, _, _, RGBColor = _ppt_imports()
    return {
        "bg": RGBColor(253, 226, 204),
        "bg2": RGBColor(255, 239, 223),
        "orange": RGBColor(244, 122, 42),
        "orange_dark": RGBColor(219, 92, 18),
        "navy": RGBColor(16, 34, 58),
        "text": RGBColor(0, 0, 0),
        "muted": RGBColor(82, 82, 82),
        "white": RGBColor(255, 255, 255),
        "blue": RGBColor(61, 132, 214),
        "teal": RGBColor(20, 184, 166),
        "green": RGBColor(34, 160, 107),
        "red": RGBColor(239, 68, 68),
        "yellow": RGBColor(245, 158, 11),
    }


def _ppt_orange_bg(slide):
    """Background close to the old Presentation DWT orange template."""
    _, _, _, _, _, _, MSO_SHAPE, Inches, _, RGBColor = _ppt_imports()
    th = _orange_theme()
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = th["bg"]
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    base.fill.solid(); base.fill.fore_color.rgb = th["bg"]; base.line.fill.background()
    # subtle light bands to approximate the source PowerPoint template
    for i in range(5):
        x = 0.4 + i * 1.15
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(0), Inches(0.05), Inches(7.5))
        line.rotation = 34
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(255, 241, 225); line.line.fill.background()


def _ppt_orange_header(slide, title: str, subtitle: str = ""):
    _, _, _, _, PP_ALIGN, _, _, Inches, Pt, RGBColor = _ppt_imports()
    th = _orange_theme()
    _ppt_orange_bg(slide)
    # Match Presentation DWT: MET on left, Dawiyat on right
    _add_logo(slide, MET_LOGO_PATH, 0.25, 0.12, 1.95, 0.82)
    _add_logo(slide, DAWIYAT_LOGO_PATH, 10.45, 0.12, 2.55, 0.68)
    _ppt_text(slide, title, 2.25, 0.44, 8.8, 0.38, size=20, bold=True, color=th["text"], align=PP_ALIGN.CENTER)
    if subtitle:
        _ppt_text(slide, subtitle, 0.35, 1.12, 12.6, 0.34, size=10.2, bold=False, color=th["text"])


def _ppt_orange_section_title(slide, text, x, y, w, h=0.38, size=14):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    th = _orange_theme()
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = th["orange"]; shp.line.color.rgb = th["white"]
    _ppt_text(slide, text, x, y+0.06, w, h-0.08, size=size, bold=True, color=th["white"], align=PP_ALIGN.CENTER)


def _ppt_orange_table(slide, headers, rows, x, y, w, h, col_widths=None, font_size=10, max_rows=8):
    _, _, _, _, PP_ALIGN, MSO_ANCHOR, _, Inches, Pt, RGBColor = _ppt_imports()
    th = _orange_theme()
    rows = rows[:max_rows]
    table = slide.shapes.add_table(len(rows)+1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths[:len(headers)]):
            table.columns[i].width = Inches(cw)
    else:
        for i in range(len(headers)):
            table.columns[i].width = Inches(w/len(headers))
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(head)
        cell.fill.solid(); cell.fill.fore_color.rgb = th["orange"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Aptos"; r.font.bold = True; r.font.size = Pt(font_size); r.font.color.rgb = th["text"]
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(252, 231, 215)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = "Aptos"; r.font.size = Pt(font_size); r.font.color.rgb = th["text"]
                    if str(val).lower() in ["grand total", "total"] or i == len(rows):
                        r.font.bold = True
    return table


def _ppt_orange_bar_list(slide, labels, values, x, y, w, h, title="", accent=None, max_items=8):
    _, _, _, _, PP_ALIGN, _, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    th = _orange_theme(); accent = accent or th["blue"]
    if title:
        _ppt_text(slide, title, x, y-0.35, w, 0.24, size=11.5, bold=True, color=th["muted"], align=PP_ALIGN.CENTER)
    vals = [float(v or 0) for v in list(values)[:max_items]]
    labs = [str(l) for l in list(labels)[:max_items]]
    max_val = max(vals + [1])
    row_h = h / max(1, len(vals))
    for i, (lab, val) in enumerate(zip(labs, vals)):
        yy = y + i*row_h
        _ppt_text(slide, lab[:24], x, yy, w*0.28, row_h*0.65, size=8.0, bold=False, color=th["text"], align=PP_ALIGN.RIGHT)
        track_x = x + w*0.31
        track_w = w*0.55
        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(track_x), Inches(yy+0.05), Inches(track_w), Inches(0.12))
        track.fill.solid(); track.fill.fore_color.rgb = RGBColor(239, 226, 216); track.line.fill.background()
        fw = max(0.03, track_w * val / max_val)
        fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(track_x), Inches(yy+0.05), Inches(fw), Inches(0.12))
        fill.fill.solid(); fill.fill.fore_color.rgb = accent; fill.line.fill.background()
        _ppt_text(slide, _fmt_money(val), x+w*0.87, yy-0.02, w*0.13, row_h*0.75, size=8.0, bold=True, color=th["text"])


def _ppt_orange_metric_bar(slide, label, value, max_value, x, y, w, color=None):
    _, _, _, _, _, _, MSO_SHAPE, Inches, _, RGBColor = _ppt_imports()
    th = _orange_theme(); color = color or th["blue"]
    _ppt_text(slide, label, x, y, 2.0, 0.18, size=8.0, bold=True, color=th["muted"])
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+2.0), Inches(y+0.03), Inches(w-3.0), Inches(0.12))
    track.fill.solid(); track.fill.fore_color.rgb = RGBColor(238, 226, 216); track.line.fill.background()
    fill_w = max(0.04, (w-3.0) * float(value or 0) / max(float(max_value or 1), 1.0))
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+2.0), Inches(y+0.03), Inches(fill_w), Inches(0.12))
    fill.fill.solid(); fill.fill.fore_color.rgb = color; fill.line.fill.background()
    _ppt_text(slide, _fmt_money(value), x+w-0.85, y-0.02, 0.85, 0.20, size=8.0, bold=True, color=th["text"])


def _slide_full_scope(prs, blank, rows, cities):
    _, _, _, _, _, _, MSO_SHAPE, Inches, _, RGBColor = _ppt_imports()
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Dawiyat Project Full Scope 2025  & 2026")
    c = cities.head(8).copy()
    body = [[r["Region"], r["City"], f'{int(r["No. of Link Codes"]):,}', _fmt_money(r["WO Amount"])] for _, r in c.iterrows()]
    body.append(["Grand Total", "", f'{int(cities["No. of Link Codes"].sum()):,}', _fmt_money(cities["WO Amount"].sum())])
    _ppt_orange_section_title(slide, "Dawiyat Project Full Scope", 0.35, 1.16, 6.55, 0.42, 14)
    _ppt_orange_table(slide, ["Region", "City", "No. of Link Codes", "WO Amount"], body, 0.35, 1.58, 6.55, 3.9, [1.55, 1.55, 1.55, 1.9], 9.2, 9)
    # chart area on City only
    area = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.15), Inches(1.16), Inches(5.85), Inches(4.3))
    area.fill.solid(); area.fill.fore_color.rgb = RGBColor(235, 235, 235); area.line.color.rgb = RGBColor(210, 210, 210)
    _add_pie_chart(slide, "Business Volume", c["City"].astype(str).tolist(), c["WO Amount"].tolist(), 7.35, 1.45, 5.45, 3.75)


def _slide_regional(prs, blank, rows, cities):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Regional Performance Summary")
    c = cities.head(10).copy()
    body = [[r["Region"], r["City"], f'{int(r["No. of Link Codes"]):,}', int(r["Completed"]), int(r["In Progress"]), int(r["Not Start"]), _fmt_pct(r["Completion %"])] for _, r in c.iterrows()]
    _ppt_orange_table(slide, ["Region", "City", "No. of Link Codes", "Completed", "In Progress", "Not Start", "Completion %"], body, 0.35, 1.45, 12.55, 3.25, [1.25,1.35,1.8,1.45,1.45,1.45,1.55], 8.6, 10)
    _ppt_orange_bar_list(slide, c["City"].astype(str).tolist(), c["WO Amount"].tolist(), 0.60, 5.25, 11.9, 1.35, "City Chart Area", th["blue"], 10)


def _slide_completion(prs, blank, rows, status_cost):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Sites Completion Analysis")
    _ppt_text(slide, "The Site Completion Analysis provides a comprehensive overview of project progress by evaluating the status of all planned FTTH sites. The analysis categorizes sites into completed, in-progress, and pending stages, enabling effective monitoring of project performance and identification of areas requiring additional attention.", 0.30, 1.30, 12.65, 0.75, size=13.5, color=th["text"])
    total_links = max(1, link_level_dataframe(rows)["Link Code"].nunique())
    body = [[r["Status"], int(r["CIVIL"]), _fmt_pct(_pct(r["CIVIL"], total_links)), int(r["FIBRE"]), _fmt_pct(_pct(r["FIBRE"], total_links)), _fmt_pct(_pct(int(r["CIVIL"])+int(r["FIBRE"]), 2*total_links))] for _, r in status_cost.iterrows()]
    _ppt_orange_table(slide, ["Status", "CIVIL", "CIVIL%", "FIBRE", "FIBRE%", "OVERALL%"], body, 1.15, 2.55, 11.0, 2.05, [1.8,1.5,1.55,1.5,1.55,2.0], 12, 3)
    _ppt_orange_bar_list(slide, status_cost["Status"].astype(str).tolist(), status_cost["Total WO Cost"].tolist(), 3.1, 5.20, 7.2, 1.3, "Site Completion Analysis", th["orange"], 3)


def _slide_cost(prs, blank, rows, status_cost):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Sites Cost Analysis")
    _ppt_text(slide, "The Site Cost Analysis provides a detailed assessment of project expenditures and cost efficiency across all implemented FTTH sites. The analysis enables management to monitor project spending, evaluate resource utilization, and ensure alignment with approved budgets and project targets.", 0.30, 1.30, 12.65, 0.75, size=13.0, color=th["text"])
    body = [[r["Status"], int(r["CIVIL"]), _fmt_money(r["WO Cost Civil"]), int(r["FIBRE"]), _fmt_money(r["WO Cost Fibre"]), _fmt_money(r["Total WO Cost"])] for _, r in status_cost.iterrows()]
    _ppt_orange_table(slide, ["Status", "CIVIL", "WO Cost", "FIBRE", "WO Cost", "Total WO Cost"], body, 1.05, 2.35, 11.2, 2.15, [1.8,1.3,2.0,1.3,2.0,2.8], 12, 3)
    _ppt_orange_bar_list(slide, status_cost["Status"].astype(str).tolist(), status_cost["Total WO Cost"].tolist(), 2.1, 5.15, 9.5, 1.35, "Cost Analysis", th["blue"], 3)


def _slide_monthly(prs, blank, rows, months):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Monthly Progress Trend 2025-2026")
    _ppt_text(slide, "Monthly progress trend demonstrates the project's consistent advancement across Civil and Fiber implementation activities. Through effective planning, resource allocation, and coordination with stakeholders, the project performance is monitored throughout the reporting period.", 0.30, 1.30, 12.65, 0.72, size=12.6, color=th["text"])
    body = [[r["Month"], _fmt_pct(r["Civil %"]), _fmt_pct(r["Fiber %"]), _fmt_pct(r["Overall %"])] for _, r in months.iterrows()]
    _ppt_orange_table(slide, ["Month", "Civil %", "Fiber %", "Overall %"], body, 0.65, 2.28, 5.55, 4.05, [2.1,1.15,1.15,1.15], 12, 9)
    _ppt_orange_bar_list(slide, months["Month"].astype(str).tolist(), months["Overall %"].tolist(), 6.75, 2.70, 5.7, 3.30, "Overall Progress", th["blue"], 9)


def _slide_financial(prs, blank, rows):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Executive Financial Report")
    _ppt_text(slide, "Financial exposure summary calculated from the active filtered scope, including cost, penalties, closure exposure, and risk concentration.", 0.35, 1.25, 12.6, 0.35, size=11.5, color=th["text"])
    penalties = penalty_total_filtered(rows)
    completed_cost = rows.loc[rows["Progress"] >= 100, "Cost"].sum()
    risk_cost = rows.loc[rows["Performance"].str.lower().str.contains("risk|off", na=False), "Cost"].sum()
    pending_closure = rows.loc[(rows["Progress"] >= 100) & (~rows["Closure Status"].str.lower().str.contains("closed", na=False)), "Cost"].sum()
    data = [["Total WO Cost", _fmt_money(rows["Cost"].sum())], ["Completed Cost Exposure", _fmt_money(completed_cost)], ["Penalty Amount", _fmt_money(penalties)], ["At Risk / Off Track Cost", _fmt_money(risk_cost)], ["Pending Closure Cost", _fmt_money(pending_closure)]]
    _ppt_orange_table(slide, ["Metric", "Value"], data, 2.15, 1.75, 9.05, 3.25, [5.7,3.35], 10.2, 5)
    vals = [parse_num(x[1]) for x in data]
    maxv = max(vals+[1])
    y = 5.25
    colors = [th["blue"], th["teal"], th["orange"], th["red"], th["green"]]
    for i, (row, val) in enumerate(zip(data, vals)):
        _ppt_orange_metric_bar(slide, row[0], val, maxv, 2.55, y+i*0.32, 8.4, colors[i % len(colors)])


def _slide_readiness(prs, blank, rows, readiness):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "Executive Readiness Report")
    _ppt_text(slide, "Document readiness summary by document type under the current filtered Link Code scope.", 0.35, 1.25, 12.6, 0.32, size=11.5, color=th["text"])
    body = [[r["Document Type"], int(r["Uploaded"]), int(r["Partial"]), int(r["Missing"])] for _, r in readiness.iterrows()] if not readiness.empty else []
    _ppt_orange_table(slide, ["Document Type", "Uploaded", "Partial", "Missing"], body, 1.3, 1.80, 10.7, 4.05, [3.6,2.35,2.35,2.35], 11.5, 7)
    totals = readiness[["Uploaded","Partial","Missing"]].sum() if not readiness.empty else {"Uploaded":0,"Partial":0,"Missing":0}
    _ppt_card(slide, "Uploaded", f'{int(totals["Uploaded"]):,}', "Total uploaded document categories", 1.35, 6.15, 3.1, 0.75, th["green"])
    _ppt_card(slide, "Partial", f'{int(totals["Partial"]):,}', "Total partial document categories", 5.10, 6.15, 3.1, 0.75, th["yellow"])
    _ppt_card(slide, "Missing", f'{int(totals["Missing"]):,}', "Total missing document categories", 8.85, 6.15, 3.1, 0.75, th["red"])


def _slide_pmo_audit(prs, blank, rows):
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "PMO Audit Summary")
    audited_wos = len(rows); audited_links = rows["Link Code"].nunique(); pm_review = int(((rows["Progress"] < 100) & (rows["Closure Status"].str.lower().str.contains("complete|closed", na=False))).sum())
    missing_review = rows["Missing MET Actual / PM Review"].astype(str).str.lower().ne("n/a").sum() if "Missing MET Actual / PM Review" in rows.columns else 0
    cards = [("Audited WOs", f"{audited_wos:,}", f"{audited_links:,} Link Codes", th["blue"]), ("PM Review", f"{pm_review:,}", "Completed status but progress < 100%", th["yellow"]), ("Missing MET Actual", f"{missing_review:,}", "Rows flagged for review", th["red"]), ("Filtered Cost", _fmt_money(rows["Cost"].sum()), "WO Cost / Cost", th["teal"])]
    for i, (lab, val, meta, color) in enumerate(cards):
        _ppt_card(slide, lab, val, meta, 0.45 + i*3.2, 1.25, 2.85, 1.05, color)
    by_stage = stage_summary_dataframe(rows).head(8)
    _ppt_orange_section_title(slide, "Audit Concentration by Stage", 0.70, 2.72, 11.85, 0.40, 13)
    body = [[r["Stage"], f'{int(r["Link Codes"]):,}', f'{int(r["WOs"]):,}', _fmt_money(r["WO Cost"]), _fmt_pct(r["Avg Progress"])] for _, r in by_stage.iterrows()]
    _ppt_orange_table(slide, ["Stage", "Link Codes", "WOs", "WO Cost", "Avg Progress"], body, 0.70, 3.12, 11.85, 3.25, [4.0,1.75,1.35,2.4,2.35], 9.0, 8)


def _slide_assistant_insights(prs, blank, rows):
    _, _, _, _, _, _, _, _, _, RGBColor = _ppt_imports()
    th = _orange_theme(); slide = prs.slides.add_slide(blank)
    _ppt_orange_header(slide, "PMO Report Assistant Insights")
    portfolio = portfolio_summary_dataframe(rows); stages = stage_summary_dataframe(rows); sor = sor_summary_dataframe(rows)
    insights = []
    if not portfolio.empty:
        p = portfolio.iloc[0]; insights.append(f'Top project concentration is {p["Project"]} with {_fmt_money(p["WO Cost"])} SAR and {_fmt_pct(p["Share"])} of filtered cost.')
    if not stages.empty:
        s = stages.iloc[0]; insights.append(f'Top stage exposure is {s["Stage"]} with {_fmt_money(s["WO Cost"])} SAR across {int(s["Link Codes"]):,} Link Codes.')
    if not sor.empty:
        ss = sor.iloc[0]; insights.append(f'SOR concentration is led by {ss["SOR Status"]} with {int(ss["Link Codes"]):,} Link Codes and {_fmt_pct(ss["Share"])} share.')
    penalties = penalty_total_filtered(rows); insights.append(f'Current filtered penalty exposure is {_fmt_money(penalties)} SAR.')
    insights.append(f'Total filtered scope includes {rows["Link Code"].nunique():,} Link Codes, {len(rows):,} WOs, and {_fmt_money(rows["Cost"].sum())} SAR cost exposure.')
    y = 1.45
    for i, txt in enumerate(insights[:7], start=1):
        _ppt_round_rect(slide, 0.85, y, 11.75, 0.65, fill_rgb=RGBColor(252,231,215), line_rgb=th["orange"])
        _ppt_chip(slide, str(i), 1.05, y+0.14, 0.35, 0.30, th["orange"], th["white"], 9)
        _ppt_text(slide, txt, 1.55, y+0.13, 10.7, 0.32, size=10.5, bold=False, color=th["text"])
        y += 0.78



# ---------------- V40 Board Snapshot PowerPoint Helpers ----------------
def _pil_font(size: int = 22, bold: bool = False):
    from PIL import ImageFont
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for c in candidates:
        try:
            return ImageFont.truetype(c, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _pil_text(draw, xy, text, size=22, fill=(18,35,58), bold=False, anchor=None, max_width=None):
    font = _pil_font(size, bold)
    txt = str(text)
    if max_width:
        # basic ellipsis fit for single-line labels
        while len(txt) > 3 and draw.textbbox((0, 0), txt, font=font)[2] > max_width:
            txt = txt[:-2]
        if txt != str(text):
            txt = txt.rstrip() + "…"
    draw.text(xy, txt, fill=fill, font=font, anchor=anchor)


def _pil_money(v):
    try:
        return f"{float(v):,.0f}"
    except Exception:
        return "0"


def _pil_pct(v, digits=1):
    try:
        return f"{float(v):.{digits}f}%"
    except Exception:
        return "0.0%"


def _pil_open_asset(path: Path):
    from PIL import Image
    try:
        if path.exists():
            return Image.open(path).convert("RGBA")
    except Exception:
        return None
    return None


def _pil_paste_logo(img, path: Path, box: tuple[int, int, int, int]):
    logo = _pil_open_asset(path)
    if logo is None:
        return
    x, y, w, h = box
    logo.thumbnail((w, h))
    lx = x + (w - logo.width) // 2
    ly = y + (h - logo.height) // 2
    img.alpha_composite(logo, (lx, ly))


def _pil_round(draw, box, radius=22, fill=(255,255,255), outline=(217,227,239), width=2):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def _pil_card(draw, x, y, w, h, label, value, meta="", color=(37,99,235)):
    _pil_round(draw, (x,y,x+w,y+h), 20, (255,255,255), (210,224,241), 2)
    draw.rounded_rectangle((x, y, x+w, y+6), radius=3, fill=color)
    _pil_text(draw, (x+18, y+25), label.upper(), 16, (83,105,135), True, max_width=w-30)
    _pil_text(draw, (x+18, y+70), value, 30, (15,31,53), True, max_width=w-34)
    if meta:
        _pil_text(draw, (x+18, y+110), meta, 15, (91,117,151), False, max_width=w-30)
    # progress track
    draw.rounded_rectangle((x+18, y+h-28, x+w-18, y+h-17), radius=6, fill=(232,240,248))
    draw.rounded_rectangle((x+18, y+h-28, x+w-55, y+h-17), radius=6, fill=color)


def _pil_table(draw, x, y, w, h, headers, rows, col_widths=None, title=None, header_fill=(246,248,251), max_rows=6, font_size=14):
    if title:
        _pil_text(draw, (x, y-36), title, 23, (18,35,58), True, max_width=w)
    _pil_round(draw, (x,y,x+w,y+h), 18, (255,255,255), (217,227,239), 2)
    if col_widths is None:
        col_widths = [1/len(headers)]*len(headers)
    sw = sum(col_widths)
    col_widths = [cw/sw*w for cw in col_widths]
    row_count = min(len(rows), max_rows)
    header_h = 48
    row_h = (h-header_h) / max(row_count, 1)
    # header
    draw.rectangle((x, y, x+w, y+header_h), fill=header_fill)
    cx = x
    for j, head in enumerate(headers):
        _pil_text(draw, (cx+12, y+17), head.upper(), max(11, font_size-2), (83,105,135), True, max_width=int(col_widths[j]-20))
        cx += col_widths[j]
    # rows
    for i, row in enumerate(rows[:max_rows]):
        yy = y+header_h+i*row_h
        draw.line((x, yy, x+w, yy), fill=(237,242,247), width=1)
        cx = x
        for j, val in enumerate(row[:len(headers)]):
            fill = (15,31,53)
            bold = j in (0, 1)
            _pil_text(draw, (cx+12, yy+row_h/2-8), val, font_size, fill, bold, max_width=int(col_widths[j]-18))
            cx += col_widths[j]
    return


def _pil_orange_bg(title: str):
    from PIL import Image, ImageDraw, ImageFilter
    img = Image.new("RGBA", (1600, 900), (255, 224, 199, 255))
    draw = ImageDraw.Draw(img)
    # subtle diagonal lines matching DWT visual style
    for x in [-120, 40, 210, 390, 570, 760]:
        draw.line((x, 840, x+500, 70), fill=(255,255,255,150), width=7)
        draw.line((x+8, 842, x+508, 72), fill=(198,159,120,70), width=3)
    _pil_paste_logo(img, MET_LOGO_PATH, (70, 35, 210, 95))
    _pil_paste_logo(img, DAWIYAT_LOGO_PATH, (1260, 35, 260, 86))
    _pil_text(draw, (800, 72), title, 33, (0,0,0), True, anchor="mm")
    return img, ImageDraw.Draw(img)


def _pil_white_bg(title: str, subtitle: str = ""):
    from PIL import Image, ImageDraw
    img = Image.new("RGBA", (1600, 900), (248, 251, 255, 255))
    draw = ImageDraw.Draw(img)
    _pil_round(draw, (16, 14, 1584, 876), 24, (255,255,255), (217,227,239), 2)
    _pil_text(draw, (36, 36), title, 30, (18,35,58), True, max_width=980)
    if subtitle:
        _pil_text(draw, (36, 76), subtitle, 16, (91,117,151), False, max_width=1100)
    return img, draw


def _pil_slide_picture(prs, blank, img):
    Presentation, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    slide = prs.slides.add_slide(blank)
    bio = io.BytesIO()
    img.convert("RGB").save(bio, format="PNG", optimize=True)
    bio.seek(0)
    slide.shapes.add_picture(bio, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
    return slide


def _pil_bar_list(draw, x, y, w, rows, label_key, value_key, color=(20,184,166), max_rows=7):
    vals = [float(r.get(value_key, 0) or 0) for r in rows[:max_rows]]
    maxv = max(vals + [1])
    for i, r in enumerate(rows[:max_rows]):
        yy = y + i*50
        label = str(r.get(label_key, "N/A"))
        val = float(r.get(value_key, 0) or 0)
        _pil_text(draw, (x, yy+8), label, 14, (18,35,58), True, max_width=250)
        draw.rounded_rectangle((x+290, yy+14, x+w-90, yy+30), radius=8, fill=(232,240,248))
        bw = max(5, int((w-390) * (val/maxv)))
        draw.rounded_rectangle((x+290, yy+14, x+290+bw, yy+30), radius=8, fill=color)
        _pil_text(draw, (x+w-80, yy+3), _pil_money(val), 14, (15,31,53), True, max_width=78)


def _pil_pie(draw, cx, cy, r, labels, values, colors=None):
    if colors is None:
        colors = [(37,99,235),(239,68,68),(20,184,166),(245,158,11),(139,92,246),(100,116,139),(34,160,107)]
    total = sum(float(v or 0) for v in values) or 1
    start = -90
    for i, v in enumerate(values):
        angle = 360*float(v or 0)/total
        draw.pieslice((cx-r, cy-r, cx+r, cy+r), start, start+angle, fill=colors[i%len(colors)])
        start += angle
    draw.ellipse((cx-r*0.45, cy-r*0.45, cx+r*0.45, cy+r*0.45), fill=(255,255,255))
    _pil_text(draw, (cx, cy-10), _pil_money(total), 24, (15,31,53), True, anchor="mm")


def _rows_from_df(df, cols, max_rows=8):
    if df is None or df.empty:
        return []
    out=[]
    for _, r in df.head(max_rows).iterrows():
        row=[]
        for c in cols:
            v = r.get(c, "")
            if isinstance(v, (int, float)) and ("Cost" in c or "Amount" in c or "WO" in c or "Share" in c):
                if "Share" in c or "%" in c or "Progress" in c:
                    row.append(_pil_pct(v))
                else:
                    row.append(_pil_money(v))
            else:
                row.append(str(v))
        out.append(row)
    return out


def _snapshot_portfolio(rows, portfolio, stages):
    img, draw = _pil_white_bg("Executive Portfolio Summary & Cost Exposure", "Project and Stage distribution is calculated from the current PPT Builder filters. Link Code counts and values are based on the same active filtered scope.")
    total_links = rows["Link Code"].nunique() if not rows.empty else 0
    total_cost = rows["Cost"].sum() if not rows.empty else 0
    top_project = portfolio.iloc[0] if portfolio is not None and not portfolio.empty else None
    top_stage = stages.iloc[0] if stages is not None and not stages.empty else None
    _pil_card(draw, 36, 150, 360, 145, "Total Link Codes", f"{total_links:,}", "active filtered scope", (37,99,235))
    _pil_card(draw, 420, 150, 360, 145, "Total WO Cost", _pil_money(total_cost), "combined value for active scope", (20,184,166))
    _pil_card(draw, 804, 150, 360, 145, "Top Project", _pil_money(top_project["WO Cost"] if top_project is not None else 0), str(top_project["Project"]) if top_project is not None else "N/A", (139,92,246))
    _pil_card(draw, 1188, 150, 360, 145, "Top Stage", _pil_money(top_stage["WO Cost"] if top_stage is not None else 0), str(top_stage["Stage"]) if top_stage is not None else "N/A", (245,158,11))
    # critical observation
    _pil_round(draw, (1215, 28, 1565, 132), 18, (235,245,255), (147,197,253), 2)
    _pil_text(draw, (1235, 50), "CRITICAL PORTFOLIO OBSERVATION", 14, (29,78,216), True)
    _pil_text(draw, (1235, 78), str(top_stage["Stage"])[:28] if top_stage is not None else "N/A", 24, (15,31,53), True)
    if top_stage is not None:
        _pil_text(draw, (1235, 110), f"Top cost stage: {_pil_pct(top_stage['Share'])} | {_pil_money(top_stage['WO Cost'])} SAR", 13, (91,117,151), False)
    p_rows = []
    for _, r in (portfolio if portfolio is not None else pd.DataFrame()).head(6).iterrows():
        p_rows.append([r.get('Project','N/A'), f"{int(r.get('Link Codes',0)):,}\n{int(r.get('WOs',0)):,} WOs", _pil_money(r.get('WO Cost',0)), _pil_pct(r.get('Avg Progress',0),1), _pil_pct(r.get('Share',0),1)])
    s_rows = []
    for _, r in (stages if stages is not None else pd.DataFrame()).head(6).iterrows():
        s_rows.append([r.get('Stage','N/A'), f"{int(r.get('Link Codes',0)):,}\n{int(r.get('WOs',0)):,} WOs", _pil_money(r.get('WO Cost',0)), _pil_pct(r.get('Avg Progress',0),1), _pil_pct(r.get('Share',0),1)])
    _pil_table(draw, 36, 370, 740, 440, ["Project","Link Codes / WOs","WO Cost","Avg Progress","Share"], p_rows, [1.8,1.2,1.2,1,0.8], "Project Distribution — Link Codes & Cost", max_rows=6, font_size=13)
    _pil_table(draw, 824, 370, 740, 440, ["Stage","Link Codes / WOs","WO Cost","Avg Progress","Share"], s_rows, [1.8,1.2,1.2,1,0.8], "Stage Distribution — Link Codes & Cost", max_rows=6, font_size=13)
    return img


def _snapshot_sor(rows, sor):
    img, draw = _pil_white_bg("Executive SOR Summary & Revenue Exposure", "Status distribution is calculated by Link Code. Percentages are based on Link Codes under the current filters.")
    colors = [(34,160,107),(37,99,235),(245,158,11),(239,68,68),(148,163,184)]
    labels = sor["SOR Status"].astype(str).tolist() if sor is not None and not sor.empty else []
    vals = sor["Link Codes"].tolist() if sor is not None and not sor.empty else []
    # cards
    x=36
    for i, (_, r) in enumerate((sor if sor is not None else pd.DataFrame()).head(5).iterrows()):
        _pil_card(draw, x+i*300, 150, 275, 150, str(r.get('SOR Status','N/A')), f"{int(r.get('Link Codes',0)):,}", f"{_pil_pct(r.get('Share',0),1)} | {_pil_money(r.get('WO Cost',0))}", colors[i%len(colors)])
    _pil_round(draw, (1185, 28, 1568, 132), 18, (255,245,245), (252,165,165), 2)
    _pil_text(draw, (1205, 50), "CRITICAL OBSERVATION", 14, (185,28,28), True)
    not_created = next((r for _, r in sor.iterrows() if 'not' in str(r.get('SOR Status','')).lower()), None) if sor is not None and not sor.empty else None
    _pil_text(draw, (1205, 82), f"{int(not_created.get('Link Codes',0)):,} Link Codes" if not_created is not None else "N/A", 24, (185,28,28), True)
    _pil_round(draw, (36, 350, 330, 780), 20, (255,255,255), (217,227,239), 2)
    _pil_text(draw, (60, 380), "Status Distribution", 20, (18,35,58), True)
    _pil_pie(draw, 185, 540, 100, labels, vals, colors)
    ly=665
    for i,(lab,val) in enumerate(zip(labels[:5], vals[:5])):
        draw.ellipse((60, ly+i*24, 72, ly+i*24+12), fill=colors[i%len(colors)])
        _pil_text(draw,(82, ly+i*24-1), f"{lab}  {int(val):,}", 13, (18,35,58), True, max_width=220)
    _pil_round(draw, (360, 350, 960, 780), 20, (255,255,255), (217,227,239), 2)
    _pil_text(draw, (385, 382), "SOR / Invoice Funnel", 20, (18,35,58), True)
    funnel_rows = [{"label": lab, "value": val} for lab,val in zip(labels, vals)]
    _pil_bar_list(draw, 390, 430, 520, funnel_rows, "label", "value", (20,184,166), max_rows=5)
    # exposure table
    link_rows = link_level_dataframe(rows)
    exp = link_rows.sort_values("Cost", ascending=False).head(5)
    trows = [[r.get('City','N/A'), r.get('Link Code',''), r.get('SOR Status','N/A'), _pil_pct(r.get('Progress',0),1), _pil_money(r.get('Cost',0))] for _, r in exp.iterrows()]
    _pil_table(draw, 990, 350, 560, 430, ["City","Link Code","SOR Status","Avg %","Exposure"], trows, [1,1.7,1.1,0.8,1.1], "SOR Exposure Link Codes", max_rows=5, font_size=12)
    return img


def _snapshot_stage(rows, stages):
    img, draw = _pil_white_bg("Executive Overall Stages Summary & Cost Exposure", "Stage distribution is calculated from the current PMO filters. Stage cost exposure is shown in one-slide executive view.")
    total_links=rows['Link Code'].nunique() if not rows.empty else 0
    total_cost=rows['Cost'].sum() if not rows.empty else 0
    top=stages.iloc[0] if stages is not None and not stages.empty else None
    active_districts=rows['District'].nunique() if not rows.empty and 'District' in rows else 0
    _pil_card(draw, 36, 150, 330, 145, "Total Link Codes", f"{total_links:,}", f"{len(rows):,} WOs", (37,99,235))
    _pil_card(draw, 396, 150, 330, 145, "Total WO Cost", _pil_money(total_cost), "combined stage exposure", (20,184,166))
    _pil_card(draw, 756, 150, 330, 145, "Top Cost Stage", _pil_money(top['WO Cost'] if top is not None else 0), str(top['Stage']) if top is not None else 'N/A', (245,158,11))
    _pil_card(draw, 1116, 150, 330, 145, "Active Districts", f"{active_districts:,}", "stages with active cost", (139,92,246))
    _pil_round(draw, (1230, 40, 1565, 120), 18, (255,247,237), (253,186,116), 2)
    _pil_text(draw, (1250, 62), "CRITICAL STAGE OBSERVATION", 14, (194,65,12), True)
    _pil_text(draw, (1250, 92), str(top['Stage'])[:28] if top is not None else 'N/A', 24, (194,65,12), True)
    labels = stages['Stage'].astype(str).tolist() if stages is not None and not stages.empty else []
    vals = stages['WO Cost'].tolist() if stages is not None and not stages.empty else []
    _pil_round(draw, (36, 345, 500, 805), 20, (255,255,255), (217,227,239), 2)
    _pil_text(draw, (60, 375), "Stage Cost Distribution", 20, (18,35,58), True)
    _pil_pie(draw, 255, 525, 105, labels, vals)
    ly=660
    for i,(lab,val) in enumerate(zip(labels[:6], vals[:6])):
        _pil_text(draw,(60, ly+i*24), f"• {lab[:22]}  {_pil_money(val)}", 13, (18,35,58), True, max_width=390)
    stage_rows = [{"label": str(r.get('Stage','')), "value": float(r.get('WO Cost',0) or 0)} for _, r in (stages if stages is not None else pd.DataFrame()).head(7).iterrows()]
    _pil_round(draw, (530, 345, 1045, 805), 20, (255,255,255), (217,227,239), 2)
    _pil_text(draw, (555, 375), "Stage Cost Ranking", 20, (18,35,58), True)
    _pil_bar_list(draw, 560, 425, 430, stage_rows, "label", "value", (239,68,68), max_rows=7)
    # drivers by district+stage
    drivers = rows.groupby(['District','Stage'], dropna=False).agg(**{'Link Codes':('Link Code','nunique'), 'WOs':('Work Order','count'), 'WO Cost':('Cost','sum')}).reset_index().sort_values('WO Cost', ascending=False).head(6)
    drows = [[r['District'], r['Stage'], f"{int(r['Link Codes']):,} LC • {int(r['WOs']):,} WOs", _pil_money(r['WO Cost'])] for _, r in drivers.iterrows()]
    _pil_table(draw, 1080, 345, 485, 460, ["District","Stage","Scope","Cost"], drows, [1.0,1.2,1.0,1.0], "Top 10 District / Stage Cost Drivers", max_rows=6, font_size=12)
    return img


def _snapshot_full_scope(cities):
    img, draw = _pil_orange_bg("Dawiyat Project Full Scope 2025  & 2026")
    t = cities.head(8) if cities is not None else pd.DataFrame()
    rows = [[r.get('Region',''), r.get('City',''), f"{int(r.get('No. of Link Codes',0)):,}", _pil_money(r.get('WO Amount',0))] for _, r in t.iterrows()]
    rows.append(["Grand Total", "", f"{int(cities['No. of Link Codes'].sum() if cities is not None and not cities.empty else 0):,}", _pil_money(cities['WO Amount'].sum() if cities is not None and not cities.empty else 0)])
    _pil_table(draw, 260, 165, 720, 500, ["Region","City","No. of Link Codes","WO Amount"], rows, [1,1,1,1.2], "Dawiyat Project Full Scope", header_fill=(247,122,42), max_rows=9, font_size=14)
    _pil_round(draw, (1010, 165, 1510, 665), 0, (229,229,229), (190,190,190), 2)
    _pil_text(draw, (1260, 235), "Business Volume", 28, (45,45,45), True, anchor="mm")
    labels = t['City'].astype(str).tolist() if not t.empty else []
    vals = t['WO Amount'].tolist() if not t.empty else []
    _pil_pie(draw, 1260, 410, 130, labels, vals)
    # legend
    colors=[(37,99,235),(239,68,68),(20,184,166),(245,158,11),(139,92,246),(100,116,139),(34,160,107)]
    lx=1110; ly=590
    for i, lab in enumerate(labels[:6]):
        draw.rectangle((lx+(i%3)*130, ly+(i//3)*28, lx+10+(i%3)*130, ly+10+(i//3)*28), fill=colors[i%len(colors)])
        _pil_text(draw, (lx+16+(i%3)*130, ly-3+(i//3)*28), lab[:12], 15, (0,0,0), False)
    return img


def _snapshot_regional(cities):
    img, draw = _pil_orange_bg("Regional Performance Summary")
    t = cities.head(8) if cities is not None else pd.DataFrame()
    rows = [[r.get('Region',''), r.get('City',''), f"{int(r.get('No. of Link Codes',0)):,}", int(r.get('Completed',0)), int(r.get('In Progress',0)), int(r.get('Not Start',0)), _pil_pct(r.get('Completion %',0),0)] for _, r in t.iterrows()]
    _pil_table(draw, 120, 180, 1360, 360, ["Region","City","Total Sites","Completed","In Progress","Not Start","Completion %"], rows, [1,1,0.8,0.8,0.8,0.8,0.8], None, header_fill=(247,122,42), max_rows=8, font_size=14)
    # chart area bottom
    chart_rows=[{"label": str(r.get('City','')), "value": float(r.get('WO Amount',0) or 0)} for _, r in t.iterrows()]
    _pil_bar_list(draw, 220, 590, 1050, chart_rows, "label", "value", (79,134,217), max_rows=6)
    return img


def _snapshot_completion(status_cost):
    img, draw = _pil_orange_bg("Sites Completion Analysis")
    _pil_text(draw, (80, 145), "The Site Completion Analysis provides a comprehensive overview of project progress by evaluating the status of all planned FTTH sites.", 21, (0,0,0), False, max_width=1430)
    t = status_cost if status_cost is not None else pd.DataFrame()
    rows=[]
    for _, r in t.iterrows():
        rows.append([r.get('Status',''), int(r.get('CIVIL',0)), int(r.get('FIBRE',0)), _pil_money(r.get('Total WO Cost',0))])
    _pil_table(draw, 250, 255, 1100, 280, ["Status","Civil","Fibre","Total WO Cost"], rows, [1.4,1,1,1.4], None, header_fill=(247,122,42), max_rows=4, font_size=20)
    labels=t['Status'].astype(str).tolist() if not t.empty else []
    vals=t['Total WO Cost'].tolist() if not t.empty else []
    _pil_pie(draw, 800, 690, 90, labels, vals)
    _pil_text(draw,(800,570),"Site Completion Analysis",18,(85,85,85),True,anchor="mm")
    return img


def _snapshot_cost(status_cost):
    img, draw = _pil_orange_bg("Sites Cost Analysis")
    _pil_text(draw, (80, 145), "The Site Cost Analysis provides a detailed assessment of project expenditures and cost efficiency across implemented FTTH sites.", 21, (0,0,0), False, max_width=1430)
    t=status_cost if status_cost is not None else pd.DataFrame()
    rows=[[r.get('Status',''), int(r.get('CIVIL',0)), _pil_money(r.get('WO Cost Civil',0)), int(r.get('FIBRE',0)), _pil_money(r.get('WO Cost Fibre',0)), _pil_money(r.get('Total WO Cost',0))] for _, r in t.iterrows()]
    _pil_table(draw, 170, 250, 1260, 290, ["Status","Civil","WO Cost","Fibre","WO Cost","Total WO Cost"], rows, [1.2,0.8,1.2,0.8,1.2,1.4], None, header_fill=(247,122,42), max_rows=4, font_size=18)
    chart_rows=[{"label": str(r.get('Status','')), "value": float(r.get('Total WO Cost',0) or 0)} for _, r in t.iterrows()]
    _pil_bar_list(draw, 360, 600, 900, chart_rows, "label", "value", (79,134,217), max_rows=4)
    return img


def _snapshot_monthly(months):
    img, draw = _pil_orange_bg("Monthly Progress Trend 2025-2026")
    _pil_text(draw, (80, 150), "Monthly progress trend demonstrates the project's advancement across Civil and Fiber implementation activities.", 21, (0,0,0), False, max_width=1430)
    t=months if months is not None else pd.DataFrame()
    rows=[[r.get('Month',''), _pil_pct(r.get('Civil %',0),0), _pil_pct(r.get('Fiber %',0),0), _pil_pct(r.get('Overall %',0),0)] for _, r in t.head(8).iterrows()]
    _pil_table(draw, 120, 265, 600, 420, ["Month","Civil %","Fiber %","Overall %"], rows, [1.8,1,1,1], None, header_fill=(255,245,225), max_rows=8, font_size=18)
    chart_rows=[{"label": str(r.get('Month','')), "value": float(r.get('Overall %',0) or 0)} for _, r in t.head(8).iterrows()]
    _pil_bar_list(draw, 800, 295, 600, chart_rows, "label", "value", (79,134,217), max_rows=8)
    return img


def _snapshot_financial(rows):
    img, draw = _pil_orange_bg("Executive Financial Report")
    penalties = penalty_total_filtered(rows)
    completed_cost = rows.loc[rows['Progress'] >= 100, 'Cost'].sum() if not rows.empty else 0
    risk_cost = rows.loc[rows['Performance'].str.lower().str.contains('risk|off', na=False), 'Cost'].sum() if not rows.empty else 0
    pending = rows.loc[(rows['Progress'] >= 100) & (~rows['Closure Status'].str.lower().str.contains('closed', na=False)), 'Cost'].sum() if not rows.empty else 0
    data=[['Total WO Cost',_pil_money(rows['Cost'].sum() if not rows.empty else 0)],['Completed Cost Exposure',_pil_money(completed_cost)],['Penalty Amount',_pil_money(penalties)],['At Risk / Off Track Cost',_pil_money(risk_cost)],['Pending Closure Cost',_pil_money(pending)]]
    _pil_text(draw, (90, 150), "Financial exposure summary calculated from the active filtered scope, including cost, penalties, closure exposure, and risk concentration.", 21, (45,45,45), False, max_width=1400)
    _pil_table(draw, 360, 230, 880, 330, ["Metric","Value"], data, [1.6,1], None, header_fill=(247,122,42), max_rows=5, font_size=18)
    chart_rows=[{"label": d[0], "value": float(str(d[1]).replace(',','') or 0)} for d in data]
    _pil_bar_list(draw, 420, 620, 760, chart_rows, "label", "value", (37,99,235), max_rows=5)
    return img


def _snapshot_readiness(readiness):
    img, draw = _pil_orange_bg("Executive Readiness Report")
    t=readiness if readiness is not None else pd.DataFrame()
    rows=[[r.get('Document Type',''), int(r.get('Uploaded',0)), int(r.get('Partial',0)), int(r.get('Missing',0))] for _, r in t.iterrows()]
    _pil_table(draw, 250, 190, 1100, 500, ["Document Type","Uploaded","Partial","Missing"], rows, [1.6,1,1,1], None, header_fill=(247,122,42), max_rows=7, font_size=19)
    return img


def _snapshot_pmo(rows, stages):
    img, draw = _pil_orange_bg("PMO Audit Summary")
    audited=len(rows) if not rows.empty else 0
    links=rows['Link Code'].nunique() if not rows.empty else 0
    pm_review=len(rows[(rows['Progress']<100)&(rows['Progress']>0)]) if not rows.empty else 0
    missing=0
    filtered_cost=rows['Cost'].sum() if not rows.empty else 0
    _pil_card(draw, 250, 160, 285, 105, "Audited WOs", f"{audited:,}", f"{links:,} Link Codes", (37,99,235))
    _pil_card(draw, 570, 160, 285, 105, "PM Review", f"{pm_review:,}", "Completed status but progress < 100%", (245,158,11))
    _pil_card(draw, 890, 160, 285, 105, "Missing MET Actual", f"{missing:,}", "Rows flagged for review", (239,68,68))
    _pil_card(draw, 1210, 160, 285, 105, "Filtered Cost", _pil_money(filtered_cost), "WO Cost / Cost", (20,184,166))
    t=stages.head(6) if stages is not None else pd.DataFrame()
    rows2=[[r.get('Stage',''), int(r.get('Link Codes',0)), int(r.get('WOs',0)), _pil_money(r.get('WO Cost',0)), _pil_pct(r.get('Avg Progress',0),0)] for _, r in t.iterrows()]
    _pil_table(draw, 280, 330, 1040, 420, ["Stage","Link Codes","WOs","WO Cost","Avg Progress"], rows2, [2,0.9,0.8,1.2,0.9], "Audit Concentration by Stage", header_fill=(247,122,42), max_rows=6, font_size=15)
    return img


def _snapshot_assistant(rows, portfolio, stages, sor):
    img, draw = _pil_orange_bg("PMO Report Assistant Insights")
    insights=[]
    if portfolio is not None and not portfolio.empty:
        p=portfolio.iloc[0]; insights.append(f"Top project is {p['Project']} with {_pil_money(p['WO Cost'])} SAR and {int(p['Link Codes']):,} Link Codes.")
    if stages is not None and not stages.empty:
        s=stages.iloc[0]; insights.append(f"Top cost stage is {s['Stage']} with {_pil_money(s['WO Cost'])} SAR and {int(s['Link Codes']):,} Link Codes.")
    if sor is not None and not sor.empty:
        ss=sor.iloc[0]; insights.append(f"SOR concentration is led by {ss['SOR Status']} with {int(ss['Link Codes']):,} Link Codes and {_pil_pct(ss['Share'])} share.")
    insights.append(f"Current filtered penalty exposure is {_pil_money(penalty_total_filtered(rows))} SAR.")
    insights.append(f"Total filtered scope includes {rows['Link Code'].nunique():,} Link Codes, {len(rows):,} WOs, and {_pil_money(rows['Cost'].sum())} SAR cost exposure.")
    y=190
    for i,txt in enumerate(insights[:7], start=1):
        _pil_round(draw,(200,y,1400,y+72),18,(252,231,215),(247,122,42),2)
        draw.ellipse((225,y+20,260,y+55), fill=(247,122,42))
        _pil_text(draw,(242,y+37),str(i),15,(255,255,255),True,anchor="mm")
        _pil_text(draw,(285,y+24),txt,18,(15,31,53),False,max_width=1050)
        y+=90
    return img


def _build_snapshot_image(report_key, rows, cities, status_cost, months, readiness, portfolio, stages, sor):
    if report_key == 'portfolio':
        return _snapshot_portfolio(rows, portfolio, stages)
    if report_key == 'sor_summary':
        return _snapshot_sor(rows, sor)
    if report_key == 'stage_summary':
        return _snapshot_stage(rows, stages)
    if report_key == 'full_scope':
        return _snapshot_full_scope(cities)
    if report_key == 'regional':
        return _snapshot_regional(cities)
    if report_key == 'completion':
        return _snapshot_completion(status_cost)
    if report_key == 'cost':
        return _snapshot_cost(status_cost)
    if report_key == 'monthly':
        return _snapshot_monthly(months)
    if report_key == 'financial':
        return _snapshot_financial(rows)
    if report_key == 'readiness':
        return _snapshot_readiness(readiness)
    if report_key == 'pmo_audit':
        return _snapshot_pmo(rows, stages)
    if report_key == 'assistant_insights':
        return _snapshot_assistant(rows, portfolio, stages, sor)
    return _snapshot_portfolio(rows, portfolio, stages)

def build_ppt_report(selected_reports: List[str], rows: pd.DataFrame | None = None, snapshot_zoom: float = 0.80) -> bytes:
    """V40: One screenshot-style image per selected report slide.
    This avoids split tables/scrollbars and keeps each report inside a single 16:9 slide.
    snapshot_zoom is kept for future tuning; current implementation fits each full report into one image.
    """
    rows = _ensure_ppt_columns(rows if rows is not None else load_ppt_workorders())
    Presentation, ChartData, XL_CHART_TYPE, XL_LEGEND_POSITION, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, Inches, Pt, RGBColor = _ppt_imports()
    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank = prs.slide_layouts[6]

    selected_reports = [r for r in (selected_reports or []) if r != "kpi_cards"]
    order = ["portfolio", "sor_summary", "stage_summary", "full_scope", "regional", "completion", "cost", "monthly", "financial", "readiness", "pmo_audit", "assistant_insights"]
    selected_set = set(selected_reports)

    cities = city_summary_dataframe(rows)
    status_cost = status_cost_dataframe(rows)
    months = monthly_progress_dataframe(rows)
    readiness = readiness_summary_dataframe(rows)
    portfolio = portfolio_summary_dataframe(rows)
    stages = stage_summary_dataframe(rows)
    sor = sor_summary_dataframe(rows)

    _add_cover_slide(prs, blank)
    for report_key in order:
        if report_key not in selected_set:
            continue
        img = _build_snapshot_image(report_key, rows, cities, status_cost, months, readiness, portfolio, stages, sor)
        _pil_slide_picture(prs, blank, img)
    _add_thanks_slide(prs, blank)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def executive_ppt_builder_page() -> None:
    if not can("export_ppt"):
        st.error("You do not have permission to generate PowerPoint presentations.")
        return

    st.title("📊 Executive PPT Builder")
    st.caption("V40 Board Snapshot Edition — every selected Executive Report is exported as one full-slide image, avoiding split tables and scrollbars.")

    if st.button("← Back to Dashboard", use_container_width=True, type="secondary"):
        st.session_state.pop("active_hidden_page", None)
        st.session_state["force_dashboard"] = True
        st.rerun()

    all_rows = load_ppt_workorders()
    if all_rows.empty:
        st.warning("No work order data is available.")
        return

    st.info(
        "Due to Streamlit iframe isolation, the native Streamlit page cannot automatically read filter values selected inside dashboard.html. "
        "For reliability, this page provides the same dashboard filters here and applies them dynamically to the PowerPoint output."
    )

    with st.expander("🎛️ Dashboard Filters for PPT", expanded=True):
        f1, f2, f3 = st.columns(3)

        with f1:
            region_sel = st.multiselect("Region", _opt_values(all_rows, "Region"), default=[])
            city_base = apply_ppt_filters(all_rows, {"Region": region_sel})
            city_sel = st.multiselect("City", _opt_values(city_base, "City"), default=[])
            district_base = apply_ppt_filters(city_base, {"City": city_sel})
            district_sel = st.multiselect("District", _opt_values(district_base, "District"), default=[])
            project_sel = st.multiselect("Project", _opt_values(all_rows, "Project"), default=[])

        with f2:
            stage_sel = st.multiselect("Stage", _opt_values(all_rows, "Stage"), default=[])
            subclass_sel = st.multiselect("Subclass", _opt_values(all_rows, "Subclass"), default=[])
            year_sel = st.multiselect("Year", _opt_values(all_rows, "Year"), default=[])
            wo_status_sel = st.multiselect("Work Order Status", _opt_values(all_rows, "Work Order Status"), default=[])

        with f3:
            sor_sel = st.multiselect("SOR Status", _opt_values(all_rows, "SOR Status"), default=[])
            sor_ref_sel = st.multiselect("SOR Reference Number", _opt_values(all_rows, "SOR Reference Number"), default=[])
            type_sel = st.multiselect("Type", _opt_values(all_rows, "Type"), default=[])
            class_sel = st.multiselect("Class", _opt_values(all_rows, "Class"), default=[])

        extra1, extra2 = st.columns(2)
        with extra1:
            scope_sel = st.multiselect("Scope Target", _opt_values(all_rows, "Scope Target"), default=[])
            invoice50_sel = st.multiselect("1st 50 Invoice Status", _opt_values(all_rows, "1st 50 Invoice Status"), default=[])
            second50_sel = st.multiselect("Second 50% status", _opt_values(all_rows, "Second 50% status"), default=[])
        with extra2:
            missing_met_sel = st.multiselect("Missing MET Actual / PM Review", _opt_values(all_rows, "Missing MET Actual / PM Review"), default=[])

        link_base = apply_ppt_filters(all_rows, {
            "Region": region_sel,
            "City": city_sel,
            "District": district_sel,
            "Project": project_sel,
            "Stage": stage_sel,
            "Subclass": subclass_sel,
            "Year": year_sel,
            "Work Order Status": wo_status_sel,
            "SOR Status": sor_sel,
            "SOR Reference Number": sor_ref_sel,
            "Type": type_sel,
            "Class": class_sel,
            "Scope Target": scope_sel,
            "1st 50 Invoice Status": invoice50_sel,
            "Second 50% status": second50_sel,
            "Missing MET Actual / PM Review": missing_met_sel,
        })
        link_sel = st.multiselect("Link Code", _opt_values(link_base, "Link Code")[:1000], default=[])

    ppt_filters = {
        "Region": region_sel,
        "City": city_sel,
        "District": district_sel,
        "Project": project_sel,
        "Stage": stage_sel,
        "Subclass": subclass_sel,
        "Year": year_sel,
        "Work Order Status": wo_status_sel,
        "SOR Status": sor_sel,
        "SOR Reference Number": sor_ref_sel,
        "Type": type_sel,
        "Class": class_sel,
        "Scope Target": scope_sel,
        "1st 50 Invoice Status": invoice50_sel,
        "Second 50% status": second50_sel,
        "Missing MET Actual / PM Review": missing_met_sel,
        "Link Code": link_sel,
    }
    rows = apply_ppt_filters(all_rows, ppt_filters)

    metrics_links = rows["Link Code"].nunique()
    metrics_cost = rows["Cost"].sum()
    c1, c2, c3 = st.columns(3)
    c1.metric("Filtered Link Codes", f"{metrics_links:,}")
    c2.metric("Filtered Work Orders", f"{len(rows):,}")
    c3.metric("Filtered WO Cost", f"{metrics_cost:,.0f}", help="Uses the same cost logic as dashboard: first positive WO Cost, otherwise Cost.")

    selected = st.multiselect(
        "Select reports to include in PowerPoint",
        options=list(PPT_REPORT_OPTIONS.keys()),
        default=list(PPT_REPORT_OPTIONS.keys()),
        format_func=lambda k: PPT_REPORT_OPTIONS[k],
    )

    st.info("Fixed slides are always included: first Cover page + last Thanks page. Footer is applied only on the first and last pages as requested. Each selected report is exported as one full-slide snapshot image to avoid splitting across slides.")

    snapshot_zoom_label = st.selectbox(
        "PowerPoint Snapshot Scale",
        options=["75%", "80%", "85%"],
        index=1,
        help="Use 75% or 80% when you want the report content to appear smaller and fit more comfortably on one slide.",
    )
    snapshot_zoom = float(snapshot_zoom_label.replace("%", "")) / 100.0

    preview_cols = st.columns(2)
    with preview_cols[0]:
        st.subheader("Selected Slide Order")
        order = ["Fixed Cover"] + [PPT_REPORT_OPTIONS[k] for k in selected] + ["Fixed Thanks"]
        st.write(pd.DataFrame({"#": range(1, len(order)+1), "Slide": order}))
    with preview_cols[1]:
        st.subheader("Filtered Data Scope")
        city_summary = city_summary_dataframe(rows)
        st.dataframe(city_summary[["Region", "City", "No. of Link Codes", "WO Amount"]].head(10), use_container_width=True, hide_index=True)

    if rows.empty:
        st.warning("No records match the selected filters.")
        return
    if not selected:
        st.warning("Please select at least one report.")
        return

    if st.button("📊 Generate PowerPoint", use_container_width=True, type="primary"):
        try:
            with st.spinner("Generating PowerPoint presentation..."):
                ppt_bytes = build_ppt_report(selected, rows, snapshot_zoom=snapshot_zoom)
            st.success("PowerPoint generated successfully.")
            st.download_button(
                "📥 Download Selected PowerPoint Presentation",
                data=ppt_bytes,
                file_name=f"Dawiyat_Executive_Presentation_{ksa_now().strftime('%Y%m%d_%H%M')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        except Exception as exc:
            st.error(f"Unable to generate PowerPoint: {exc}")
            st.info("Please try selecting fewer report sections or contact the PMO System Administrator if the issue continues.")


def reports_page() -> None:
    if not can("reports"):
        st.error("You do not have permission to access Executive Reports.")
        return

    st.title("📄 Executive Reports")
    st.caption("Download PDF Executive Report based on current CSV data.")

    pdf_bytes = build_pdf_report()
    st.download_button(
        "Download PDF Executive Report",
        pdf_bytes,
        "Dawiyat_PMO_Executive_Report.pdf",
        "application/pdf",
        use_container_width=True,
    )

    st.info("The full HTML dashboard PDF can still be exported from the dashboard Export PDF button.")


def admin_page() -> None:
    if str(st.session_state.get("username", "")).strip().lower() != "ahmedfekry" or not can("admin"):
        st.error("You do not have permission to access Admin Board.")
        return

    back_col, spacer_col = st.columns([1.5, 6])
    with back_col:
        if st.button("← Back to Dashboard", use_container_width=True, key="admin_back_to_dashboard"):
            # Do not assign st.session_state["main_nav"] here because the sidebar radio
            # with key="main_nav" has already been created in this rerun.
            # The flag below is applied safely before the radio is created on the next rerun.
            st.session_state.pop("active_hidden_page", None)
            st.session_state["force_dashboard"] = True
            st.rerun()

    st.title("⚙️ Admin Board")
    st.caption("User-Based Permissions Only: users, page access, component access, exports, and permissions are controlled from data/permissions.xlsx.")

    st.subheader("Permission Workbook")
    c1, c2, c3 = st.columns(3)
    c1.metric("permissions.xlsx", "Found" if PERMISSIONS_XLSX_PATH.exists() else "Missing")
    c2.metric("Last Modified", _permissions_last_modified_text())
    c3.metric("Users", f"{len(get_users()):,}")

    st.caption(f"Permission file signature: {_permissions_signature()} | Last loaded this session: {st.session_state.get('permissions_last_loaded', 'Not loaded yet')}")
    reload_col1, reload_col2 = st.columns([1, 4])
    with reload_col1:
        if st.button("🔄 Reload Permissions", use_container_width=True):
            st.session_state["permission_runtime_signature"] = ""
            st.session_state["permissions_manual_reload_at"] = ksa_now().strftime("%Y-%m-%d %H:%M:%S")
            st.rerun()
    with reload_col2:
        st.info("Permissions are read from data/permissions.xlsx on each normal Streamlit rerun. No timed auto-refresh is used. Users see changes after browser refresh or Logout/Login.")

    if PERMISSIONS_XLSX_PATH.exists():
        with open(PERMISSIONS_XLSX_PATH, "rb") as f:
            st.download_button(
                "⬇️ Download Current permissions.xlsx",
                data=f.read(),
                file_name="permissions.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True,
            )

    uploaded_perm = st.file_uploader("Upload / Replace permissions.xlsx", type=["xlsx"], key="permission_xlsx_upload")
    if uploaded_perm is not None:
        DATA_DIR.mkdir(exist_ok=True)
        PERMISSIONS_XLSX_PATH.write_bytes(uploaded_perm.getbuffer())
        st.success("permissions.xlsx uploaded successfully. Permissions will reload on rerun. Existing users should refresh browser or Logout/Login if page access changed.")
        st.rerun()

    sheets = _read_permissions_excel()

    users_table = []
    for username, data in get_users().items():
        users_table.append({
            "username": username,
            "password": str(data.get("password", "")),
            "department": str(data.get("department", data.get("role", "user"))).title(),
            "full_name": str(data.get("full_name", data.get("department", data.get("role", "user")))) or username,
        })
    users_df = pd.DataFrame(users_table)

    user_options = ["All Users"] + sorted(users_df["username"].dropna().astype(str).unique().tolist()) if not users_df.empty else ["All Users"]
    selected_admin_user = st.selectbox(
        "Filter Admin Board by User",
        user_options,
        index=0,
        key="admin_board_user_filter",
        help="Filters Active Users, Page Access, and Component Access tables for the selected username.",
    )

    def _filter_admin_user_df(df: pd.DataFrame) -> pd.DataFrame:
        if selected_admin_user == "All Users" or df.empty:
            return df
        username_col = None
        for candidate in ["username", "Username", "USER", "User"]:
            if candidate in df.columns:
                username_col = candidate
                break
        if username_col is None:
            for col in df.columns:
                if str(col).strip().lower() == "username":
                    username_col = col
                    break
        if username_col is None:
            return df
        return df[df[username_col].astype(str).str.strip().str.lower() == selected_admin_user.strip().lower()]

    st.subheader("Active Users")
    st.caption("Edit users directly here, then press Save. Active=No disables login after the user refreshes the browser or logs in again.")
    users_editor_df = _filter_admin_user_df(sheets.get("Users", users_df).fillna("")) if "Users" in sheets else _filter_admin_user_df(users_df)
    edited_users_df = st.data_editor(
        users_editor_df,
        use_container_width=True,
        hide_index=True,
        num_rows="dynamic",
        key="admin_users_editor",
    )
    if st.button("💾 Save Active Users", use_container_width=True, key="save_admin_users"):
        _update_permission_sheet_from_editor("Users", edited_users_df, selected_admin_user)
        st.success("Users saved. Other sessions will pick up the change after browser refresh or Logout/Login.")
        st.rerun()

    st.subheader("Page Access")
    if "User_Page_Access" in sheets and not sheets["User_Page_Access"].empty:
        page_df = _filter_admin_user_df(sheets["User_Page_Access"].fillna(""))
        edited_page_df = st.data_editor(
            page_df,
            use_container_width=True,
            hide_index=True,
            num_rows="dynamic",
            key="admin_page_access_editor",
        )
        if st.button("💾 Save Page Access", use_container_width=True, key="save_admin_page_access"):
            _update_permission_sheet_from_editor("User_Page_Access", edited_page_df, selected_admin_user)
            st.success("Page Access saved. Other users will see the update after browser refresh or Logout/Login.")
            st.rerun()
    else:
        st.warning("User_Page_Access sheet is missing from permissions.xlsx.")

    st.subheader("Component Access")
    if "User_Component_Access" in sheets and not sheets["User_Component_Access"].empty:
        comp_df = _filter_admin_user_df(sheets["User_Component_Access"].fillna(""))
        edited_comp_df = st.data_editor(
            comp_df,
            use_container_width=True,
            hide_index=True,
            num_rows="dynamic",
            key="admin_component_access_editor",
        )
        if st.button("💾 Save Component Access", use_container_width=True, key="save_admin_component_access"):
            _update_permission_sheet_from_editor("User_Component_Access", edited_comp_df, selected_admin_user)
            st.success("Component Access saved. Other users will see the update after browser refresh or Logout/Login.")
            st.rerun()
    else:
        st.warning("User_Component_Access sheet is missing from permissions.xlsx.")

    st.subheader("Excel Sheet Structure")
    st.markdown("""
The Excel permission workbook must contain these sheets:

- **Users**: Username, Password, Department / Display Role, Active, Full Name / Department
- **User_Page_Access**: Yes/No page permissions by Username
- **User_Component_Access**: Yes/No Show + Export Excel/PDF/PPT by Username and Component/Table
- **Reference_Lists** and **How_To_Use**: helper sheets

Role_Page_Access, Role_Component_Access, and User_Override are ignored. The system now runs on **User-Based Permissions Only**.
""")

    st.subheader("Data Files")
    rows = []
    for name, path in {**DATA_FILES, "permissions.xlsx": PERMISSIONS_XLSX_PATH}.items():
        rows.append({
            "File": name,
            "Exists": path.exists(),
            "Size KB": round(path.stat().st_size / 1024, 1) if path.exists() else 0,
            "Modified": datetime.fromtimestamp(path.stat().st_mtime).strftime("%Y-%m-%d %H:%M:%S") if path.exists() else "",
        })
    st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

    st.info("Permission changes can be edited here or in permissions.xlsx. Users see changes after browser refresh or Logout/Login; no forced timed refresh is used.")


def render_session_bar() -> None:
    """Visible authenticated session bar shown above every page."""
    username = str(st.session_state.get("username", "") or "User")
    role = str(st.session_state.get("role", "viewer") or "viewer").lower()
    role_label = role.title()
    last_login = str(st.session_state.get("last_login", "") or _format_login_time())
    unread_count = unread_notifications_count(username) if st.session_state.get("authenticated") else 0

    st.markdown(
        f"""
        <style>
        .session-bar {{
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 14px;
            padding: 14px 18px;
            margin: 0 0 16px 0;
            border-radius: 18px;
            background: linear-gradient(135deg, #0f172a 0%, #1e3a8a 55%, #0f766e 100%);
            color: #ffffff;
            box-shadow: 0 10px 30px rgba(15, 23, 42, 0.22);
            border: 1px solid rgba(255,255,255,0.20);
        }}
        .session-main {{
            display: flex;
            flex-wrap: wrap;
            gap: 10px 18px;
            align-items: center;
            font-weight: 800;
            font-size: 16px;
            line-height: 1.4;
        }}
        .session-pill {{
            display: inline-flex;
            align-items: center;
            gap: 7px;
            background: rgba(255,255,255,0.13);
            border: 1px solid rgba(255,255,255,0.20);
            padding: 8px 12px;
            border-radius: 999px;
            color: #f8fafc;
            white-space: nowrap;
        }}
        .session-pill strong {{
            color: #ffd400;
            font-weight: 900;
        }}
        .session-active {{
            color: #86efac;
            font-weight: 900;
        }}
        </style>
        <div class="session-bar">
            <div class="session-main">
                <span class="session-pill">👤 <strong>{username}</strong></span>
                <span class="session-pill">👔 Role: <strong>{role_label}</strong></span>
                <span class="session-pill">🕒 Last Login: <strong>{last_login}</strong></span>
                <span class="session-pill">🔔 Notifications: <strong>{unread_count}</strong></span>
                <span class="session-pill session-active">🔐 Session Active</span>
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    col_spacer, col_logout = st.columns([8.7, 1.3])
    with col_logout:
        if st.button("🚪 Logout", use_container_width=True, key="top_logout"):
            _clear_login_query_params()
            st.session_state.clear()
            st.rerun()



def main() -> None:
    if not login_page():
        return
    validate_current_authenticated_user()

    # No timed auto-refresh: permissions update only on browser refresh, rerun, or Logout/Login.

    role = st.session_state.get("role", "user")

    with st.sidebar:
        st.markdown("### Dawiyat PMO Portal V3")
        st.caption(f"User: {st.session_state.get('username','')}")

        all_allowed_pages = allowed_pages_for_current_user()
        hidden_allowed_pages = [p for p in all_allowed_pages if p in HIDDEN_ACTION_PAGES]
        pages = [p for p in all_allowed_pages if p not in HIDDEN_ACTION_PAGES]
        if not pages:
            pages = ["No Access"]

        # Hidden action pages are intentionally excluded from the sidebar. They remain
        # accessible only through Dashboard action buttons and only when the user has permission.
        active_hidden = st.session_state.get("active_hidden_page")
        if active_hidden and active_hidden not in hidden_allowed_pages:
            st.session_state.pop("active_hidden_page", None)

        # Support old force flags by converting them to hidden-page routing.
        if st.session_state.pop("force_document_upload_center", False) and "📤 Document Upload Center" in hidden_allowed_pages:
            st.session_state["active_hidden_page"] = "📤 Document Upload Center"
        if st.session_state.pop("force_ppt_builder", False) and "📊 Executive PPT Builder" in hidden_allowed_pages:
            st.session_state["active_hidden_page"] = "📊 Executive PPT Builder"
        if st.session_state.pop("force_admin_board", False) and "Admin Board" in hidden_allowed_pages:
            st.session_state["active_hidden_page"] = "Admin Board"

        if st.session_state.get("force_dashboard") and "Dashboard" in pages:
            st.session_state.pop("active_hidden_page", None)
            st.session_state["main_nav"] = "Dashboard"
            st.session_state["force_dashboard"] = False

        if st.session_state.get("main_nav") not in pages:
            st.session_state["main_nav"] = pages[0] if pages else "No Access"

        page = st.radio(
            "Navigation",
            pages,
            key="main_nav",
            label_visibility="collapsed",
        )

    render_session_bar()

    active_hidden_page = st.session_state.get("active_hidden_page")
    if active_hidden_page == "📤 Document Upload Center":
        document_upload_page()
        return
    if active_hidden_page == "📊 Executive PPT Builder":
        executive_ppt_builder_page()
        return
    if active_hidden_page == "Admin Board":
        admin_page()
        return
    if active_hidden_page == "Project Updates Center":
        project_updates_center_page()
        return
    if active_hidden_page == "Data Update Agent":
        data_update_agent_page()
        return
    if active_hidden_page == "Notification Center 🔔":
        notification_center_page()
        return
    if active_hidden_page == "Executive Daily Digest":
        executive_daily_digest_page()
        return
    if active_hidden_page == "WhatsApp Agent":
        whatsapp_agent_page()
        return

    if page == "No Access":
        st.error("No pages are currently assigned to your username in permissions.xlsx. Please contact the PMO System Administrator.")
        return
    if page == "Dashboard":
        render_dashboard()
    elif page == "Project Updates Center":
        project_updates_center_page()
    elif page == "Data Update Agent":
        data_update_agent_page()
    elif page == "Notification Center 🔔":
        notification_center_page()
    elif page == "Executive Daily Digest":
        executive_daily_digest_page()
    elif page == "WhatsApp Agent":
        whatsapp_agent_page()
    elif page == "AI Executive Assistant":
        ai_assistant_page()
    elif page == "Smart Alerts":
        smart_alerts_page()
    elif page == "Executive Reports":
        reports_page()
    elif page == "Upload CSV":
        upload_data_page()


if __name__ == "__main__":
    main()
